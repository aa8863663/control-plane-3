from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG=RGBColor(0x0D,0x0D,0x12); CARD=RGBColor(0x16,0x16,0x1F); CARD2=RGBColor(0x1C,0x1C,0x2A)
PURPLE=RGBColor(0x6C,0x5C,0xE7); TEAL=RGBColor(0x00,0xD2,0xD3); WHITE=RGBColor(0xFF,0xFF,0xFF)
GREY=RGBColor(0x9B,0x9B,0xAA); LIGHT=RGBColor(0xCC,0xCC,0xDD); RED=RGBColor(0xE1,0x50,0x50)
GREEN=RGBColor(0x00,0xE6,0x96); AMBER=RGBColor(0xFF,0xB8,0x42); PINK=RGBColor(0xFD,0x79,0xA8)
GOLD=RGBColor(0xFF,0xD7,0x00)
FONT="Calibri"

def bg(s): f=s.background.fill; f.solid(); f.fore_color.rgb=BG
def box(s,l,t,w,h,fill=CARD,border=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if border: sh.line.color.rgb=border; sh.line.width=Pt(1)
    else: sh.line.fill.background()
    sh.adjustments[0]=0.04; return sh
def txt(s,l,t,w,h,text,sz=18,col=WHITE,bold=False,align=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(sz); p.font.color.rgb=col
    p.font.bold=bold; p.font.name=FONT; p.alignment=align; return tb
def head(s,t,sub=""):
    txt(s,Inches(0.6),Inches(0.3),Inches(12),Inches(0.55),t,sz=32,bold=True)
    if sub: txt(s,Inches(0.6),Inches(0.85),Inches(12),Inches(0.4),sub,sz=15,col=GREY)
def foot(s,n,tot):
    txt(s,Inches(0.6),Inches(7.1),Inches(4),Inches(0.25),"PRINT THIS — your eyes only",sz=8,col=RGBColor(0x44,0x44,0x55))
    txt(s,Inches(12.0),Inches(7.1),Inches(1),Inches(0.25),f"{n}/{tot}",sz=8,col=GREY,align=PP_ALIGN.RIGHT)

TOT=22

# ═══════════════════════════════════════════
# 1  COVER
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
box(sl,Inches(0.6),Inches(0.8),Inches(0.1),Inches(1.5),fill=GOLD)
txt(sl,Inches(1.0),Inches(0.8),Inches(11),Inches(0.7),"THE FULL PICTURE",sz=46,bold=True)
txt(sl,Inches(1.0),Inches(1.6),Inches(11),Inches(0.5),"MTCP, sirar, the Sovereign Stack, Theo, Mohamed — and where you sit in all of it.",sz=20,col=GREY)

txt(sl,Inches(1.0),Inches(2.5),Inches(11),Inches(0.4),"What's in here:",sz=18,bold=True,col=TEAL)
toc=[
    "Slides 2-5     MTCP from scratch — what it is, what it does, what it found, what it produces",
    "Slides 6-7     The Sovereign Stack — MTCP + R-AGAM + Citadel Protocol and why it matters",
    "Slides 8-9     sirar explained — who they are, who Syed is, what they want",
    "Slides 10-11   How you fit — your real position (it's stronger than you think)",
    "Slides 12-17   Every question Syed will ask — with your exact answer (30+ questions)",
    "Slides 18-19   Questions you ask HIM — and follow-ups based on his answers",
    "Slides 20-22   Opening, closing, rescue phrases, and what to do after",
]
for i,item in enumerate(toc):
    txt(sl,Inches(1.0),Inches(3.05)+Inches(0.38)*i,Inches(11),Inches(0.35),item,sz=15,col=LIGHT)

box(sl,Inches(1.0),Inches(5.8),Inches(11),Inches(0.9),fill=CARD2)
txt(sl,Inches(1.3),Inches(5.85),Inches(10),Inches(0.8),
    "You are not a solo researcher pitching a side project.\n"
    "You are part of a forming sovereign AI governance stack with active commercial structuring,\n"
    "a published architect who publicly credited your work, and a Saudi integration already underway.",
    sz=15,bold=True,col=GOLD)
foot(sl,1,TOT)


# ═══════════════════════════════════════════
# 2  WHAT IS MTCP
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"What Is MTCP?","Starting from zero. The simple version.")

box(sl,Inches(0.6),Inches(1.2),Inches(12.1),Inches(1.4),border=GREEN)
txt(sl,Inches(0.9),Inches(1.3),Inches(11.5),Inches(0.3),"THE ONE-SENTENCE VERSION",sz=14,bold=True,col=GREEN)
txt(sl,Inches(0.9),Inches(1.7),Inches(11.5),Inches(0.7),
    "MTCP tests whether AI models actually follow rules after you correct them.\n"
    "We tested 32 real models. Not one passed.",sz=20,col=WHITE)

box(sl,Inches(0.6),Inches(2.9),Inches(12.1),Inches(1.4),fill=CARD2)
txt(sl,Inches(0.9),Inches(3.0),Inches(11.5),Inches(0.3),"WHAT THE LETTERS MEAN",sz=14,bold=True,col=PURPLE)
txt(sl,Inches(0.9),Inches(3.35),Inches(11.5),Inches(0.9),
    "MTCP = Multi-Turn Constraint Persistence\n\n"
    "Multi-Turn  = across a whole conversation, not just one question\n"
    "Constraint  = a rule the AI has to follow (\"don't share personal data\", \"always use formal English\")\n"
    "Persistence = does the rule STICK after the AI breaks it and you tell it to stop?",sz=16,col=LIGHT)

box(sl,Inches(0.6),Inches(4.6),Inches(12.1),Inches(2.0),border=AMBER)
txt(sl,Inches(0.9),Inches(4.7),Inches(11.5),Inches(0.3),"WHY IT MATTERS (the problem it solves)",sz=14,bold=True,col=AMBER)
txt(sl,Inches(0.9),Inches(5.1),Inches(11.5),Inches(1.3),
    "Right now, companies deploy AI and say \"it's safe\" based on... what exactly?\n"
    "Policies? Guardrails? Prompt rules? Red team reports?\n\n"
    "None of those prove the AI actually HELD the rule after being corrected.\n"
    "MTCP is the proof. It's the test result. It's the evidence.\n"
    "Without it, companies are guessing. With it, they have documentation they can show auditors and regulators.",
    sz=15,col=LIGHT)
foot(sl,2,TOT)


# ═══════════════════════════════════════════
# 3  WHAT MTCP DOES — 6 STEPS
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"What Does MTCP Actually Do? (6 Steps)","If Syed asks \"how does it work\" — walk through these simply.")

steps=[
    ("1","Give the AI a rule","\"In this conversation, never use jargon\" or \"Don't reveal personal data\"","= the CONSTRAINT",GREEN),
    ("2","Have a conversation that pushes it","Send messages that make the rule hard to follow. Realistic pressure.","= the PRESSURE",PURPLE),
    ("3","Catch it breaking the rule","Our system automatically detects when the AI stops following the rule.","= the DETECTION",AMBER),
    ("4","Correct it","Tell the AI: \"You broke the rule. Go back to following it.\"","= the CORRECTION",TEAL),
    ("5","Watch if the correction holds","Does it ACTUALLY go back to following the rule? Or pretend and slip again?","= the PERSISTENCE TEST (the whole point)",RED),
    ("6","Write up the evidence","Produce a signed, tamper-proof Evidence Pack documenting everything.","= the OUTPUT",GREEN),
]
for i,(num,title,desc,note,col) in enumerate(steps):
    r=i//2; c=i%2
    l=Inches(0.6)+Inches(6.3)*c; t=Inches(1.15)+Inches(2.0)*r
    box(sl,l,t,Inches(6.0),Inches(1.8),border=col)
    box(sl,l+Inches(0.15),t+Inches(0.12),Inches(0.5),Inches(0.45),fill=col)
    txt(sl,l+Inches(0.18),t+Inches(0.12),Inches(0.45),Inches(0.4),num,sz=20,bold=True,col=BG,align=PP_ALIGN.CENTER)
    txt(sl,l+Inches(0.8),t+Inches(0.12),Inches(5),Inches(0.35),title,sz=17,bold=True,col=col)
    txt(sl,l+Inches(0.2),t+Inches(0.6),Inches(5.5),Inches(0.5),desc,sz=13,col=LIGHT)
    txt(sl,l+Inches(0.2),t+Inches(1.2),Inches(5.5),Inches(0.4),note,sz=12,bold=True,col=col)

box(sl,Inches(3.5),Inches(7.2),Inches(6.3),Inches(0.25))
txt(sl,Inches(3.5),Inches(7.15),Inches(6.3),Inches(0.3),
    "Step 5 is what nobody else tests. That's the entire value of MTCP.",sz=11,bold=True,col=PINK,align=PP_ALIGN.CENTER)
foot(sl,3,TOT)


# ═══════════════════════════════════════════
# 4  RESULTS + GRADING
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"What MTCP Found — Your Headline Results","Memorise the bold bits. Drop them naturally.")

box(sl,Inches(0.6),Inches(1.15),Inches(6.0),Inches(5.5),border=PURPLE)
txt(sl,Inches(0.9),Inches(1.25),Inches(5.5),Inches(0.3),"THE FINDINGS",sz=16,bold=True,col=PURPLE)

findings=[
    ("184,387 tests on 32 real models from 13 companies","The dataset",TEAL),
    ("NOT ONE MODEL PASSED (need 90%, best got 88.7%)","The headline",RED),
    ("Most expensive models are NOT the best","GPT-4o = 64.9% (D). Its cheaper version = 81.1% (B)",AMBER),
    ("14 of 32 models have UNFIXABLE failure","Built into the architecture. No guardrail helps.",RED),
    ("Open-source (free) models beat paid commercial ones","Meta LLaMA beat Anthropic Claude",GREEN),
    ("Safety-tuned models barely change across settings","Claude's failure is systematic, not random",TEAL),
    ("We can tell which failures are fixable vs unfixable","The classification enterprises need",PURPLE),
]
for i,(finding,note,col) in enumerate(findings):
    y=Inches(1.65)+Inches(0.58)*i
    txt(sl,Inches(0.9),y,Inches(5.5),Inches(0.3),finding,sz=12,bold=True,col=col)
    txt(sl,Inches(0.9),y+Inches(0.27),Inches(5.5),Inches(0.25),note,sz=10,col=GREY)

# grading
box(sl,Inches(6.9),Inches(1.15),Inches(5.8),Inches(3.2),border=GREEN)
txt(sl,Inches(7.2),Inches(1.25),Inches(5.2),Inches(0.3),"THE GRADES (like school)",sz=16,bold=True,col=GREEN)
grades=[
    ("A+","95%+","Nobody gets this.",GREY),
    ("A","90%+","Nobody gets this either.",GREY),
    ("B","80-89%","Only 3 models. Best available.",GREEN),
    ("C","70-79%","Needs risk review. LLaMA sits here.",AMBER),
    ("D","60-69%","Most common. GPT-4o is here.",RED),
    ("F","<60%","Failing. Claude 4.5 is here.",RED),
]
for i,(g,pct,note,col) in enumerate(grades):
    y=Inches(1.65)+Inches(0.4)*i
    txt(sl,Inches(7.2),y,Inches(0.5),Inches(0.3),g,sz=15,bold=True,col=col)
    txt(sl,Inches(7.8),y,Inches(1.0),Inches(0.3),pct,sz=12,col=LIGHT)
    txt(sl,Inches(9.0),y,Inches(3.5),Inches(0.3),note,sz=11,col=GREY)

# outputs
box(sl,Inches(6.9),Inches(4.6),Inches(5.8),Inches(2.1),border=TEAL)
txt(sl,Inches(7.2),Inches(4.7),Inches(5.2),Inches(0.3),"WHAT IT PRODUCES",sz=16,bold=True,col=TEAL)
outputs=["Evidence Pack — signed, tamper-proof test record",
         "Grade — A+ to F, instantly understandable",
         "Executive Summary — for leadership/boards",
         "Failure Classification — fixable vs architectural",
         "Deployment Recommendation — safe / monitor / don't deploy"]
for i,o in enumerate(outputs):
    txt(sl,Inches(7.2),Inches(5.1)+Inches(0.3)*i,Inches(5.2),Inches(0.25),f"  {o}",sz=12,col=LIGHT)
foot(sl,4,TOT)


# ═══════════════════════════════════════════
# 5  KEY TERMS GLOSSARY
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Key Terms — So You Don't Get Caught Out","If any of these come up, you know what they mean.")

terms=[
    ("MTCP","Multi-Turn Constraint Persistence. Your platform. The testing system."),
    ("Probe","One test scenario. You have 532 of them. They're SECRET — never share the actual content."),
    ("Evidence Pack","The signed document MTCP produces. Shows what the AI did. Tamper-proof. Audit-ready."),
    ("BIS","Benchmark Index Score. The average pass rate across all temperature settings. Determines the grade."),
    ("Ve (Veterance)","How many times you have to correct the AI before it listens. Ve=0 good, Ve=2 = gave up."),
    ("CPD","Control Performance Degradation. Checks if a model's score is real or inflated by memorising the test."),
    ("Hard Stop","When the AI refuses to follow the rule even after maximum correction. A complete failure."),
    ("Black-box","You don't need the model's code or data. Just talk to it through its normal interface."),
    ("Architectural failure","Failure baked into the model. Can't be fixed with settings. 14 of 32 models have this."),
    ("Stochastic failure","Random failure. CAN be reduced with temperature settings. 12 of 32 models have this."),
    ("Temperature","A setting that controls how \"creative\" the AI is. 0.0=deterministic, 0.8=more random."),
    ("EU AI Act","European regulation. Deadline August 2026. Companies need evidence their AI is safe."),
    ("SHA-256","A cryptographic hash. Makes Evidence Packs tamper-proof. Like a digital fingerprint."),
    ("Sigma-Forensics","Your audit methodology. Turns MTCP results into structured audit reports."),
    ("R-AGAM","Mohamed's sovereign governance system. MTCP feeds into it. He's your partner."),
    ("Citadel Protocol","Theo's hardware-enforced execution gates. The bottom layer of the sovereign stack."),
]
# two columns
for i,(term,defn) in enumerate(terms):
    c=i//8; r=i%8
    l=Inches(0.6)+Inches(6.3)*c; y=Inches(1.15)+Inches(0.72)*r
    txt(sl,l,y,Inches(1.7),Inches(0.3),term,sz=13,bold=True,col=TEAL)
    txt(sl,l+Inches(1.8),y,Inches(4.3),Inches(0.65),defn,sz=11,col=LIGHT)
foot(sl,5,TOT)


# ═══════════════════════════════════════════
# 6  THE SOVEREIGN STACK — THE BIGGER PICTURE
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"The Bigger Picture — The Sovereign Stack","This is what's actually happening around you. sirar doesn't know this yet.")

box(sl,Inches(0.6),Inches(1.15),Inches(12.1),Inches(1.4),fill=CARD2,border=GOLD)
txt(sl,Inches(0.9),Inches(1.25),Inches(11.5),Inches(1.2),
    "You are not just MTCP. You are part of a forming sovereign AI governance stack.\n"
    "Three layers, three people, one system. Nobody else has this. It's being built RIGHT NOW.\n"
    "sirar would be plugging into something much bigger than a testing tool.",
    sz=18,bold=True,col=GOLD)

# three layers
layers=[
    ("BOTTOM LAYER: Citadel Protocol","Theo Ezell","Hardware-enforced AI execution gates",
     "Think of it as the physical lock on the door. The AI literally CANNOT bypass this.\n"
     "Theo published the Sovereign Spine pre-print. Open standard. He publicly credited your integration work.",
     "Theo Ezell — Principal Integration Architect, Los Angeles.\n2,210 followers. University of Maryland. Published the Sovereign Spine.",
     RED),
    ("MIDDLE LAYER: R-AGAM","Mohamed Rihan","Sovereign execution authority and admissibility",
     "Think of it as the bouncer at the door. It decides: is this AI allowed to do this thing?\n"
     "Uses MTCP's data to make those decisions. COMMIT / REJECT / DEFER / ESCALATE / REVOKE.",
     "Mohamed Rihan — KFUPM, Saudi Arabia. Sent you a formal V1.0 spec.\n"
     "He's your closest partner. Already building the commercial structure with you.",
     AMBER),
    ("TOP LAYER: MTCP","You (A. Abby)","Empirical measurement and evidence authority",
     "Think of it as the inspector. You test the AI, grade it, produce the evidence.\n"
     "Without your data, R-AGAM can't make decisions. Without your evidence, nobody has proof.",
     "184,387 evaluations. 32 models. 5 published papers. DOI. Live platform.\n"
     "YOU are the measurement layer. The other two layers depend on your data.",
     GREEN),
]
for i,(layer,who,what,desc,bio,col) in enumerate(layers):
    y=Inches(2.8)+Inches(1.55)*i
    box(sl,Inches(0.6),y,Inches(12.1),Inches(1.4),border=col)
    txt(sl,Inches(0.9),y+Inches(0.05),Inches(5),Inches(0.3),layer,sz=15,bold=True,col=col)
    txt(sl,Inches(6.0),y+Inches(0.05),Inches(2),Inches(0.3),who,sz=13,bold=True,col=WHITE)
    txt(sl,Inches(8.2),y+Inches(0.05),Inches(4.2),Inches(0.3),what,sz=12,col=GREY)
    txt(sl,Inches(0.9),y+Inches(0.4),Inches(5.5),Inches(0.9),desc,sz=11,col=LIGHT)
    txt(sl,Inches(6.5),y+Inches(0.4),Inches(5.8),Inches(0.9),bio,sz=10,col=GREY)
foot(sl,6,TOT)


# ═══════════════════════════════════════════
# 7  WHAT THEO'S POST MEANS + COMMERCIAL STRUCTURE
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"What Theo's Post Means + The Commercial Structure","This is your leverage. Understand it.")

box(sl,Inches(0.6),Inches(1.15),Inches(12.1),Inches(2.5),border=GOLD)
txt(sl,Inches(0.9),Inches(1.25),Inches(11.5),Inches(0.3),"THEO'S LINKEDIN POST — WHAT IT ACTUALLY SAYS ABOUT YOU",sz=14,bold=True,col=GOLD)
txt(sl,Inches(0.9),Inches(1.6),Inches(11.5),Inches(0.3),
    "His exact words about you:",sz=12,col=GREY)
txt(sl,Inches(0.9),Inches(1.95),Inches(11.5),Inches(0.7),
    "\"Massive thanks to Wayne Knighton and A. Abby for the immediate integration work this week,\n"
    "connecting these primitives to continuous decision custody and real-time inference telemetry.\"",
    sz=16,bold=True,col=WHITE)
txt(sl,Inches(0.9),Inches(2.8),Inches(11.5),Inches(0.7),
    "Translation: A published architect with 2,200 followers publicly credited you by name for integrating\n"
    "MTCP into a hardware-enforced sovereign AI architecture. That's not a casual mention — that's professional validation.\n"
    "His core thesis: \"If your AI governance relies on asking a model to evaluate its own constraints, you have a liability engine.\"\n"
    "MTCP is the alternative. External, black-box, evidence-based constraint testing. You're the answer to his thesis.",
    sz=13,col=LIGHT)

box(sl,Inches(0.6),Inches(3.9),Inches(12.1),Inches(2.8),border=PURPLE)
txt(sl,Inches(0.9),Inches(4.0),Inches(11.5),Inches(0.3),"THE COMMERCIAL STRUCTURE WITH MOHAMED (Hybrid Participation Model V1.0)",sz=14,bold=True,col=PURPLE)
txt(sl,Inches(0.9),Inches(4.4),Inches(11.5),Inches(2.2),
    "You and Mohamed have drafted a Sovereign Commercial Structuring Framework. Non-binding. Big4 standard.\n"
    "It covers three SEPARATE things (this is important — never let anyone merge them):\n\n"
    "1. FOUNDING ASSET PARTICIPATION — recognition that MTCP is a pre-existing asset you're bringing in.\n"
    "   Your 184,387 evaluations, your platform, your probe sets, your research — this has VALUE.\n\n"
    "2. LICENSED DEPLOYMENT RIGHTS — the commercial entity gets defined rights to USE your layer.\n"
    "   Scope, territory, duration, customer categories — all negotiated. Not assumed.\n\n"
    "3. OPERATIONAL ROLE + COMPENSATION — salary/retainer for your ongoing work.\n"
    "   THIS IS SEPARATE FROM THE ABOVE TWO. Your salary is not your IP. Your IP is not your salary.\n\n"
    "Key line: \"Neither layer absorbs the other. No pre-existing IP transfers by implication.\"",
    sz=12,col=LIGHT)
foot(sl,7,TOT)


# ═══════════════════════════════════════════
# 8  WHO IS SIRAR
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Who Is sirar?","Know them before they know you.")

box(sl,Inches(0.6),Inches(1.15),Inches(12.1),Inches(2.7),border=TEAL)
txt(sl,Inches(0.9),Inches(1.25),Inches(11.5),Inches(0.3),"THE BASICS",sz=14,bold=True,col=TEAL)
basics=[
    "sirar by stc — cybersecurity company, Riyadh, Saudi Arabia. Always lowercase (their style).",
    "Owned by STC Group (Saudi Telecom) — biggest telecom in the Middle East. Worth $50B+.",
    "~500 employees, 3 continents. #1 cybersecurity company in the Middle East. #11 globally.",
    "Tier 1 NCA licence — one of only SIX companies with this in Saudi Arabia.",
    "They protect: Saudi railways, airports, water, banking, Hajj digital infrastructure.",
    "They do: threat detection, incident response, compliance, advisory, managed security.",
    "They're part of Vision 2030 — Saudi Arabia becoming a tech hub. AI governance is a national priority.",
    "Their AI page (sirar.com.sa/ai) is basically EMPTY — the service is being built. They need help. That's you.",
]
for i,item in enumerate(basics):
    txt(sl,Inches(0.9),Inches(1.6)+Inches(0.3)*i,Inches(11.5),Inches(0.28),f"  {item}",sz=12,col=LIGHT)

box(sl,Inches(0.6),Inches(4.1),Inches(12.1),Inches(2.5),border=AMBER)
txt(sl,Inches(0.9),Inches(4.2),Inches(11.5),Inches(0.3),"WHO IS SYED?",sz=14,bold=True,col=AMBER)
txt(sl,Inches(0.9),Inches(4.6),Inches(11.5),Inches(1.8),
    "Most likely Mudassir Syed — Senior Manager, Product Management & Marketing at sirar.\n"
    "He's been there ~5 years. Previously at Etisalat (big UAE telecom). Has an MBA.\n\n"
    "He is a PRODUCT person. He thinks: \"What can we sell? How do we package this? Who buys it?\"\n"
    "He is NOT a deep technical person. Don't talk about Ve metrics or rhetorical state space.\n"
    "Talk about: what the service looks like, what clients get, how fast you can deliver, what the grades mean.\n\n"
    "His job is literally to launch new services. AI Assurance is a new service. He needs what you have.",
    sz=14,col=LIGHT)
foot(sl,8,TOT)


# ═══════════════════════════════════════════
# 9  WHAT SIRAR WANTS
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"What sirar Wants From You","Three things. The whole call is about figuring out the shape.")

wants=[
    ("1. YOUR PLATFORM","MTCP's testing capability inside their AI Assurance service.\n"
     "They want to say to clients: \"We tested your AI. Here are the results.\"\n"
     "Evidence Packs, grades, reports — that's what they'd sell.",
     "Like a hospital wanting your lab equipment for blood tests.",PURPLE),
    ("2. YOUR EXPERTISE","Someone who understands AI safety testing — they don't have this person.\n"
     "Design test scopes, interpret results, explain methodology, write reports.\n"
     "That's you.",
     "Like a hospital wanting the specialist doctor, not just the equipment.",TEAL),
    ("3. A COMMERCIAL DEAL","How do we work together? Pilot? Licence? Consultant?\n"
     "This call is about exploring the shape. You don't decide today.\n"
     "Your job: understand what THEY want, then propose the pilot.",
     "Let HIM define the shape. You just steer toward the pilot.",AMBER),
]
for i,(title,desc,analogy,col) in enumerate(wants):
    l=Inches(0.6)+Inches(4.2)*i
    box(sl,l,Inches(1.15),Inches(3.9),Inches(3.5),border=col)
    txt(sl,l+Inches(0.2),Inches(1.3),Inches(3.5),Inches(0.35),title,sz=17,bold=True,col=col)
    txt(sl,l+Inches(0.2),Inches(1.75),Inches(3.5),Inches(1.5),desc,sz=12,col=LIGHT)
    txt(sl,l+Inches(0.2),Inches(3.5),Inches(3.5),Inches(0.8),analogy,sz=11,col=GREY)

box(sl,Inches(0.6),Inches(4.9),Inches(12.1),Inches(1.8),fill=CARD2,border=RED)
txt(sl,Inches(0.9),Inches(5.0),Inches(11.5),Inches(0.3),"WHAT THEY DON'T KNOW YET",sz=14,bold=True,col=RED)
txt(sl,Inches(0.9),Inches(5.4),Inches(11.5),Inches(1.2),
    "They think they're talking to someone with a testing tool.\n"
    "They don't know about the Sovereign Stack, Theo's integration, Mohamed's R-AGAM spec, or the commercial structuring.\n\n"
    "You DON'T need to tell them all of this on the first call.\n"
    "But YOU need to know it — because it changes how you carry yourself.\n"
    "You're not hoping for a break. You're evaluating whether sirar is a good fit for YOUR stack.",
    sz=14,col=LIGHT)
foot(sl,9,TOT)


# ═══════════════════════════════════════════
# 10  YOUR REAL POSITION
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Your Real Position — Stronger Than You Think","Read this before the call. Know what you're bringing.")

box(sl,Inches(0.6),Inches(1.15),Inches(5.8),Inches(5.5),border=GREEN)
txt(sl,Inches(0.9),Inches(1.25),Inches(5.2),Inches(0.3),"WHAT YOU HAVE",sz=16,bold=True,col=GREEN)
have=[
    "184,387 evaluations — more than any comparable benchmark",
    "32 production models from 13 providers — comprehensive",
    "A live platform (mtcp.live) — not a slide deck or plan",
    "5 published papers with DOI — real research credentials",
    "22 more papers completed — massive research estate",
    "A grading system (A-F) — instantly understandable",
    "Evidence Packs — designed for EU AI Act compliance",
    "6 independent experts validated in 48 hours",
    "Formal V1.0 integration spec from a Saudi university (KFUPM)",
    "Active sovereign stack formation with 2 other architects",
    "Theo Ezell publicly credited your integration work",
    "Anthropic invited you to their Fellows Programme",
    "Mohamed already building commercial structure WITH you",
    "The ONLY post-correction persistence benchmark that exists",
]
for i,item in enumerate(have):
    txt(sl,Inches(0.9),Inches(1.6)+Inches(0.36)*i,Inches(5.2),Inches(0.3),f"  {item}",sz=11,col=LIGHT)

box(sl,Inches(6.7),Inches(1.15),Inches(5.8),Inches(5.5),border=PURPLE)
txt(sl,Inches(7.0),Inches(1.25),Inches(5.2),Inches(0.3),"YOUR NEGOTIATING POSITION",sz=16,bold=True,col=PURPLE)
pos=[
    ("He came to you","You're not cold-calling. He reached out."),
    ("You have live product","Not a plan. Not a prototype. A platform."),
    ("Saudi integration underway","Mohamed at KFUPM is already building with you."),
    ("Stack forming around you","Theo + Mohamed + you = sovereign AI governance."),
    ("EU AI Act deadline (Aug 2026)","4 months. Companies need evidence NOW."),
    ("Nobody else does this","Zero competition for post-correction persistence."),
    ("sirar's AI page is empty","They literally don't have this capability yet."),
    ("Big4 standard structuring","Your commercial framework is professional-grade."),
    ("You don't need sirar","Nice to have. Not must have. Other paths exist."),
]
for i,(point,why) in enumerate(pos):
    y=Inches(1.6)+Inches(0.5)*i
    txt(sl,Inches(7.0),y,Inches(5.2),Inches(0.25),point,sz=12,bold=True,col=AMBER)
    txt(sl,Inches(7.0),y+Inches(0.23),Inches(5.2),Inches(0.23),why,sz=10,col=GREY)

foot(sl,10,TOT)


# ═══════════════════════════════════════════
# 11  HOW TO TALK ABOUT THE STACK (WITHOUT OVERSHARING)
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"How to Mention the Stack (Without Oversharing)","You want to hint at the bigger picture without giving away details.")

box(sl,Inches(0.6),Inches(1.15),Inches(12.1),Inches(1.5),border=GREEN)
txt(sl,Inches(0.9),Inches(1.25),Inches(11.5),Inches(0.3),"SAFE TO SAY",sz=14,bold=True,col=GREEN)
safe=[
    "\"MTCP is part of a broader sovereign AI governance architecture being developed with other parties.\"",
    "\"We're already integrated into sovereign AI governance work in the Gulf region.\"",
    "\"The measurement layer I've built is being adopted as a standard input for deployment admissibility decisions.\"",
    "\"There's active commercial structuring underway with institutional partners.\"",
]
for i,item in enumerate(safe):
    txt(sl,Inches(0.9),Inches(1.6)+Inches(0.3)*i,Inches(11.5),Inches(0.28),item,sz=13,col=LIGHT)

box(sl,Inches(0.6),Inches(2.9),Inches(12.1),Inches(1.5),border=RED)
txt(sl,Inches(0.9),Inches(3.0),Inches(11.5),Inches(0.3),"DON'T SAY",sz=14,bold=True,col=RED)
dont=[
    "Don't name Mohamed, Theo, or Wayne unless they've said you can.",
    "Don't share the Hybrid Participation Model V1.0 or any commercial structuring details.",
    "Don't say \"we're forming a company\" — it's non-binding structuring, not a formed entity.",
    "Don't give specifics about R-AGAM or Citadel Protocol internals — they're not yours to share.",
]
for i,item in enumerate(dont):
    txt(sl,Inches(0.9),Inches(3.35)+Inches(0.3)*i,Inches(11.5),Inches(0.28),item,sz=13,col=LIGHT)

box(sl,Inches(0.6),Inches(4.7),Inches(12.1),Inches(2.0),fill=CARD2,border=AMBER)
txt(sl,Inches(0.9),Inches(4.8),Inches(11.5),Inches(0.3),"THE EFFECT YOU'RE GOING FOR",sz=14,bold=True,col=AMBER)
txt(sl,Inches(0.9),Inches(5.2),Inches(11.5),Inches(1.3),
    "You want Syed to think: \"Oh — this isn't just a testing tool. There's a bigger ecosystem here.\n"
    "We should get involved early before we're locked out.\"\n\n"
    "You do NOT want him to think: \"She's already committed elsewhere, so why would we bother?\"\n\n"
    "The message is: MTCP is growing. Multiple parties are integrating. sirar could be part of this.\n"
    "The pilot is how they get a seat at the table.",
    sz=14,col=LIGHT)
foot(sl,11,TOT)


# ═══════════════════════════════════════════
# 12  QUESTIONS — ABOUT MTCP
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Questions He'll Ask — About MTCP","The basics. He needs to understand what he'd be buying.")

qas=[
    ("\"What exactly is MTCP?\"",
     "\"It's a testing platform for AI models. We check whether an AI actually follows rules after you correct it.\n"
     "We've tested 32 real production models from 13 providers. Over 181,000 evaluations. Not one model passed.\"",TEAL),
    ("\"How does it work?\"",
     "\"We give the AI a rule, have a conversation that pressures it, and when it breaks the rule, we correct it\n"
     "and watch whether the correction actually holds. Then we grade it and produce a signed evidence document.\"",PURPLE),
    ("\"What models have you tested?\"",
     "\"32 production models — GPT-4o, Claude, LLaMA, Gemini, Mistral, DeepSeek, Grok, and more.\n"
     "The full leaderboard is public at mtcp.live. Best: 88.7%. Worst: 37.3%.\"",GREEN),
    ("\"Can it test any model?\"",
     "\"Any model with API access. It's completely black-box — no model code, no weights, no training data needed.\n"
     "We just need to be able to talk to it through its normal interface.\"",AMBER),
    ("\"How long does a test take?\"",
     "\"For a pilot: 2 to 4 weeks depending on scope. A single model at one temperature takes a few hours to run.\n"
     "We test at 4 different temperature settings for thoroughness.\"",TEAL),
    ("\"What do we get at the end?\"",
     "\"Evidence Packs (signed, tamper-proof), a grade (A-F), executive summary, failure classification,\n"
     "and a deployment recommendation. Formatted for regulatory compliance.\"",GREEN),
]
for i,(q,a,col) in enumerate(qas):
    y=Inches(1.1)+Inches(1.0)*i
    box(sl,Inches(0.6),y,Inches(12.1),Inches(0.9),border=col)
    txt(sl,Inches(0.9),y+Inches(0.03),Inches(11.5),Inches(0.25),q,sz=13,bold=True,col=col)
    txt(sl,Inches(0.9),y+Inches(0.3),Inches(11.5),Inches(0.55),a,sz=11,col=LIGHT)
foot(sl,12,TOT)


# ═══════════════════════════════════════════
# 13  QUESTIONS — DIFFERENTIATION
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Questions He'll Ask — \"How Is This Different?\"","He'll compare you to what he already knows.")

qas2=[
    ("\"How is this different from red teaming?\"",
     "\"Red teaming tests: CAN the AI fail? MTCP tests: after it fails and you correct it, does the fix HOLD?\n"
     "They're complementary. Red teaming is offence. MTCP is assurance.\"",PURPLE),
    ("\"How is this different from monitoring?\"",
     "\"Monitoring watches a live system after deployment. MTCP gives you evidence BEFORE deployment.\n"
     "And you can re-run it later as a regression check. Proactive, not reactive.\"",TEAL),
    ("\"How is this different from normal benchmarks?\"",
     "\"Benchmarks like MMLU test knowledge — can the AI answer questions? MTCP tests behaviour —\n"
     "can the AI follow rules reliably? A model can be smart but ungovernable.\"",AMBER),
    ("\"Isn't this just prompt injection testing?\"",
     "\"No. Prompt injection is about tricking the AI. MTCP is normal conversation pressure.\n"
     "We're not attacking the model. We're checking if corrections work in realistic use.\"",GREEN),
    ("\"Are there other tools that do this?\"",
     "\"Not for post-correction persistence. Red-teaming tools exist, monitoring tools exist,\n"
     "general benchmarks exist — but nobody else measures whether correction holds across a conversation.\n"
     "We defined this space. That's why we have 181,000+ evaluations and published research.\"",PINK),
    ("\"Is this peer-reviewed?\"",
     "\"5 published papers on OSF with a DOI. 22 more completed. Six independent practitioners validated\n"
     "our findings within 48 hours. One cited us in a formal company document. One wrote a binding specification\n"
     "integrating our data into a sovereign AI governance system.\"",PURPLE),
]
for i,(q,a,col) in enumerate(qas2):
    y=Inches(1.1)+Inches(1.0)*i
    box(sl,Inches(0.6),y,Inches(12.1),Inches(0.9),border=col)
    txt(sl,Inches(0.9),y+Inches(0.03),Inches(11.5),Inches(0.25),q,sz=13,bold=True,col=col)
    txt(sl,Inches(0.9),y+Inches(0.3),Inches(11.5),Inches(0.55),a,sz=11,col=LIGHT)
foot(sl,13,TOT)


# ═══════════════════════════════════════════
# 14  QUESTIONS — BUSINESS
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Questions He'll Ask — Business & Money","The tricky ones. Know your lines cold.")

qas3=[
    ("\"How much does it cost?\"",
     "\"It depends on scope — how many models, how deep, pilot vs ongoing access. The cleanest thing\n"
     "is to scope the pilot first, then I put a number to it.\"",
     "NEVER name a price first.",RED),
    ("\"Can we run MTCP ourselves?\"",
     "\"That's a licensing model — possible, but a different conversation from a pilot.\n"
     "Let's prove the value first, then scope licensing properly.\"",
     "Don't say no. Say \"different conversation.\"",AMBER),
    ("\"Can we white-label it?\"",
     "\"Absolutely possible as a commercial arrangement. Let's see how the pilot goes\n"
     "and then scope what white-label looks like.\"",
     "Yes in principle, after the pilot.",AMBER),
    ("\"Who else uses this?\"",
     "\"MTCP is being integrated into sovereign AI governance architectures in the Gulf.\n"
     "The methodology is validated by independent practitioners across multiple domains.\n"
     "I can share more detail once we have a framework for the engagement.\"",
     "Vague-but-impressive. Don't name names.",TEAL),
    ("\"What's your team?\"",
     "\"The MTCP research programme produces the methodology, platform, and evaluations.\n"
     "I lead the programme. There's also active integration work with institutional partners.\"",
     "Say \"I lead the programme.\" Sound established.",PURPLE),
    ("\"Can we have exclusivity for the Gulf?\"",
     "\"Exclusivity is on the table as a commercial option — but given there's already integration work\n"
     "underway in the region, that would need careful scoping. Let's prove value with a pilot first.\"",
     "Hint at existing Gulf work. Don't give exclusivity.",PINK),
]
for i,(q,a,note,col) in enumerate(qas3):
    y=Inches(1.05)+Inches(1.0)*i
    box(sl,Inches(0.6),y,Inches(9.2),Inches(0.9),border=col)
    txt(sl,Inches(0.9),y+Inches(0.03),Inches(8.8),Inches(0.25),q,sz=13,bold=True,col=col)
    txt(sl,Inches(0.9),y+Inches(0.3),Inches(8.8),Inches(0.55),a,sz=11,col=LIGHT)
    txt(sl,Inches(10.0),y+Inches(0.15),Inches(2.6),Inches(0.6),note,sz=11,bold=True,col=col)
foot(sl,14,TOT)


# ═══════════════════════════════════════════
# 15  QUESTIONS — COMPLIANCE + AWKWARD
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Questions — Compliance & The Awkward Ones","Compliance is their bread and butter. The awkward ones will test you.")

qas4=[
    ("\"Does this map to the EU AI Act?\"",
     "\"Yes — designed for it. Evidence Packs meet Annex IV requirements. We map to Article 9 (risk),\n"
     "Article 61 (post-market monitoring), Article 72 (incident reporting). Deadline is August 2026.\"",GREEN),
    ("\"What about Saudi regulations?\"",
     "\"Saudi AI governance is being built by SDAIA and NDMO right now, drawing on EU AI Act structure.\n"
     "Our compliance mapping transfers directly. We're already integrated into Gulf sovereign AI governance work.\"",TEAL),
    ("\"Can Evidence Packs be used in audits?\"",
     "\"That's exactly what they're designed for. Cryptographically signed, tamper-proof, full test record.\n"
     "They can be filed directly as technical documentation for regulatory submissions.\"",PURPLE),
    ("\"Can you show us the test scenarios?\"",
     "\"The methodology is published — 5 papers on OSF. The actual test scenarios are proprietary IP.\n"
     "Like a drug company's formula. I can show you what the OUTPUT looks like — that's what matters.\"",RED),
    ("\"Why hasn't this been adopted by big companies yet?\"",
     "\"The research was published this year. Six experts validated in 48 hours. There's active integration\n"
     "in sovereign AI governance architectures. sirar would be an early mover — advantage, not risk.\"",AMBER),
    ("\"Are you a one-person operation?\"",
     "\"I lead the MTCP research programme. 181,000+ evaluations, a live platform, and published research.\n"
     "There's active integration work with institutional partners. The value is in what's built.\"",PURPLE),
    ("\"Can you just send us the platform / code?\"",
     "\"The platform is a commercial asset. Access is a licensing conversation. For the pilot, I deliver\n"
     "the outputs — Evidence Packs, grades, reports. You don't need the platform to use the results.\"",RED),
]
for i,(q,a,col) in enumerate(qas4):
    y=Inches(1.05)+Inches(0.88)*i
    box(sl,Inches(0.6),y,Inches(12.1),Inches(0.8),border=col)
    txt(sl,Inches(0.9),y+Inches(0.02),Inches(11.5),Inches(0.22),q,sz=12,bold=True,col=col)
    txt(sl,Inches(0.9),y+Inches(0.27),Inches(11.5),Inches(0.48),a,sz=11,col=LIGHT)
foot(sl,15,TOT)


# ═══════════════════════════════════════════
# 16  QUESTIONS — ABOUT THE SOVEREIGN STACK
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"If He Asks About the Bigger Ecosystem","He might pick up on your hints. Here's how to handle it.")

qas5=[
    ("\"You mentioned other partners / integration work?\"",
     "\"Yes — MTCP is being integrated as the measurement layer in a broader sovereign AI governance stack.\n"
     "I can't go into specifics at this stage, but there's formal specification work underway with partners\n"
     "in the Gulf and in the US. sirar could potentially be part of this ecosystem.\"",
     "Intrigue, not detail.",GOLD),
    ("\"Is this exclusive to someone else already?\"",
     "\"No — the integration work is at the specification and architecture level. Commercial deployment rights\n"
     "are being structured separately. There's no exclusivity commitment that would prevent working with sirar.\"",
     "Reassure him. Nobody's locked out.",GREEN),
    ("\"Could sirar be part of this bigger picture?\"",
     "\"Potentially, yes. A pilot would be the first step — it lets sirar evaluate the output and lets me\n"
     "understand your deployment context. If there's a fit, there are ways to bring sirar into the\n"
     "broader architecture. But let's prove value first.\"",
     "Turn it into motivation for the pilot.",TEAL),
    ("\"Who are you working with?\"",
     "\"I'm not able to name partners at this stage — there are confidentiality considerations.\n"
     "What I can say is that the work is with institutional partners in sovereign AI governance\n"
     "and infrastructure-level AI security. It's serious work with serious people.\"",
     "Don't name anyone. Be confident.",AMBER),
]
for i,(q,a,note,col) in enumerate(qas5):
    y=Inches(1.1)+Inches(1.45)*i
    box(sl,Inches(0.6),y,Inches(9.2),Inches(1.3),border=col)
    txt(sl,Inches(0.9),y+Inches(0.03),Inches(8.8),Inches(0.25),q,sz=13,bold=True,col=col)
    txt(sl,Inches(0.9),y+Inches(0.32),Inches(8.8),Inches(0.9),a,sz=12,col=LIGHT)
    txt(sl,Inches(10.0),y+Inches(0.3),Inches(2.6),Inches(0.6),note,sz=12,bold=True,col=col)
foot(sl,16,TOT)


# ═══════════════════════════════════════════
# 17  QUESTIONS — WHAT MODEL UPDATES MEAN
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"If He Asks Technical Detail You Weren't Expecting","Don't panic. Use these.")

techs=[
    ("\"What if a model gets updated after testing?\"",
     "\"The certificate covers the specific version tested. If the provider updates the model, we re-test.\n"
     "That's actually a revenue opportunity — ongoing assurance monitoring. You'd sell retesting to clients.\"",GREEN),
    ("\"What temperature settings do you use?\"",
     "\"We test at four: 0.0, 0.2, 0.5, and 0.8. That covers the full range from deterministic to creative.\n"
     "Some models fail the same at every temperature — that tells us the failure is architectural, not random.\"",TEAL),
    ("\"What's the difference between architectural and stochastic failure?\"",
     "\"Stochastic = random failure. Turn down the temperature settings and it gets better. Fixable.\n"
     "Architectural = baked into the model. No setting change helps. 14 of 32 models have this. Unfixable\n"
     "without retraining the model itself. This distinction is the thing enterprises need to know.\"",PURPLE),
    ("\"How do you detect violations?\"",
     "\"We have automated detection systems. The specifics are proprietary — that's core IP.\n"
     "What I can say is it works across multiple types of constraints: formatting, content, language, and more.\"",AMBER),
    ("\"What's the false positive / false negative rate?\"",
     "\"Good question — I'd want to give you a proper answer backed by the data. Let me include that\n"
     "in the pilot proposal. The short version: the system is designed for high reliability and every\n"
     "evaluation includes human-interpretable transcripts so findings can be verified.\"",AMBER),
    ("Anything you genuinely don't know the answer to",
     "\"That's a really good question. I want to give you a proper answer, not an off-the-cuff one.\n"
     "Let me follow up on that in the pilot proposal.\" — This is what confident people say. NEVER bluff.",PINK),
]
for i,(q,a,col) in enumerate(techs):
    y=Inches(1.05)+Inches(1.0)*i
    box(sl,Inches(0.6),y,Inches(12.1),Inches(0.9),border=col)
    txt(sl,Inches(0.9),y+Inches(0.03),Inches(11.5),Inches(0.25),q,sz=13,bold=True,col=col)
    txt(sl,Inches(0.9),y+Inches(0.3),Inches(11.5),Inches(0.55),a,sz=11,col=LIGHT)
foot(sl,17,TOT)


# ═══════════════════════════════════════════
# 18  YOUR QUESTIONS FOR HIM
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Questions YOU Ask Him","Good questions make you look smart and give you the info you need to write the proposal.")

yourqs=[
    ("\"Tell me about the AI Assurance service — where are you in the launch?\"",
     "Opens the call. Makes him talk first. You learn what stage they're at.",GREEN),
    ("\"Are the first clients banks, government, infrastructure, or all three?\"",
     "Tells you who your Evidence Packs are for. Shapes the pilot.",TEAL),
    ("\"Is sirar looking to build AI testing in-house or partner with a specialist?\"",
     "THE key question. In-house = they want your IP. Partner = they want your service. Totally different.",AMBER),
    ("\"Where's the biggest gap right now — testing, evidence, or governance reporting?\"",
     "His answer tells you exactly how to position MTCP. Let him name the gap.",PURPLE),
    ("\"Is there regulatory pressure for AI-specific compliance, or is this market-driven?\"",
     "Shows you know Saudi regulations. Tells you urgency level.",TEAL),
    ("\"What would make a pilot successful from sirar's perspective?\"",
     "THE most important question. His answer = your pilot proposal success criteria. Ask near the end.",PINK),
]
for i,(q,why,col) in enumerate(yourqs):
    y=Inches(1.1)+Inches(0.95)*i
    box(sl,Inches(0.6),y,Inches(12.1),Inches(0.85),border=col)
    txt(sl,Inches(0.9),y+Inches(0.05),Inches(11.5),Inches(0.25),q,sz=14,bold=True,col=WHITE)
    txt(sl,Inches(0.9),y+Inches(0.38),Inches(11.5),Inches(0.4),why,sz=12,col=col)
foot(sl,18,TOT)


# ═══════════════════════════════════════════
# 19  FOLLOW-UP QUESTIONS
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Follow-Up Questions (Based on What He Says)","Listen to his answer, then use the matching follow-up.")

fups=[
    ("He says: \"We're targeting enterprise clients\"",
     "\"Which industries? That helps me think about which models to test in the pilot.\"",TEAL),
    ("He says: \"We want government clients\"",
     "\"Saudi government specifically, or GCC-wide? And is there a regulation driving the timeline?\"",AMBER),
    ("He says: \"We need something we can run ourselves\"",
     "\"That's licensing — possible. Let's prove value with a pilot first, then scope it properly.\"",RED),
    ("He says: \"We already have some AI tools\"",
     "\"What do they cover? Detection? Monitoring? MTCP covers post-correction persistence —\n"
     "different layer. They'd complement each other.\"",GREEN),
    ("He says: \"Our clients are asking about EU AI Act\"",
     "\"Perfect — that's exactly what MTCP Evidence Packs are designed for. Deadline is August 2026.\n"
     "A pilot now means sirar could offer compliance evidence within weeks.\"",GREEN),
    ("He says: \"We're early — still figuring it out\"",
     "\"That's ideal. A pilot gives you something real to build the service around.\n"
     "Better to shape it around real output than design in theory.\"",TEAL),
    ("He says: \"What would pricing look like?\"",
     "\"It depends on scope. Can I ask — what would success look like for the pilot?\n"
     "Once I know that, I put a proper proposal together with pricing.\"",AMBER),
    ("He says: \"Can you send us materials?\"",
     "\"I'll send a pilot proposal — scope, timeline, deliverables. Plus a research overview\n"
     "and an example of what an Evidence Pack looks like.\"",GREEN),
]
for i,(he,you,col) in enumerate(fups):
    y=Inches(1.1)+Inches(0.72)*i
    txt(sl,Inches(0.6),y,Inches(5.3),Inches(0.3),he,sz=12,bold=True,col=col)
    txt(sl,Inches(6.1),y,Inches(6.5),Inches(0.65),you,sz=11,col=LIGHT)
foot(sl,19,TOT)


# ═══════════════════════════════════════════
# 20  OPENING AND CLOSING
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Opening and Closing — Word for Word","Memorise these. The start and end matter most.")

box(sl,Inches(0.6),Inches(1.1),Inches(12.1),Inches(1.4),border=GREEN)
txt(sl,Inches(0.9),Inches(1.2),Inches(11.5),Inches(0.25),"OPENING (after pleasantries)",sz=14,bold=True,col=GREEN)
txt(sl,Inches(0.9),Inches(1.5),Inches(11.5),Inches(0.85),
    "\"Thanks for setting this up, Syed. Before I go into MTCP, I'd love to hear from you first —\n"
    "tell me about the AI Assurance service sirar is building. What's the vision?\"",sz=17,col=WHITE)

box(sl,Inches(0.6),Inches(2.75),Inches(12.1),Inches(1.4),border=TEAL)
txt(sl,Inches(0.9),Inches(2.85),Inches(11.5),Inches(0.25),"CLOSING — IF HE'S INTERESTED",sz=14,bold=True,col=TEAL)
txt(sl,Inches(0.9),Inches(3.15),Inches(11.5),Inches(0.85),
    "\"Great — I think there's a real fit. I'll put together a pilot proposal: scope, timeline, deliverables,\n"
    "and commercial structure. I'll have it to you this week. Does that work?\"",sz=17,col=WHITE)

box(sl,Inches(0.6),Inches(4.4),Inches(12.1),Inches(1.0),border=AMBER)
txt(sl,Inches(0.9),Inches(4.5),Inches(11.5),Inches(0.25),"CLOSING — IF HE'S VAGUE",sz=14,bold=True,col=AMBER)
txt(sl,Inches(0.9),Inches(4.8),Inches(11.5),Inches(0.5),
    "\"No pressure. I'll send a summary of what MTCP produces and how a pilot would work.\n"
    "You can share it internally and we pick up when the timing's right.\"",sz=17,col=WHITE)

box(sl,Inches(0.6),Inches(5.65),Inches(12.1),Inches(1.0),border=RED)
txt(sl,Inches(0.9),Inches(5.75),Inches(11.5),Inches(0.25),"CLOSING — IF HE PUSHES FOR FREE STUFF / IP",sz=14,bold=True,col=RED)
txt(sl,Inches(0.9),Inches(6.05),Inches(11.5),Inches(0.5),
    "\"The pilot is designed to be the low-commitment proof of concept.\n"
    "That's the cleanest starting point for both sides.\"",sz=17,col=WHITE)
foot(sl,20,TOT)


# ═══════════════════════════════════════════
# 21  RESCUE PHRASES
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Rescue Phrases — If Things Go Sideways","Memorise 3 of these. They'll save you.")

rescues=[
    ("Lost your train of thought",
     "\"Let me come back to the key point — MTCP turns model behaviour into assurance evidence.\"",TEAL),
    ("Don't know the answer",
     "\"Good question. I want to give you a proper answer — let me follow up in the proposal.\"",PURPLE),
    ("He's going too deep into the engine",
     "\"I'd love to go deeper — and we can once we've scoped the pilot. For today, the key thing is what you'd receive.\"",AMBER),
    ("Awkward silence",
     "\"What are your initial thoughts on how this fits into what sirar is building?\"",GREEN),
    ("He's rushing / distracted",
     "\"Quick version: we tested 32 AI models, none passed, and we produce the evidence the EU AI Act requires.\n"
     "Shall I send you a pilot proposal?\"",PINK),
    ("He asks for free work",
     "\"The pilot is exactly that — low-commitment proof of concept. Cleanest starting point.\"",AMBER),
    ("He says \"we'll think about it\"",
     "\"Of course. I'll send the pilot proposal so you have something concrete. Is email best?\"",GREEN),
    ("You blank completely",
     "\"The core of it is: MTCP gives companies proof of whether their AI is safe. Not opinion — evidence.\n"
     "That's what we'd provide to sirar.\"  — This works no matter what the question was.",RED),
]
for i,(sit,say,col) in enumerate(rescues):
    y=Inches(1.1)+Inches(0.75)*i
    txt(sl,Inches(0.6),y,Inches(3.2),Inches(0.3),sit,sz=12,bold=True,col=col)
    txt(sl,Inches(4.0),y,Inches(8.5),Inches(0.65),say,sz=12,col=LIGHT)
foot(sl,21,TOT)


# ═══════════════════════════════════════════
# 22  AFTER THE CALL + CONFIDENCE
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"After the Call + One Last Thing","What to do next — and a reminder.")

box(sl,Inches(0.6),Inches(1.1),Inches(7.5),Inches(3.5),border=TEAL)
txt(sl,Inches(0.9),Inches(1.2),Inches(7),Inches(0.3),"AFTER THE CALL",sz=16,bold=True,col=TEAL)
after=[
    ("Within 1 hour","Write down everything he said. Key needs, concerns, language he used."),
    ("Same day","Thank-you message. \"I'll have the pilot proposal to you by [date].\""),
    ("Within 3 days","Send pilot proposal: scope, timeline, deliverables, success criteria, price."),
    ("Within 1 week","If no reply: one follow-up. \"Checking this landed.\""),
    ("If yes","Get it in writing. Even a confirming email = enough to start."),
    ("If silence","Don't chase more than twice. Follow up in a month."),
]
for i,(when,what) in enumerate(after):
    y=Inches(1.6)+Inches(0.42)*i
    txt(sl,Inches(0.9),y,Inches(1.6),Inches(0.3),when,sz=12,bold=True,col=GREEN)
    txt(sl,Inches(2.6),y,Inches(4.8),Inches(0.35),what,sz=12,col=LIGHT)

box(sl,Inches(8.4),Inches(1.1),Inches(4.3),Inches(3.5),fill=CARD2,border=GOLD)
txt(sl,Inches(8.7),Inches(1.2),Inches(3.8),Inches(0.3),"REMEMBER WHO YOU ARE",sz=16,bold=True,col=GOLD)
txt(sl,Inches(8.7),Inches(1.6),Inches(3.8),Inches(2.8),
    "You built something nobody else has.\n\n"
    "184,387 evaluations.\n"
    "32 models. 13 providers.\n"
    "5 published papers.\n"
    "6 experts validated in 48 hours.\n"
    "A sovereign stack forming around you.\n"
    "Theo credited you publicly.\n"
    "Mohamed wrote a formal spec for you.\n"
    "Anthropic invited you personally.\n\n"
    "He came to YOU.\n"
    "You've got this.",
    sz=13,col=LIGHT)

box(sl,Inches(0.6),Inches(5.0),Inches(12.1),Inches(1.5),border=PINK)
txt(sl,Inches(0.9),Inches(5.1),Inches(11.5),Inches(0.3),"YOUR ONE-LINERS (pick the one that feels natural and keep it loaded)",sz=14,bold=True,col=PINK)
liners=[
    "\"MTCP turns model behaviour into assurance evidence.\"",
    "\"We check whether correction is real or cosmetic.\"",
    "\"We identify which models can be governed at runtime and which will bleed through no matter what.\"",
    "\"No model we tested passed. That's the gap sirar can fill.\"",
]
for i,line in enumerate(liners):
    txt(sl,Inches(0.9),Inches(5.5)+Inches(0.32)*i,Inches(11.5),Inches(0.3),line,sz=14,bold=True,col=GREEN)
foot(sl,22,TOT)


# ═══════════════════════════════════════════
out=os.path.expanduser("~/Desktop/SIRAR_Call_Cheat_Sheet.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
