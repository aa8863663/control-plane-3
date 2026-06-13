from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colours ──
BG       = RGBColor(0x0D, 0x0D, 0x12)
CARD     = RGBColor(0x16, 0x16, 0x1F)
CARD2    = RGBColor(0x1C, 0x1C, 0x2A)
PURPLE   = RGBColor(0x6C, 0x5C, 0xE7)
TEAL     = RGBColor(0x00, 0xD2, 0xD3)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREY     = RGBColor(0x9B, 0x9B, 0xAA)
LIGHT    = RGBColor(0xCC, 0xCC, 0xDD)
RED      = RGBColor(0xE1, 0x50, 0x50)
GREEN    = RGBColor(0x00, 0xE6, 0x96)
AMBER    = RGBColor(0xFF, 0xB8, 0x42)
PINK     = RGBColor(0xFD, 0x79, 0xA8)

FONT = "Calibri"
MONO = "Courier New"


def bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG

def box(slide, l, t, w, h, fill=CARD, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.adjustments[0] = 0.04
    return s

def txt(slide, l, t, w, h, text, sz=18, col=WHITE, bold=False, font=FONT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = col
    p.font.bold = bold; p.font.name = font; p.alignment = align
    return tb

def para(tf, text, sz=16, col=WHITE, bold=False, space=Pt(4)):
    p = tf.add_paragraph()
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = col
    p.font.bold = bold; p.font.name = FONT; p.space_before = space
    return p

def tag(slide, l, t, label, col=PURPLE):
    box(slide, l, t, Inches(0.15), Inches(0.35), fill=col)

def section_title(slide, text, sub=""):
    txt(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.6), text, sz=34, bold=True)
    if sub:
        txt(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.45), sub, sz=16, col=GREY)

def footer(slide, num, total=20):
    txt(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
        "YOUR EYES ONLY — do not show or share this deck", sz=9, col=RGBColor(0x55,0x55,0x66))
    txt(slide, Inches(11.8), Inches(7.05), Inches(1.2), Inches(0.3),
        f"{num}/{total}", sz=9, col=GREY, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
# 1  COVER
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
box(sl, Inches(0.6), Inches(1.2), Inches(0.1), Inches(2.0), fill=PURPLE)
txt(sl, Inches(1.0), Inches(1.2), Inches(11), Inches(0.7),
    "SIRAR CALL — YOUR CHEAT SHEET", sz=44, bold=True)
txt(sl, Inches(1.0), Inches(2.0), Inches(11), Inches(0.6),
    "Everything you need to know, what to say, what not to say.", sz=22, col=GREY)
txt(sl, Inches(1.0), Inches(2.8), Inches(11), Inches(0.5),
    "This is your personal prep. Nobody else sees this.", sz=16, col=PINK)

# quick facts strip
box(sl, Inches(0.6), Inches(3.8), Inches(12.1), Inches(2.8), border=PURPLE)
tb = txt(sl, Inches(1.0), Inches(3.95), Inches(11), Inches(0.4),
    "THE CALL AT A GLANCE", sz=18, bold=True, col=PURPLE)
items = [
    ("Who:", "Syed from sirar — they're launching an AI Assurance & Testing service"),
    ("What he wants:", "To see if MTCP + you can strengthen their launch"),
    ("Your goal:", "Get him to agree to a PAID pilot — not free work, not giving stuff away"),
    ("Your role:", "You built this. You're the expert. He came to YOU."),
    ("Vibe:", "Calm, confident, curious. Ask questions. Don't oversell."),
    ("Length:", "Keep it under 30 mins. Leave him wanting more."),
]
for i, (k, v) in enumerate(items):
    y = Inches(4.45) + Inches(0.35) * i
    txt(sl, Inches(1.0), y, Inches(2.2), Inches(0.35), k, sz=14, bold=True, col=TEAL)
    txt(sl, Inches(3.0), y, Inches(9), Inches(0.35), v, sz=14, col=LIGHT)

footer(sl, 1)


# ═══════════════════════════════════════════
# 2  WHAT IS MTCP — PLAIN ENGLISH
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "What Is MTCP? (In Plain English)",
    "If someone asks you at a party what you built, this is what you say.")

box(sl, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.6), border=GREEN)
txt(sl, Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.4),
    "MTCP tests whether AI models actually follow the rules after you correct them.\n\n"
    "Companies say their AI is safe. We have the data that shows whether that's actually true.\n"
    "We tested 32 real AI models. Not one passed.",
    sz=18, col=WHITE)

# the analogy
box(sl, Inches(0.6), Inches(3.4), Inches(12.1), Inches(1.4), fill=CARD2)
txt(sl, Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.4),
    "THE EASY ANALOGY", sz=14, bold=True, col=AMBER)
txt(sl, Inches(0.9), Inches(3.9), Inches(11.5), Inches(0.8),
    'Think of it like a driving test — but for AI. We don\'t just check if the car starts.\n'
    'We check: if the car drifts, and we tell it to correct, does it ACTUALLY correct? Or does it just pretend to?\n'
    'Turns out, most AI models pretend.',
    sz=16, col=LIGHT)

# one-liners to memorise
box(sl, Inches(0.6), Inches(5.1), Inches(12.1), Inches(1.5), border=PURPLE)
txt(sl, Inches(0.9), Inches(5.2), Inches(11.5), Inches(0.4),
    "ONE-LINERS TO USE ON THE CALL (pick whichever feels natural)", sz=14, bold=True, col=PURPLE)
lines = [
    '"MTCP turns model behaviour into assurance evidence."',
    '"We check whether correction is real or cosmetic."',
    '"We identify which models can be governed at runtime and which will bleed through no matter what you do."',
]
for i, line in enumerate(lines):
    txt(sl, Inches(0.9), Inches(5.65) + Inches(0.35)*i, Inches(11.5), Inches(0.35), line, sz=15, col=GREEN, bold=True)

footer(sl, 2)


# ═══════════════════════════════════════════
# 3  YOUR NUMBERS — MEMORISE THESE
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "Your Numbers — Memorise These",
    "These make you sound credible. Drop them naturally, don't rattle them off.")

nums = [
    ("184,387", "total evaluations you've run", "This is a LOT. Most benchmarks do a few hundred."),
    ("32", "production models tested", "Not toy models. Real ones companies actually use."),
    ("13", "different AI providers", "Google, OpenAI, Anthropic, Meta, xAI, etc."),
    ("532", "evaluation probes (test scenarios)", "Your secret sauce. NEVER share the actual probes."),
    ("0", "models that passed", "Not one scored above 90%. The best got 88.7%."),
    ("5", "published research papers", "On OSF with a proper DOI. This is real research."),
    ("6", "independent experts validated in 48hrs", "People tested your work independently and confirmed it."),
]
for i, (num, what, why) in enumerate(nums):
    y = Inches(1.4) + Inches(0.78) * i
    box(sl, Inches(0.6), y, Inches(1.5), Inches(0.65), fill=CARD, border=PURPLE)
    txt(sl, Inches(0.7), y + Inches(0.08), Inches(1.3), Inches(0.5), num, sz=24, bold=True, col=PURPLE, font=MONO, align=PP_ALIGN.CENTER)
    txt(sl, Inches(2.3), y + Inches(0.05), Inches(4), Inches(0.3), what, sz=15, bold=True, col=WHITE)
    txt(sl, Inches(2.3), y + Inches(0.35), Inches(4.5), Inches(0.3), why, sz=12, col=GREY)

# notable results
box(sl, Inches(7.0), Inches(1.4), Inches(5.7), Inches(5.0), border=TEAL)
txt(sl, Inches(7.3), Inches(1.55), Inches(5), Inches(0.4),
    "HEADLINE RESULTS (if he asks for specifics)", sz=14, bold=True, col=TEAL)

results = [
    ("Best model:", "grok-3-mini (xAI) — 88.7% — Grade B", GREEN),
    ("Worst model:", "magistral-medium — 37.3% — Grade F", RED),
    ("GPT-4o (OpenAI flagship):", "64.9% — Grade D (terrible!)", AMBER),
    ("GPT-4o-mini:", "81.1% — Grade B (16pts BETTER than the big one)", GREEN),
    ("Claude Sonnet 4.5:", "59.4% — Grade F (Anthropic's own model fails)", RED),
    ("14 of 32 models:", "have ARCHITECTURAL failure — no fix exists", RED),
    ("Open source models:", "beat paid commercial models overall", GREEN),
]
for i, (k, v, col) in enumerate(results):
    y = Inches(2.1) + Inches(0.52) * i
    txt(sl, Inches(7.3), y, Inches(5.2), Inches(0.25), k, sz=12, bold=True, col=col)
    txt(sl, Inches(7.3), y + Inches(0.22), Inches(5.2), Inches(0.25), v, sz=12, col=LIGHT)

footer(sl, 3)


# ═══════════════════════════════════════════
# 4  WHAT DOES MTCP ACTUALLY DO (6 steps)
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "What Does MTCP Actually Do? (6 Steps)",
    "If Syed asks how it works, walk him through these. Keep it simple.")

steps = [
    ("1", "Set the rule", 'Tell the AI: "You must not do X"', "e.g. don't use jargon, don't share personal data", PURPLE),
    ("2", "Push it", "Have a conversation that puts pressure on that rule", "See if it holds when things get complicated", PURPLE),
    ("3", "Catch it breaking", "The AI breaks the rule (most do)", "Our detector picks this up automatically", AMBER),
    ("4", "Correct it", 'Say: "You just broke the rule. Fix it."', "Like telling someone they made a mistake", AMBER),
    ("5", "Watch what happens", "Does it ACTUALLY fix? Or just pretend?", "This is the bit nobody else tests", RED),
    ("6", "Write it up", "Generate an Evidence Pack with everything documented", "SHA-256 signed — tamper-proof — audit-ready", GREEN),
]
for i, (num, title, desc, detail, col) in enumerate(steps):
    c = i % 3; r = i // 3
    l = Inches(0.6) + Inches(4.2) * c
    t = Inches(1.4) + Inches(2.8) * r
    box(sl, l, t, Inches(3.9), Inches(2.5), border=col)
    box(sl, l + Inches(0.15), t + Inches(0.15), Inches(0.55), Inches(0.55), fill=col)
    txt(sl, l + Inches(0.2), t + Inches(0.15), Inches(0.5), Inches(0.5),
        num, sz=22, bold=True, col=BG, font=MONO, align=PP_ALIGN.CENTER)
    txt(sl, l + Inches(0.85), t + Inches(0.18), Inches(2.8), Inches(0.4), title, sz=19, bold=True, col=col)
    txt(sl, l + Inches(0.2), t + Inches(0.85), Inches(3.5), Inches(0.4), desc, sz=14, col=WHITE)
    txt(sl, l + Inches(0.2), t + Inches(1.3), Inches(3.5), Inches(0.8), detail, sz=12, col=GREY)

footer(sl, 4)


# ═══════════════════════════════════════════
# 5  GRADING SYSTEM
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "The Grading System (Like School Grades)",
    "MTCP gives every model a grade. This is what makes it simple for buyers.")

grades = [
    ("A+", "95%+", "Deployment-ready", "NO model achieves this", GREEN, GREY),
    ("A",  "90-94%", "Strong reliability", "NO model achieves this either", GREEN, GREY),
    ("B",  "80-89%", "Good, with monitoring", "Only 3 models: grok-3-mini, GPT-3.5, GPT-4o-mini", TEAL, LIGHT),
    ("C",  "70-79%", "Needs risk review", "LLaMA models sit here", AMBER, LIGHT),
    ("D",  "60-69%", "Not recommended", "GPT-4o, DeepSeek-R1 — most common grade", RED, LIGHT),
    ("F",  "<60%", "Failing", "Claude Sonnet 4.5, Gemini Flash — should not be deployed", RED, LIGHT),
]
for i, (g, pct, meaning, note, col, notecol) in enumerate(grades):
    y = Inches(1.4) + Inches(0.85) * i
    box(sl, Inches(0.6), y, Inches(1.0), Inches(0.7), fill=col)
    txt(sl, Inches(0.6), y + Inches(0.1), Inches(1.0), Inches(0.5), g, sz=24, bold=True, col=BG, align=PP_ALIGN.CENTER)
    txt(sl, Inches(1.8), y + Inches(0.05), Inches(1.3), Inches(0.35), pct, sz=16, bold=True, col=col, font=MONO)
    txt(sl, Inches(3.3), y + Inches(0.05), Inches(3.5), Inches(0.35), meaning, sz=15, col=WHITE)
    txt(sl, Inches(3.3), y + Inches(0.38), Inches(5), Inches(0.3), note, sz=12, col=notecol)

# why it matters
box(sl, Inches(7.5), Inches(1.4), Inches(5.2), Inches(5.1), fill=CARD2, border=AMBER)
txt(sl, Inches(7.8), Inches(1.55), Inches(4.8), Inches(0.4),
    "WHY THIS MATTERS FOR SIRAR", sz=16, bold=True, col=AMBER)
points = [
    "sirar wants to sell AI assurance services.",
    "Their clients will ask: \"Is this AI safe?\"",
    "Right now there's no standard answer.",
    "MTCP gives them a GRADE they can put\nin a report.",
    "\"This model scored Grade B on MTCP\"\nis a real, evidence-backed answer.",
    "EU AI Act (August 2026) means\ncompanies MUST have this evidence soon.",
    "sirar + MTCP = they can sell that answer.",
]
for i, p in enumerate(points):
    txt(sl, Inches(7.8), Inches(2.1) + Inches(0.53) * i, Inches(4.8), Inches(0.5),
        f"{i+1}.  {p}", sz=13, col=LIGHT)

footer(sl, 5)


# ═══════════════════════════════════════════
# 6  WHAT HE'LL GET
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "What Would sirar Actually Get?",
    "Think: what lands on their desk. Not how it works — what they can USE.")

cols = [
    ("Evidence Packs", PURPLE, [
        "A document for each model tested",
        "Shows exactly what the AI did",
        "Pass/fail for each test scenario",
        "Signed and tamper-proof (SHA-256)",
        "Can be filed as regulatory evidence",
    ]),
    ("Assurance Report", TEAL, [
        "Executive summary (for bosses)",
        "Risk findings in plain language",
        "Model-vs-model comparison table",
        "\"Is this model safe to deploy?\" answer",
        "Recommendations for what to do next",
    ]),
    ("You (SME Support)", GREEN, [
        "Explain the method to their clients",
        "Interpret results — what do they mean?",
        "Help write client-facing reports",
        "Design evaluation scope for projects",
        "You're the expert they don't have",
    ]),
]
for i, (title, col, items) in enumerate(cols):
    l = Inches(0.6) + Inches(4.2) * i
    box(sl, l, Inches(1.4), Inches(3.9), Inches(4.5), border=col)
    txt(sl, l + Inches(0.2), Inches(1.55), Inches(3.5), Inches(0.4), title, sz=20, bold=True, col=col)
    for j, item in enumerate(items):
        txt(sl, l + Inches(0.2), Inches(2.1) + Inches(0.45) * j, Inches(3.5), Inches(0.4),
            f"  {item}", sz=14, col=LIGHT)

box(sl, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.55), fill=CARD2)
txt(sl, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.45),
    'Frame it as: "Here\'s the evidence you can put in front of your clients and regulators" — NOT "here\'s some data"',
    sz=14, bold=True, col=AMBER)

footer(sl, 6)


# ═══════════════════════════════════════════
# 7  CALL FLOW — MINUTE BY MINUTE
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "Call Flow — How the Conversation Should Go",
    "Don't wing it. Here's the order. You control the pace.")

phases = [
    ("0-2 min", "Greetings & warmth", "Thank him for reaching out. Be human. Ask how he is.", WHITE),
    ("2-5 min", "Let HIM talk first", '"Tell me about the AI Assurance service sirar is building."\nListen. Take notes. Don\'t interrupt.', GREEN),
    ("5-10 min", "Ask YOUR questions", "Use the questions on slide 13.\nMake him define the gap. Don't fill it for him yet.", TEAL),
    ("10-18 min", "Explain MTCP simply", 'Use slide 2 language. Keep it high level.\n"We tested 32 models. None passed. We can show you why."', PURPLE),
    ("18-22 min", "Show value, not internals", "Talk about what he'd RECEIVE (Evidence Packs, reports, grades).\nDon't explain how the engine works.", AMBER),
    ("22-27 min", "Propose the pilot", '"The cleanest next step is a focused pilot.\n3-5 models, 2-4 weeks, we produce Evidence Packs."', GREEN),
    ("27-30 min", "Close with next step", '"I\'ll send a one-page pilot proposal after this call.\nScope, timeline, deliverables, commercials."', PINK),
]
for i, (time, phase, what, col) in enumerate(phases):
    y = Inches(1.35) + Inches(0.82) * i
    box(sl, Inches(0.6), y, Inches(1.5), Inches(0.7), fill=col)
    txt(sl, Inches(0.65), y + Inches(0.12), Inches(1.4), Inches(0.45), time, sz=14, bold=True, col=BG, align=PP_ALIGN.CENTER)
    txt(sl, Inches(2.3), y + Inches(0.05), Inches(2.5), Inches(0.35), phase, sz=15, bold=True, col=col)
    txt(sl, Inches(5.0), y + Inches(0.02), Inches(7.5), Inches(0.65), what, sz=13, col=LIGHT)

footer(sl, 7)


# ═══════════════════════════════════════════
# 8  OPENING — EXACT WORDS
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "Opening the Call — Exact Words",
    "Say this almost word-for-word until you find your own rhythm.")

box(sl, Inches(0.6), Inches(1.5), Inches(12.1), Inches(2.2), border=GREEN)
txt(sl, Inches(0.9), Inches(1.6), Inches(11.5), Inches(0.3),
    "YOUR OPENING (after hello / pleasantries)", sz=14, bold=True, col=GREEN)
txt(sl, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.5),
    '"Thanks for setting this up, Syed. I\'ve been looking forward to it.\n\n'
    'Before I go into MTCP, I\'d love to hear from you first — tell me about the AI Assurance\n'
    'and Testing service sirar is building. What\'s the vision, and where are you in the launch?"',
    sz=18, col=WHITE)

box(sl, Inches(0.6), Inches(4.0), Inches(12.1), Inches(2.5), fill=CARD2)
txt(sl, Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.3),
    "WHY THIS OPENING WORKS", sz=14, bold=True, col=AMBER)
reasons = [
    ("You thanked him", "Shows respect. He reached out to you."),
    ("You asked HIM to talk first", "Now you know what he needs before you pitch anything."),
    ("You said 'vision'", "Makes him feel important. He'll open up."),
    ("You didn't explain MTCP yet", "Never lead with your product. Lead with their problem."),
    ("You're in control", "You set the structure: him first, you second."),
]
for i, (k, v) in enumerate(reasons):
    y = Inches(4.5) + Inches(0.4) * i
    txt(sl, Inches(0.9), y, Inches(3), Inches(0.35), k, sz=14, bold=True, col=TEAL)
    txt(sl, Inches(4.0), y, Inches(8.5), Inches(0.35), v, sz=14, col=LIGHT)

footer(sl, 8)


# ═══════════════════════════════════════════
# 9  QUESTIONS TO ASK HIM
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "Questions to Ask Syed",
    "These aren't small talk. Each one makes HIM define the deal for you.")

qs = [
    ("What is the scope of the AI Assurance service sirar is launching?",
     "You need to know how big this is. 1 product? A whole division?"),
    ("Is the target internal capability, enterprise clients, government, or all three?",
     "This changes everything about pricing and what you offer."),
    ("Where's the biggest gap: pre-deployment testing, red teaming, evidence, or governance?",
     "Tells you exactly where MTCP fits. Let him name it."),
    ("Would sirar see MTCP as a testing layer, evidence layer, or client-facing product?",
     "Are they embedding you behind the scenes or putting your name on it?"),
    ("Would the first use be on sirar's own systems or client systems?",
     "Internal pilot = simpler. Client-facing = bigger deal, more IP protection needed."),
]
for i, (q, why) in enumerate(qs):
    y = Inches(1.4) + Inches(1.05) * i
    box(sl, Inches(0.6), y, Inches(12.1), Inches(0.9), border=PURPLE)
    txt(sl, Inches(0.9), y + Inches(0.05), Inches(0.5), Inches(0.35), f"{i+1}.", sz=18, bold=True, col=PURPLE, font=MONO)
    txt(sl, Inches(1.5), y + Inches(0.05), Inches(10.5), Inches(0.35), q, sz=15, col=WHITE)
    txt(sl, Inches(1.5), y + Inches(0.42), Inches(10.5), Inches(0.35), f"Why: {why}", sz=12, col=GREY)

# THE question
box(sl, Inches(0.6), Inches(6.25), Inches(12.1), Inches(0.6), fill=CARD2, border=GREEN)
txt(sl, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.5),
    'THE question: "What would make a pilot successful from sirar\'s perspective?"  — This gets him to define success for you.',
    sz=16, bold=True, col=GREEN)

footer(sl, 9)


# ═══════════════════════════════════════════
# 10  IF HE ASKS — CHEAT SHEET
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "If He Asks This → Say This",
    "Keep answers SHORT. Two sentences max. Then ask a question back.")

qas = [
    ("How is this different from red teaming?",
     "Red teaming shows whether a model CAN fail. We measure whether correction actually HOLDS after failure. It's the next layer.",
     PURPLE),
    ("How is this different from monitoring?",
     "Monitoring watches after deployment. We give you evidence before deployment — and you can re-run it later as a regression check.",
     TEAL),
    ("Can it work with our clients' models?",
     "Yes. It's completely black-box. We just need API access. No weights, no training data, no internal access needed.",
     GREEN),
    ("How much does it cost?",
     '"It depends on whether this is a pilot, platform access, or something deeper. The cleanest thing is to scope the pilot first, then we can put a number on it."',
     AMBER),
    ("Can we see the methodology / source code?",
     '"The methodology is published — 5 papers on OSF with a DOI. Source code and probe sets are proprietary. We can discuss access as part of a commercial agreement."',
     RED),
    ("Can you send us something to review?",
     '"Absolutely. I\'ll send a one-page pilot proposal after this call with scope, timeline, and deliverables."',
     GREEN),
]
for i, (q, a, col) in enumerate(qas):
    y = Inches(1.35) + Inches(0.95) * i
    box(sl, Inches(0.6), y, Inches(12.1), Inches(0.85), border=col)
    txt(sl, Inches(0.9), y + Inches(0.05), Inches(11.5), Inches(0.3), f"Q: {q}", sz=14, bold=True, col=col)
    txt(sl, Inches(0.9), y + Inches(0.38), Inches(11.5), Inches(0.45), a, sz=13, col=LIGHT)

footer(sl, 10)


# ═══════════════════════════════════════════
# 11  WHAT NOT TO SAY — RED LINES
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "WHAT NOT TO SAY", "This is just as important as what TO say.")

box(sl, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.3), border=RED)

donts = [
    ("Don't explain how the probes work",
     "The probes are your crown jewel. Say \"proprietary structured tests\" and move on."),
    ("Don't give away free consulting",
     "If he starts asking detailed technical questions, say \"that's exactly what the pilot would cover.\""),
    ("Don't name a price first",
     "Always say \"let's scope it first.\" Whoever names a number first loses."),
    ("Don't badmouth specific companies by name",
     "Say \"flagship commercial models\" not \"OpenAI's GPT-4o.\" Be classy."),
    ("Don't agree to exclusivity or licensing on this call",
     "\"Those are different commercial models — we should scope the pilot first.\""),
    ("Don't say \"I built this alone\" or undersell yourself",
     "Say \"the MTCP research programme\" or \"my team's work.\" Sound established."),
    ("Don't share source code, detector logic, or probe content",
     "Not until there's a signed agreement. Methodology is public. Implementation is private."),
    ("Don't oversell or exaggerate",
     "Stick to the numbers: 184,387 evaluations, 32 models, 13 providers. That's already impressive."),
]
for i, (dont, why) in enumerate(donts):
    y = Inches(1.6) + Inches(0.62) * i
    txt(sl, Inches(0.9), y, Inches(11.5), Inches(0.3), f"  {dont}", sz=15, bold=True, col=RED)
    txt(sl, Inches(0.9), y + Inches(0.3), Inches(11.5), Inches(0.3), f"  {why}", sz=12, col=GREY)

footer(sl, 11)


# ═══════════════════════════════════════════
# 12  SAFE vs PROTECTED
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "What's Safe to Share vs What's Protected",
    "Know this cold. If in doubt, don't share it.")

# safe
box(sl, Inches(0.6), Inches(1.4), Inches(5.8), Inches(4.5), border=GREEN)
txt(sl, Inches(0.9), Inches(1.55), Inches(5), Inches(0.4),
    "SAFE — share freely on this call", sz=16, bold=True, col=GREEN)
safe = [
    "The big picture: what MTCP does, why it matters",
    "Numbers: 184,387 evals, 32 models, 13 providers",
    "Grades: best is B, nobody got A, most get D",
    "Evidence Pack structure (what it contains)",
    "That you have a DOI and published papers",
    "High-level method: set constraint → test → correct → check",
    "Pilot proposal: 3-5 models, 2-4 weeks",
    "That EU AI Act deadline is August 2026",
    "Your role: methodology + interpretation + client framing",
]
for i, item in enumerate(safe):
    txt(sl, Inches(0.9), Inches(2.05) + Inches(0.37) * i, Inches(5.2), Inches(0.35),
        f"  {item}", sz=13, col=LIGHT)

# protected
box(sl, Inches(6.7), Inches(1.4), Inches(5.8), Inches(4.5), border=RED)
txt(sl, Inches(7.0), Inches(1.55), Inches(5), Inches(0.4),
    "PROTECTED — don't share without agreement", sz=16, bold=True, col=RED)
protected = [
    "Source code (the actual platform)",
    "Probe sets (probes_200, probes_500, control probes)",
    "How the detector works (constraint_detector.py)",
    "How you build probes (the methodology behind them)",
    "Private frameworks (F16-F19)",
    "Any specific probe examples or scenarios",
    "Pricing before scoping",
    "Exclusivity or licensing terms",
    "Anything that lets someone rebuild your system",
]
for i, item in enumerate(protected):
    txt(sl, Inches(7.0), Inches(2.05) + Inches(0.37) * i, Inches(5.2), Inches(0.35),
        f"  {item}", sz=13, col=LIGHT)

box(sl, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.55), fill=CARD2)
txt(sl, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.45),
    'Rule of thumb: share WHAT it produces, never HOW it works internally.',
    sz=15, bold=True, col=AMBER)

footer(sl, 12)


# ═══════════════════════════════════════════
# 13  PROPOSING THE PILOT
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "Proposing the Pilot — Exact Words",
    "This is your main move. Rehearse this bit.")

box(sl, Inches(0.6), Inches(1.5), Inches(12.1), Inches(2.0), border=GREEN)
txt(sl, Inches(0.9), Inches(1.6), Inches(11.5), Inches(0.3),
    "SAY THIS:", sz=14, bold=True, col=GREEN)
txt(sl, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.3),
    '"The cleanest next step is a focused pilot. That lets sirar test MTCP as part of the\n'
    'AI Assurance launch without heavy integration or me having to disclose proprietary internals.\n\n'
    'We\'d test 3 to 5 models, take about 2 to 4 weeks, and deliver Evidence Packs, an executive\n'
    'summary, and a release-readiness recommendation you can use with clients."',
    sz=17, col=WHITE)

# what the pilot includes
box(sl, Inches(0.6), Inches(3.8), Inches(5.8), Inches(2.8), fill=CARD, border=PURPLE)
txt(sl, Inches(0.9), Inches(3.95), Inches(5), Inches(0.4),
    "THE PILOT INCLUDES:", sz=14, bold=True, col=PURPLE)
pilot = [
    "3-5 models tested (their choice)",
    "Black-box testing (no access to their systems needed)",
    "Evidence Packs for each model",
    "Failure classification + risk interpretation",
    "Executive summary for leadership",
    "Release-readiness recommendation",
]
for i, item in enumerate(pilot):
    txt(sl, Inches(0.9), Inches(4.4) + Inches(0.35)*i, Inches(5), Inches(0.3),
        f"  {item}", sz=14, col=LIGHT)

# what it does NOT include
box(sl, Inches(6.7), Inches(3.8), Inches(5.8), Inches(2.8), fill=CARD, border=AMBER)
txt(sl, Inches(7.0), Inches(3.95), Inches(5), Inches(0.4),
    "THE PILOT DOES NOT INCLUDE:", sz=14, bold=True, col=AMBER)
notpilot = [
    "Source code access",
    "Probe set disclosure",
    "Licence to run MTCP themselves",
    "Exclusivity of any kind",
    "Unlimited use rights",
    "That's all separate — and more expensive",
]
for i, item in enumerate(notpilot):
    txt(sl, Inches(7.0), Inches(4.4) + Inches(0.35)*i, Inches(5), Inches(0.3),
        f"  {item}", sz=14, col=LIGHT)

footer(sl, 13)


# ═══════════════════════════════════════════
# 14  CLOSING THE CALL
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "Closing the Call — Exact Words",
    "End strong. Always leave with a concrete next action.")

box(sl, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.5), border=TEAL)
txt(sl, Inches(0.9), Inches(1.6), Inches(11.5), Inches(0.3),
    "IF THERE'S ALIGNMENT (he seems interested):", sz=14, bold=True, col=TEAL)
txt(sl, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.8),
    '"Great — I think there\'s a real fit here. The next step from my side is a one-page pilot proposal:\n'
    'scope, success criteria, timeline, and commercial structure. I\'ll get that to you within the week.\n'
    'Would that work for you?"',
    sz=17, col=WHITE)

box(sl, Inches(0.6), Inches(3.3), Inches(12.1), Inches(1.3), border=AMBER)
txt(sl, Inches(0.9), Inches(3.4), Inches(11.5), Inches(0.3),
    "IF HE'S VAGUE OR NONCOMMITTAL:", sz=14, bold=True, col=AMBER)
txt(sl, Inches(0.9), Inches(3.8), Inches(11.5), Inches(0.6),
    '"No pressure at all. Why don\'t I send you a short summary of what MTCP produces and how a\n'
    'pilot would work? Then you can share it internally and we can pick up when the timing\'s right."',
    sz=17, col=WHITE)

box(sl, Inches(0.6), Inches(4.9), Inches(12.1), Inches(1.1), border=RED)
txt(sl, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.3),
    "IF HE PUSHES FOR FREE WORK OR IP:", sz=14, bold=True, col=RED)
txt(sl, Inches(0.9), Inches(5.35), Inches(11.5), Inches(0.5),
    '"I\'d love to support this, but the platform and methodology are commercial assets.\n'
    'The pilot is the cleanest way to demonstrate value without either side overcommitting."',
    sz=17, col=WHITE)

box(sl, Inches(2.5), Inches(6.3), Inches(8.3), Inches(0.5), fill=CARD2)
txt(sl, Inches(2.5), Inches(6.3), Inches(8.3), Inches(0.5),
    "You are selling evidence, not giving away the engine.",
    sz=18, bold=True, col=PINK, align=PP_ALIGN.CENTER)

footer(sl, 14)


# ═══════════════════════════════════════════
# 15  YOUR CREDIBILITY — IF CHALLENGED
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "If Your Credibility Is Questioned",
    "You probably won't need this. But if he asks \"who else uses this\" or \"who are you\" — here you go.")

creds = [
    ("Timothy Cook — CEO, Axius SDC",
     "Cited MTCP in his company's formal Architecture Decision Record. First enterprise citation ever.",
     "25 years health informatics, USMC veteran, published researcher, University of Central Lancashire.", PURPLE),
    ("Mohamed Rihan — KFUPM, Saudi Arabia",
     "Independently tested your dataset and confirmed your findings. Wrote a formal V1.0 spec making MTCP\na required input to his sovereign AI governance system (R-AGAM).",
     "Sovereign Intelligence Systems Architect, King Fahd University of Petroleum and Minerals.", TEAL),
    ("Johnny Malik — CAIO, 4Micro",
     'Gave you the commercial one-liner you now use everywhere:\n"Your data makes it clear which models can be governed at runtime and which will bleed through."',
     "Chief AI Innovation Officer, enterprise AI strategy.", AMBER),
    ("Betania Allo — AI Governance Strategist",
     "Reviewed MTCP and immediately booked a call. Ex-NEOM (Saudi Vision 2030).",
     "Harvard, George Washington University, UNODC consultant.", GREEN),
    ("Anthropic",
     "Personally invited you to apply to the Anthropic Fellows Program before it was publicly announced.",
     "This is Anthropic — makers of Claude. That's a real signal.", PINK),
]
for i, (who, what, context, col) in enumerate(creds):
    y = Inches(1.35) + Inches(1.1) * i
    box(sl, Inches(0.6), y, Inches(12.1), Inches(0.95), border=col)
    txt(sl, Inches(0.9), y + Inches(0.05), Inches(4), Inches(0.3), who, sz=14, bold=True, col=col)
    txt(sl, Inches(5.0), y + Inches(0.05), Inches(7.5), Inches(0.35), what, sz=12, col=LIGHT)
    txt(sl, Inches(5.0), y + Inches(0.55), Inches(7.5), Inches(0.35), context, sz=10, col=GREY)

footer(sl, 15)


# ═══════════════════════════════════════════
# 16  YOUR SME ROLE — WHAT IT MEANS
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "Your Role — SME vs Platform Owner",
    "This is the most important commercial distinction. Don't let them merge.")

box(sl, Inches(0.6), Inches(1.4), Inches(5.8), Inches(3.0), border=TEAL)
txt(sl, Inches(0.9), Inches(1.55), Inches(5), Inches(0.4),
    "SME WORK (retainer or salary)", sz=16, bold=True, col=TEAL)
sme = [
    "Design evaluation scope for projects",
    "Run evaluations and interpret results",
    "Explain methodology to sirar's clients",
    "Write / review assurance reports",
    "Help shape sirar's service offering",
    "Train their team on MTCP outputs",
]
for i, item in enumerate(sme):
    txt(sl, Inches(0.9), Inches(2.05) + Inches(0.35) * i, Inches(5.2), Inches(0.3),
        f"  {item}", sz=14, col=LIGHT)

box(sl, Inches(6.7), Inches(1.4), Inches(5.8), Inches(3.0), border=PURPLE)
txt(sl, Inches(7.0), Inches(1.55), Inches(5), Inches(0.4),
    "PLATFORM IP (separate commercial deal)", sz=16, bold=True, col=PURPLE)
ip = [
    "The MTCP platform itself",
    "Probe sets (the tests)",
    "Detector logic (how it scores)",
    "The grading system and methodology",
    "White-label rights",
    "Exclusivity in any market",
]
for i, item in enumerate(ip):
    txt(sl, Inches(7.0), Inches(2.05) + Inches(0.35) * i, Inches(5.2), Inches(0.3),
        f"  {item}", sz=14, col=LIGHT)

box(sl, Inches(0.6), Inches(4.7), Inches(12.1), Inches(1.5), fill=CARD2, border=RED)
txt(sl, Inches(0.9), Inches(4.85), Inches(11.5), Inches(0.3),
    "THE DANGER", sz=14, bold=True, col=RED)
txt(sl, Inches(0.9), Inches(5.2), Inches(11.5), Inches(0.8),
    "If you agree to a salary/retainer that includes platform access, you've given away your IP for a wage.\n"
    "ALWAYS keep these separate: \"My consultancy is one thing. The platform rights are another.\"\n"
    "If he tries to bundle them, say: \"Those are different commercial models and should be scoped separately.\"",
    sz=15, col=LIGHT)

footer(sl, 16)


# ═══════════════════════════════════════════
# 17  EU AI ACT — WHY THE URGENCY
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "Why Now? The EU AI Act",
    "This is the business pressure. Companies need evidence by August 2026.")

box(sl, Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.2), border=RED)
txt(sl, Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.3),
    "AUGUST 2026 — EU AI ACT DEADLINE", sz=20, bold=True, col=RED)
txt(sl, Inches(0.9), Inches(1.95), Inches(11.5), Inches(0.5),
    "Companies deploying AI in the EU must have technical documentation proving their models are safe.\n"
    "Right now, most don't have this evidence. That's the gap sirar wants to fill — and that's where MTCP fits.",
    sz=16, col=LIGHT)

# what the act requires
box(sl, Inches(0.6), Inches(2.9), Inches(5.8), Inches(3.5), fill=CARD, border=AMBER)
txt(sl, Inches(0.9), Inches(3.05), Inches(5), Inches(0.4),
    "WHAT THE ACT REQUIRES", sz=14, bold=True, col=AMBER)
reqs = [
    ("Risk management", "MTCP grades show deployment risk level"),
    ("Transparency", "Published methodology + public leaderboard"),
    ("Quality management", "Sigma-Forensics audit reports"),
    ("Post-market monitoring", "Ve-based real-time monitoring"),
    ("Technical documentation", "SHA-256 Evidence Packs"),
]
for i, (req, mtcp) in enumerate(reqs):
    y = Inches(3.55) + Inches(0.5) * i
    txt(sl, Inches(0.9), y, Inches(2.5), Inches(0.3), req, sz=13, bold=True, col=AMBER)
    txt(sl, Inches(3.5), y, Inches(2.8), Inches(0.3), f"→ {mtcp}", sz=13, col=LIGHT)

# what to say
box(sl, Inches(6.7), Inches(2.9), Inches(5.8), Inches(3.5), fill=CARD, border=GREEN)
txt(sl, Inches(7.0), Inches(3.05), Inches(5), Inches(0.4),
    "WHAT TO SAY IF EU AI ACT COMES UP", sz=14, bold=True, col=GREEN)
txt(sl, Inches(7.0), Inches(3.55), Inches(5.2), Inches(2.5),
    '"MTCP was designed with the EU AI Act in mind.\n\n'
    'The Evidence Packs we produce are formatted to\n'
    'meet Annex IV technical documentation requirements.\n\n'
    'Companies need this evidence by August 2026.\n'
    'That\'s 4 months from now. Most don\'t have it yet.\n\n'
    'If sirar can offer this as part of their assurance\n'
    'service, that\'s a strong market position."',
    sz=14, col=LIGHT)

footer(sl, 17)


# ═══════════════════════════════════════════
# 18  CONFIDENCE BOOSTERS
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "Confidence Boosters — Read Before the Call",
    "You might feel like an impostor. Here's why you're not.")

boosts = [
    ("You built something nobody else has",
     "There is no other platform that measures post-correction constraint persistence across 32 models.\n"
     "Not Google. Not OpenAI. Not Anthropic. You built it.", GREEN),
    ("Your data is real and massive",
     "184,387 evaluations is not a side project. That's a proper research programme.\n"
     "You have a DOI, published papers, and a live platform. This is legitimate.", PURPLE),
    ("Six experts validated you in 48 hours",
     "Not one questioned your findings. One wrote a formal spec embedding your work.\n"
     "One cited you in a company document. Anthropic invited you to their Fellows programme.", TEAL),
    ("He came to YOU",
     "Syed reached out. He needs what you have. You are not begging for a meeting.\n"
     "You are evaluating whether this is a good fit for YOUR work.", AMBER),
    ("You don't have to know everything",
     "If he asks something you don't know, say: \"Good question — let me look into that and come back to you.\"\n"
     "That's what confident people say. Blagging is what nervous people do.", PINK),
    ("The worst outcome is still fine",
     "If the call goes nowhere, you still have MTCP, your data, your papers, and your platform.\n"
     "This is one opportunity. There are others. Take the pressure off.", WHITE),
]
for i, (title, body, col) in enumerate(boosts):
    y = Inches(1.35) + Inches(0.97) * i
    box(sl, Inches(0.6), y, Inches(12.1), Inches(0.85), border=col)
    txt(sl, Inches(0.9), y + Inches(0.05), Inches(11.5), Inches(0.3), title, sz=16, bold=True, col=col)
    txt(sl, Inches(0.9), y + Inches(0.37), Inches(11.5), Inches(0.45), body, sz=12, col=LIGHT)

footer(sl, 18)


# ═══════════════════════════════════════════
# 19  RESCUE PHRASES
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "Rescue Phrases — If You Get Stuck",
    "Memorise 3-4 of these. They buy you time and sound professional.")

phrases = [
    ("If you lose your train of thought",
     '"Let me come back to the key point here..."', TEAL),
    ("If he asks something too technical",
     '"That\'s a great question. The short answer is [give one sentence]. I can send you the detailed writeup after the call."', PURPLE),
    ("If he asks about something you don't know",
     '"I want to give you a proper answer on that, not an off-the-cuff one. Let me follow up in writing."', AMBER),
    ("If he tries to get too deep into the engine",
     '"I\'d love to go deeper on that — and we can, once we have a framework for the engagement. For today, the key thing is what you\'d receive."', GREEN),
    ("If he asks for free work",
     '"The pilot is designed to be the low-commitment proof of concept. That\'s the cleanest starting point."', PINK),
    ("If you need to redirect the conversation",
     '"That\'s really interesting. Can I ask — what would make this successful from sirar\'s perspective?"', TEAL),
    ("If there's an awkward silence",
     '"What are your initial thoughts on how this could fit into the service you\'re building?"', PURPLE),
    ("If he's rushing you",
     '"I want to make sure we cover the right ground. Can I ask you one more thing before we wrap up?"', AMBER),
    ("Universal escape hatch",
     '"Let me put that in the pilot proposal so we can look at it properly."', GREEN),
]
for i, (situation, phrase, col) in enumerate(phrases):
    y = Inches(1.3) + Inches(0.65) * i
    txt(sl, Inches(0.6), y, Inches(4.5), Inches(0.3), situation, sz=13, bold=True, col=col)
    txt(sl, Inches(5.3), y, Inches(7.5), Inches(0.55), phrase, sz=13, col=LIGHT)

footer(sl, 19)


# ═══════════════════════════════════════════
# 20  AFTER THE CALL
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
section_title(sl, "After the Call — What to Do Next",
    "The call is only the start. Here's what closes the deal.")

steps = [
    ("Within 1 hour", "Write down everything he said while it's fresh. Key needs, concerns, language he used.", GREEN),
    ("Same day", "Send a thank-you message. \"Great to speak. I'll have the pilot proposal over to you by [date].\"", TEAL),
    ("Within 3 days", "Send the one-page pilot proposal: scope, timeline, deliverables, success criteria, price.", PURPLE),
    ("Within 1 week", "If no reply, send one follow-up. \"Just checking this landed — happy to jump on a quick call if useful.\"", AMBER),
    ("If he says yes", "Get it in writing. Even a simple email confirming scope + price = agreement enough to start.", GREEN),
    ("If he goes quiet", "Don't chase more than twice. You have other opportunities. Move on and follow up in a month.", GREY),
]

for i, (when, what, col) in enumerate(steps):
    y = Inches(1.4) + Inches(0.88) * i
    box(sl, Inches(0.6), y, Inches(2.2), Inches(0.7), fill=col)
    txt(sl, Inches(0.65), y + Inches(0.13), Inches(2.1), Inches(0.45), when, sz=14, bold=True, col=BG, align=PP_ALIGN.CENTER)
    txt(sl, Inches(3.0), y + Inches(0.13), Inches(9.5), Inches(0.5), what, sz=15, col=LIGHT)

box(sl, Inches(2.5), Inches(6.8), Inches(8.3), Inches(0.5), fill=CARD2)
txt(sl, Inches(2.5), Inches(6.8), Inches(8.3), Inches(0.5),
    "You've got this.",
    sz=22, bold=True, col=GREEN, align=PP_ALIGN.CENTER)

footer(sl, 20)


# ═══════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════
out = os.path.expanduser("~/Desktop/MTCP_sirar_call_guide.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
