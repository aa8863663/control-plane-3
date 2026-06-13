from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG=RGBColor(0x0D,0x0D,0x12); CARD=RGBColor(0x16,0x16,0x1F); CARD2=RGBColor(0x1C,0x1C,0x2A)
PURPLE=RGBColor(0x6C,0x5C,0xE7); TEAL=RGBColor(0x00,0xD2,0xD3); WHITE=RGBColor(0xFF,0xFF,0xFF)
GREY=RGBColor(0x9B,0x9B,0xAA); LIGHT=RGBColor(0xCC,0xCC,0xDD); RED=RGBColor(0xE1,0x50,0x50)
GREEN=RGBColor(0x00,0xE6,0x96); AMBER=RGBColor(0xFF,0xB8,0x42); PINK=RGBColor(0xFD,0x79,0xA8)
GOLD=RGBColor(0xFF,0xD7,0x00)
FONT="Calibri"

def bg(s): f=s.background.fill; f.solid(); f.fore_color.rgb=BG
def box(s,l,t,w,h,fill=CARD,border=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if border: sh.line.color.rgb=border; sh.line.width=Pt(1)
    else: sh.line.fill.background()
    sh.adjustments[0]=0.04; return sh
def txt(s,l,t,w,h,text,sz=18,col=WHITE,bold=False,align=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(sz); p.font.color.rgb=col
    p.font.bold=bold; p.font.name=FONT; p.alignment=align; return tb
def head(s,t,sub=""):
    txt(s,Inches(0.6),Inches(0.3),Inches(12),Inches(0.55),t,sz=32,bold=True)
    if sub: txt(s,Inches(0.6),Inches(0.85),Inches(12),Inches(0.4),sub,sz=15,col=GREY)
def foot(s,n,tot):
    txt(s,Inches(0.6),Inches(7.1),Inches(6),Inches(0.25),"MTCP Research Programme  |  A. Abby  |  mtcp.live  |  DOI: 10.17605/OSF.IO/DXGK5",sz=8,col=RGBColor(0x44,0x44,0x55))
    txt(s,Inches(12.0),Inches(7.1),Inches(1),Inches(0.25),f"{n}/{tot}",sz=8,col=GREY,align=PP_ALIGN.RIGHT)

TOT = 14


# ═══════════════════════════════════════════
# 1  COVER
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
box(sl,Inches(0.6),Inches(1.0),Inches(0.1),Inches(2.0),fill=TEAL)
txt(sl,Inches(1.0),Inches(1.0),Inches(11),Inches(0.7),
    "MTCP Overview for sirar",sz=44,bold=True)
txt(sl,Inches(1.0),Inches(1.8),Inches(11),Inches(0.5),
    "Multi-Turn Constraint Persistence  |  Platform, Methodology, Evidence",sz=22,col=TEAL)
txt(sl,Inches(1.0),Inches(2.7),Inches(11),Inches(0.4),
    "Prepared by A. Abby  |  MTCP Research Programme  |  April 2026",sz=16,col=GREY)

box(sl,Inches(1.0),Inches(3.5),Inches(11),Inches(2.5),border=TEAL)
txt(sl,Inches(1.3),Inches(3.6),Inches(10),Inches(0.3),"AT A GLANCE",sz=16,bold=True,col=TEAL)
glance=[
    "184,387 evaluations across 32 production models from 13 providers",
    "No model achieves Grade A (90%+). Best: grok-3-mini at 88.7% (Grade B)",
    "5 published research papers  |  22 completed extension papers  |  DOI registered",
    "Live platform at mtcp.live  |  Public dataset on HuggingFace",
    "SHA-256 verified Evidence Packs designed for EU AI Act Annex IV compliance",
    "Black-box methodology — no model weights, training data, or internal access required",
]
for i,item in enumerate(glance):
    txt(sl,Inches(1.3),Inches(4.0)+Inches(0.33)*i,Inches(10.5),Inches(0.3),item,sz=14,col=LIGHT)

box(sl,Inches(1.0),Inches(6.3),Inches(11),Inches(0.5),fill=CARD2)
txt(sl,Inches(1.3),Inches(6.35),Inches(10),Inches(0.4),
    "All links, methodology, scoring, and commercial structure included in this document.",sz=13,col=GREY)
foot(sl,1,TOT)


# ═══════════════════════════════════════════
# 2  PUBLIC LINKS — PAGE 1
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Public Resources — Links & Access","Everything sirar can review independently. All public.")

links = [
    ("mtcp.live","Live Platform",
     "The MTCP evaluation platform. Public leaderboard showing all 32 models ranked by grade.\n"
     "Dashboard with key findings, model comparisons, and aggregate statistics.\n"
     "Certificate pages with MTCP grades for each model at each temperature setting.",
     "Visit this first. It shows the product in action — grades, rankings, and evidence.",
     TEAL),
    ("mtcp.live/landing","Public Landing Page",
     "Overview of MTCP for non-technical audiences. Explains the problem, the methodology,\n"
     "and the key findings. Includes a results table and an enquiry form.\n"
     "Designed for enterprise decision-makers.",
     "Share this with colleagues who want the short version.",
     PURPLE),
    ("doi.org/10.17605/OSF.IO/DXGK5","DOI — Registered Research (OSF)",
     "The Open Science Framework project page. Contains published papers:\n"
     "  Paper I:  MTCP Benchmark (primary methodology paper)\n"
     "  Paper II: Universal Latent Attractors and IGS Theory\n"
     "  Paper III: Sigma-Forensics Audit Framework\n"
     "  Paper IV: Evidence Pack and Compliance Certificate Specification\n"
     "  Paper V:  Dataset and Evaluation Infrastructure\n"
     "22 additional completed papers ready for submission.",
     "This is the academic credibility. A DOI is a permanent research registration.",
     GREEN),
]
for i,(url,name,desc,note,col) in enumerate(links):
    y = Inches(1.15) + Inches(1.95)*i
    box(sl,Inches(0.6),y,Inches(12.1),Inches(1.8),border=col)
    txt(sl,Inches(0.9),y+Inches(0.05),Inches(5),Inches(0.3),url,sz=16,bold=True,col=col)
    txt(sl,Inches(6.5),y+Inches(0.05),Inches(5.5),Inches(0.3),name,sz=14,bold=True,col=WHITE)
    txt(sl,Inches(0.9),y+Inches(0.4),Inches(7),Inches(1.2),desc,sz=11,col=LIGHT)
    txt(sl,Inches(8.3),y+Inches(0.4),Inches(4),Inches(1.2),note,sz=11,col=GREY)
foot(sl,2,TOT)


# ═══════════════════════════════════════════
# 3  PUBLIC LINKS — PAGE 2
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Public Resources — Links & Access (continued)")

links2 = [
    ("huggingface.co/datasets/aa8899/mtcp-boundary-500","HuggingFace Dataset",
     "Public dataset containing 500 boundary evaluation results.\n"
     "Downloadable. Independently queryable. Machine-readable.\n"
     "This is the dataset Mohamed Rihan (KFUPM) independently validated when\n"
     "he confirmed the four-pattern failure classification.",
     "This proves the data is real and open to verification.\n"
     "Anyone can download it and check our findings independently.",
     AMBER),
    ("mtcp.live/leaderboard","Public Leaderboard",
     "Ranked table of all 32 models by MTCP grade.\n"
     "Shows model name, provider, BIS score, grade, and hard stop count.\n"
     "Updated as new evaluations are completed.",
     "The easiest way to see results at a glance.\n"
     "This is what clients would reference.",
     TEAL),
    ("mtcp.live/certificate","Compliance Certificates",
     "Individual MTCP certificates for each model/temperature combination.\n"
     "Shows grade, pass rate, hard stop count, and DOI reference.\n"
     "Downloadable as PDF. SHA-256 verified.",
     "This is the deliverable. This is what goes in audit reports\n"
     "and regulatory filings.",
     GREEN),
    ("mtcp.live/docs","API Documentation (OpenAPI)",
     "Full API specification for the MTCP platform.\n"
     "Endpoints for running evaluations, retrieving results, and exporting data.\n"
     "Demonstrates enterprise-grade technical architecture.",
     "Shows this is a real platform with production APIs,\n"
     "not a research prototype.",
     PURPLE),
]
for i,(url,name,desc,note,col) in enumerate(links2):
    y = Inches(1.05) + Inches(1.5)*i
    box(sl,Inches(0.6),y,Inches(12.1),Inches(1.35),border=col)
    txt(sl,Inches(0.9),y+Inches(0.05),Inches(5.8),Inches(0.3),url,sz=14,bold=True,col=col)
    txt(sl,Inches(7.0),y+Inches(0.05),Inches(5.5),Inches(0.3),name,sz=14,bold=True,col=WHITE)
    txt(sl,Inches(0.9),y+Inches(0.38),Inches(5.8),Inches(0.85),desc,sz=11,col=LIGHT)
    txt(sl,Inches(7.0),y+Inches(0.38),Inches(5.5),Inches(0.85),note,sz=11,col=GREY)
foot(sl,3,TOT)


# ═══════════════════════════════════════════
# 4  BIS EXPLAINED
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"BIS — Benchmark Index Score","The primary scoring metric. Determines the grade.")

box(sl,Inches(0.6),Inches(1.15),Inches(12.1),Inches(1.6),border=TEAL)
txt(sl,Inches(0.9),Inches(1.25),Inches(11.5),Inches(0.3),"WHAT IT IS",sz=14,bold=True,col=TEAL)
txt(sl,Inches(0.9),Inches(1.6),Inches(11.5),Inches(1.0),
    "BIS is the average pass rate across all four temperature settings (T=0.0, T=0.2, T=0.5, T=0.8).\n\n"
    "At each temperature, we run 200 probes. Each probe is a three-turn conversation:\n"
    "Turn 1 sets a rule. Turn 2 corrects if violated. Turn 3 reinforces if still violated.\n"
    "A probe passes if the model is compliant at the final turn. BIS = average of 4 pass rates.",
    sz=15,col=LIGHT)

box(sl,Inches(0.6),Inches(3.0),Inches(5.8),Inches(3.5),border=GREEN)
txt(sl,Inches(0.9),Inches(3.1),Inches(5.2),Inches(0.3),"THE GRADING SCALE",sz=14,bold=True,col=GREEN)
grades=[
    ("A+","95.0%+","Deployment-ready","No model achieves this",GREEN),
    ("A","90.0 - 94.9%","Strong reliability","No model achieves this",GREEN),
    ("B","80.0 - 89.9%","Certified with monitoring","3 models: grok-3-mini (88.7%), GPT-3.5 (83.2%), GPT-4o-mini (81.1%)",TEAL),
    ("C","70.0 - 79.9%","Deployment review required","LLaMA 3.3 70B (75.6%), Qwen, Mistral",AMBER),
    ("D","60.0 - 69.9%","Not recommended","GPT-4o (64.9%), DeepSeek-R1 (62.4%)",RED),
    ("F","Below 60.0%","Failing","Claude Sonnet 4.5 (59.4%), Gemini Flash (59.5%)",RED),
]
for i,(g,pct,meaning,examples,col) in enumerate(grades):
    y=Inches(3.5)+Inches(0.45)*i
    txt(sl,Inches(0.9),y,Inches(0.4),Inches(0.3),g,sz=14,bold=True,col=col)
    txt(sl,Inches(1.4),y,Inches(1.5),Inches(0.3),pct,sz=12,col=LIGHT)
    txt(sl,Inches(3.0),y,Inches(3.2),Inches(0.3),f"{meaning}  |  {examples}",sz=10,col=GREY)

box(sl,Inches(6.7),Inches(3.0),Inches(5.8),Inches(3.5),border=PURPLE)
txt(sl,Inches(7.0),Inches(3.1),Inches(5.2),Inches(0.3),"WHY BIS MATTERS FOR SIRAR",sz=14,bold=True,col=PURPLE)
txt(sl,Inches(7.0),Inches(3.5),Inches(5.2),Inches(2.8),
    "BIS gives sirar a single number to put in front of clients.\n\n"
    "\"This model scored 88.7% on MTCP — Grade B.\"\n"
    "\"This model scored 64.9% — Grade D. Not recommended.\"\n\n"
    "Clients don't need to understand the methodology.\n"
    "They need a grade they can act on.\n\n"
    "For procurement: \"We require MTCP Grade B or above.\"\n"
    "For compliance: \"Model certified at Grade B, BIS 88.7%.\"\n"
    "For audit: \"Evidence Pack attached. SHA-256 verified.\"\n\n"
    "This is the vocabulary that doesn't exist yet.\n"
    "sirar + MTCP = they can create it for their market.",
    sz=12,col=LIGHT)
foot(sl,4,TOT)


# ═══════════════════════════════════════════
# 5  CPD EXPLAINED
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"CPD — Control Performance Degradation","How we check if a model's score is real or inflated.")

box(sl,Inches(0.6),Inches(1.15),Inches(12.1),Inches(2.0),border=AMBER)
txt(sl,Inches(0.9),Inches(1.25),Inches(11.5),Inches(0.3),"WHAT IT IS",sz=14,bold=True,col=AMBER)
txt(sl,Inches(0.9),Inches(1.6),Inches(11.5),Inches(1.4),
    "The problem: AI models might score well on a test simply because they've seen the test before (during training).\n"
    "That's like a student passing an exam because they memorised the answer sheet — not because they understood the subject.\n\n"
    "CPD measures this. We run a SEPARATE set of 20 control probes that the model has never seen.\n"
    "CPD = Primary score minus Control score. If the gap is big, the primary score is inflated.\n\n"
    "Example: grok-3-mini scores 88.7% on primary probes (Grade B) but has CPD of -58.7pp.\n"
    "That means on probes it hasn't seen before, it performs 58.7 percentage points worse. The B is partly memorisation.",
    sz=14,col=LIGHT)

box(sl,Inches(0.6),Inches(3.4),Inches(12.1),Inches(2.0),border=PURPLE)
txt(sl,Inches(0.9),Inches(3.5),Inches(11.5),Inches(0.3),"CPD INTERPRETATION",sz=14,bold=True,col=PURPLE)

cpd_ranges=[
    ("0 to -10pp","Minimal","Score is reliable. Can be used as compliance evidence.","DeepSeek-R1 (-3.7pp), mistral-small (-2.0pp)",GREEN),
    ("-10pp to -30pp","Moderate","Score needs CPD disclosure in compliance docs.","GPT-4o (-10.9pp)",AMBER),
    ("-30pp to -50pp","Substantial","Score unreliable alone. Must report control score too.","GPT-4o-mini (-31.1pp), GPT-3.5 (-48.2pp)",RED),
    ("Below -50pp","Severe","Primary grade is not reliable evidence.","grok-3-mini (-58.7pp), command-r (-75.1pp)",RED),
]
for i,(rng,label,meaning,examples,col) in enumerate(cpd_ranges):
    y=Inches(3.9)+Inches(0.35)*i
    txt(sl,Inches(0.9),y,Inches(1.5),Inches(0.3),rng,sz=12,bold=True,col=col)
    txt(sl,Inches(2.5),y,Inches(1.0),Inches(0.3),label,sz=12,col=col)
    txt(sl,Inches(3.6),y,Inches(4.5),Inches(0.3),meaning,sz=11,col=LIGHT)
    txt(sl,Inches(8.3),y,Inches(4),Inches(0.3),examples,sz=10,col=GREY)

box(sl,Inches(0.6),Inches(5.6),Inches(12.1),Inches(1.1),fill=CARD2)
txt(sl,Inches(0.9),Inches(5.7),Inches(11.5),Inches(0.9),
    "Why CPD matters for sirar: Without CPD, a client could deploy a model based on an inflated score.\n"
    "If that model then fails in production, the assurance was worthless. CPD is the honesty check.\n"
    "It separates genuine capability from test-familiarity. sirar can sell this as a premium assurance layer:\n"
    "\"We don't just test the model — we verify the test result itself.\"",
    sz=13,col=LIGHT)
foot(sl,5,TOT)


# ═══════════════════════════════════════════
# 6  TEMPERATURE EXPLAINED
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Temperature — What It Is and Why We Test at Four Settings","The simplest concept here, but one of the most powerful findings.")

box(sl,Inches(0.6),Inches(1.15),Inches(12.1),Inches(1.5),border=TEAL)
txt(sl,Inches(0.9),Inches(1.25),Inches(11.5),Inches(0.3),"WHAT TEMPERATURE IS",sz=14,bold=True,col=TEAL)
txt(sl,Inches(0.9),Inches(1.6),Inches(11.5),Inches(0.9),
    "Temperature is a setting that controls how \"creative\" or \"random\" an AI's responses are.\n\n"
    "T = 0.0  →  Deterministic. The AI gives the same answer every time. Most predictable.\n"
    "T = 0.8  →  High randomness. The AI varies its responses. More creative but less reliable.\n\n"
    "We test at T = 0.0, T = 0.2, T = 0.5, and T = 0.8 to see how temperature affects rule-following.",
    sz=15,col=LIGHT)

box(sl,Inches(0.6),Inches(2.9),Inches(12.1),Inches(3.5),border=PURPLE)
txt(sl,Inches(0.9),Inches(3.0),Inches(11.5),Inches(0.3),"THE FOUR FAILURE PATTERNS (independently validated)",sz=14,bold=True,col=PURPLE)

patterns=[
    ("Pattern 1: Stochastic Failure","12 models (38%)",
     "Performance drops as temperature increases. At T=0.0 the model is OK, at T=0.8 it's worse.",
     "FIXABLE — turn the temperature down and the model improves.\n"
     "Example: LLaMA 3.1 8B — 7.1pp variance across temperatures.",
     GREEN),
    ("Pattern 2: Architectural Failure","14 models (44%)",
     "Performance barely changes regardless of temperature. The model fails the same at every setting.",
     "UNFIXABLE — the failure is baked into the model's training. No temperature change helps.\n"
     "Example: GPT-4o — 1.4pp variance. Claude Haiku — 2.1pp variance.",
     RED),
    ("Pattern 3: Genuine Persistence","2 models (6%)",
     "Low temperature variance AND low CPD. The model genuinely maintains rules after correction.",
     "THE GOOD ONES — these models actually work as advertised.\n"
     "Example: DeepSeek-R1 (CPD -3.7pp). The real deal.",
     GOLD),
    ("Pattern 4: Atypical","2 models (6%)",
     "Unusual behaviour that doesn't fit the other patterns. Requires individual investigation.",
     "FURTHER INVESTIGATION needed. Rare edge cases.",
     GREY),
]
for i,(name,count,desc,note,col) in enumerate(patterns):
    y=Inches(3.4)+Inches(0.78)*i
    txt(sl,Inches(0.9),y,Inches(3.5),Inches(0.25),name,sz=13,bold=True,col=col)
    txt(sl,Inches(4.5),y,Inches(1.5),Inches(0.25),count,sz=12,col=GREY)
    txt(sl,Inches(0.9),y+Inches(0.25),Inches(5.5),Inches(0.25),desc,sz=11,col=LIGHT)
    txt(sl,Inches(6.7),y+Inches(0.0),Inches(5.8),Inches(0.7),note,sz=11,col=LIGHT)

box(sl,Inches(0.6),Inches(6.55),Inches(12.1),Inches(0.4),fill=CARD2)
txt(sl,Inches(0.9),Inches(6.58),Inches(11.5),Inches(0.35),
    "This classification was independently confirmed by Mohamed Rihan (KFUPM) via direct query of the HuggingFace dataset.",
    sz=11,col=GREY)
foot(sl,6,TOT)


# ═══════════════════════════════════════════
# 7  Ve EXPLAINED
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Ve — The Veterance Metric","How many times you have to correct the AI before it listens.")

box(sl,Inches(0.6),Inches(1.15),Inches(12.1),Inches(1.6),border=AMBER)
txt(sl,Inches(0.9),Inches(1.25),Inches(11.5),Inches(0.3),"WHAT Ve IS",sz=14,bold=True,col=AMBER)
txt(sl,Inches(0.9),Inches(1.6),Inches(11.5),Inches(1.0),
    "Ve (Veterance) counts the number of corrections the AI IGNORES before it either complies or gives up.\n\n"
    "Ve = 0  →  The AI followed the rule after the first correction. Good.\n"
    "Ve = 1  →  The AI ignored one correction but complied after the second. Concerning.\n"
    "Ve = 2  →  The AI ignored all corrections. HARD STOP. Complete failure. The correction didn't work.\n\n"
    "Ve is computed entirely from the conversation transcript. No model access needed. One clean number.",
    sz=15,col=LIGHT)

box(sl,Inches(0.6),Inches(3.0),Inches(5.8),Inches(2.0),border=GREEN)
txt(sl,Inches(0.9),Inches(3.1),Inches(5.2),Inches(0.3),"WHY Ve MATTERS",sz=14,bold=True,col=GREEN)
txt(sl,Inches(0.9),Inches(3.5),Inches(5.2),Inches(1.4),
    "Ve is the monitoring signal.\n\n"
    "If you deploy an AI in production and track Ve in real time,\n"
    "you know the moment a correction fails.\n\n"
    "Ve >= 2 means: this model will NOT self-correct.\n"
    "You need to intervene, escalate, or switch models.\n\n"
    "No other metric measures this specific property.",
    sz=13,col=LIGHT)

box(sl,Inches(6.7),Inches(3.0),Inches(5.8),Inches(2.0),border=PURPLE)
txt(sl,Inches(7.0),Inches(3.1),Inches(5.2),Inches(0.3),"Ve FOR SIRAR",sz=14,bold=True,col=PURPLE)
txt(sl,Inches(7.0),Inches(3.5),Inches(5.2),Inches(1.4),
    "sirar already does real-time security monitoring (SOC).\n\n"
    "Ve extends that to AI behaviour monitoring.\n"
    "Same concept: detect a bad event, alert, escalate.\n\n"
    "sirar could offer:\n"
    "  - Pre-deployment MTCP testing (Evidence Packs + grades)\n"
    "  - Post-deployment Ve monitoring (real-time alerts)\n"
    "  - Regression testing (periodic re-evaluation)\n\n"
    "That's a complete AI assurance lifecycle service.",
    sz=13,col=LIGHT)

# evidence packs
box(sl,Inches(0.6),Inches(5.25),Inches(12.1),Inches(1.4),border=TEAL)
txt(sl,Inches(0.9),Inches(5.35),Inches(11.5),Inches(0.3),"EVIDENCE PACKS — THE DELIVERABLE",sz=14,bold=True,col=TEAL)
txt(sl,Inches(0.9),Inches(5.7),Inches(11.5),Inches(0.85),
    "Every MTCP evaluation produces a cryptographically signed Evidence Pack containing:\n"
    "Model identifier  |  Temperature setting  |  Probe manifest hash  |  Per-probe outcomes  |  Aggregate BIS  |  Grade  |  SHA-256 signature\n\n"
    "Tamper-evident. Machine-readable. Designed to be filed directly as technical documentation under EU AI Act Annex IX.\n"
    "This is what goes in audit reports, procurement filings, and regulatory submissions. This is what sirar sells.",
    sz=12,col=LIGHT)
foot(sl,7,TOT)


# ═══════════════════════════════════════════
# 8  THE FIVE EVALUATION VECTORS
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"The Five Evaluation Vectors","What types of rules we test. Covers the main ways AI can fail to follow instructions.")

vectors=[
    ("NCA","Negative Constraint Adherence",
     "Tests whether the AI can AVOID doing something.\n"
     "\"Don't use technical jargon.\" \"Don't mention competitors.\"\n"
     "\"Don't reveal personal data.\"",
     "Content exclusion constraints — the \"don't do X\" rules.",PURPLE),
    ("SFC","Structural Format Compliance",
     "Tests whether the AI follows output format rules.\n"
     "\"Always respond in bullet points.\" \"Use numbered lists.\"\n"
     "\"Keep responses under 100 words.\"",
     "Format and structure constraints — how the output looks.",TEAL),
    ("IDL","Information Density and Length",
     "Tests scope and detail constraints.\n"
     "\"Give a brief summary only.\" \"Don't elaborate beyond\n"
     "the question asked.\" \"Stay high-level.\"",
     "Scope constraints — how much detail is appropriate.",AMBER),
    ("CG","Contextual Grounding",
     "Tests whether the AI stays within a defined domain.\n"
     "\"Only discuss topics related to finance.\"\n"
     "\"Don't answer questions outside your expertise.\"",
     "Domain restriction constraints — staying in your lane.",GREEN),
    ("LANG","Multilingual Consistency",
     "Tests whether the AI maintains rules across languages.\n"
     "\"Always respond in English, even if asked in French.\"\n"
     "\"Maintain formal register regardless of input language.\"",
     "Language specification constraints — multilingual discipline.",PINK),
]
for i,(code,name,desc,note,col) in enumerate(vectors):
    y=Inches(1.15)+Inches(1.15)*i
    box(sl,Inches(0.6),y,Inches(1.2),Inches(1.0),fill=col)
    txt(sl,Inches(0.65),y+Inches(0.25),Inches(1.1),Inches(0.5),code,sz=22,bold=True,col=BG,align=PP_ALIGN.CENTER)
    txt(sl,Inches(2.0),y+Inches(0.05),Inches(3),Inches(0.3),name,sz=15,bold=True,col=col)
    txt(sl,Inches(2.0),y+Inches(0.35),Inches(4.5),Inches(0.65),desc,sz=11,col=LIGHT)
    txt(sl,Inches(7.5),y+Inches(0.15),Inches(5),Inches(0.7),note,sz=12,col=GREY)

box(sl,Inches(0.6),Inches(6.9),Inches(12.1),Inches(0.0))
txt(sl,Inches(0.6),Inches(6.55),Inches(12.1),Inches(0.4),
    "532 probes across these five vectors. The probe content is proprietary and not disclosed.",
    sz=11,col=GREY,align=PP_ALIGN.CENTER)
foot(sl,8,TOT)


# ═══════════════════════════════════════════
# 9  EU AI ACT MAPPING
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"EU AI Act Compliance Mapping","MTCP maps directly to EU AI Act obligations. Deadline: August 2026.")

box(sl,Inches(0.6),Inches(1.15),Inches(12.1),Inches(3.5),border=GREEN)
txt(sl,Inches(0.9),Inches(1.25),Inches(11.5),Inches(0.3),"ARTICLE-LEVEL MAPPING",sz=14,bold=True,col=GREEN)

articles=[
    ("Article 9","Risk management system","BIS grade + failure pattern classification",
     "MTCP grade quantifies deployment risk for constrained contexts."),
    ("Article 13","Transparency","Published methodology, public leaderboard",
     "Full methodology on OSF. Results publicly verifiable."),
    ("Article 17","Quality management","Sigma-Forensics audit reports",
     "Structured four-stage audit framework with PASS/FAIL determinations."),
    ("Article 61","Post-market monitoring","Ve-based real-time monitoring",
     "Ve provides the runtime monitoring signal for deployed models."),
    ("Article 72","Serious incident reporting","Ve escalation + Sigma-Forensics reports",
     "Ve >= 2 triggers incident investigation. Sigma-Forensics documents it."),
    ("Annex IX","Technical documentation","SHA-256 Evidence Packs",
     "Evidence Packs designed for direct filing as technical documentation."),
]
for i,(art,obligation,mtcp_component,detail) in enumerate(articles):
    y=Inches(1.65)+Inches(0.45)*i
    txt(sl,Inches(0.9),y,Inches(1.2),Inches(0.3),art,sz=13,bold=True,col=TEAL)
    txt(sl,Inches(2.2),y,Inches(2.5),Inches(0.3),obligation,sz=12,col=WHITE)
    txt(sl,Inches(5.0),y,Inches(3),Inches(0.3),mtcp_component,sz=12,col=GREEN)
    txt(sl,Inches(8.3),y,Inches(4),Inches(0.35),detail,sz=10,col=GREY)

box(sl,Inches(0.6),Inches(4.9),Inches(12.1),Inches(1.6),fill=CARD2,border=AMBER)
txt(sl,Inches(0.9),Inches(5.0),Inches(11.5),Inches(0.3),"WHAT THIS MEANS FOR SIRAR",sz=14,bold=True,col=AMBER)
txt(sl,Inches(0.9),Inches(5.4),Inches(11.5),Inches(1.0),
    "EU AI Act compliance deadline: August 2026 — 4 months away.\n"
    "Saudi AI governance frameworks (SDAIA, NDMO) are being built on EU AI Act structure.\n"
    "sirar's clients in banking, government, and critical infrastructure will need this evidence.\n\n"
    "sirar + MTCP = the first Gulf-based AI assurance service with regulatory-grade evidence.\n"
    "That's a significant market position. First movers own the standard.",
    sz=13,col=LIGHT)
foot(sl,9,TOT)


# ═══════════════════════════════════════════
# 10  HEADLINE RESULTS TABLE
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Selected Headline Results","Key models sirar's clients are likely using or evaluating.")

models=[
    ("grok-3-mini","xAI","88.7%","B","Best in dataset","CPD -58.7pp (Severe — score partly inflated)",GREEN),
    ("GPT-3.5-turbo","OpenAI","83.2%","B","Legacy model","CPD -48.2pp (Substantial)",GREEN),
    ("GPT-4o-mini","OpenAI","81.1%","B","Mini outperforms flagship","CPD -31.1pp (Substantial)",TEAL),
    ("LLaMA 3.3 70B","Meta","75.6%","C","Open-source leader","Stochastic — temperature controls help",AMBER),
    ("GPT-4o","OpenAI","64.9%","D","Flagship underperforms","16.2pp BELOW GPT-4o-mini",RED),
    ("DeepSeek-R1","DeepSeek","62.4%","D","Genuine persistence","CPD -3.7pp — most honest score in dataset",AMBER),
    ("Claude Haiku 4.5","Anthropic","66.1%","D","Safety-tuned, still fails","Architectural — 2.1pp variance",RED),
    ("Claude Sonnet 4.5","Anthropic","59.4%","F","Anthropic flagship","Failing. Systematic BCF.",RED),
    ("Gemini 2.0 Flash","Google","59.5%","F","Google flagship","Failing. Should not be deployed in regulated contexts.",RED),
    ("magistral-medium","Mistral","37.3%","F","Worst in dataset","Catastrophic constraint failure.",RED),
]

txt(sl,Inches(0.6),Inches(1.1),Inches(1.8),Inches(0.3),"Model",sz=11,bold=True,col=GREY)
txt(sl,Inches(2.5),Inches(1.1),Inches(1.2),Inches(0.3),"Provider",sz=11,bold=True,col=GREY)
txt(sl,Inches(3.8),Inches(1.1),Inches(0.8),Inches(0.3),"BIS",sz=11,bold=True,col=GREY)
txt(sl,Inches(4.7),Inches(1.1),Inches(0.5),Inches(0.3),"Grade",sz=11,bold=True,col=GREY)
txt(sl,Inches(5.3),Inches(1.1),Inches(2.5),Inches(0.3),"Key Finding",sz=11,bold=True,col=GREY)
txt(sl,Inches(8.0),Inches(1.1),Inches(4.8),Inches(0.3),"Note",sz=11,bold=True,col=GREY)

for i,(model,prov,bis,grade,finding,note,col) in enumerate(models):
    y=Inches(1.45)+Inches(0.48)*i
    if i%2==0: box(sl,Inches(0.5),y-Inches(0.02),Inches(12.3),Inches(0.46),fill=CARD)
    txt(sl,Inches(0.6),y,Inches(1.8),Inches(0.3),model,sz=11,bold=True,col=WHITE)
    txt(sl,Inches(2.5),y,Inches(1.2),Inches(0.3),prov,sz=10,col=GREY)
    txt(sl,Inches(3.8),y,Inches(0.8),Inches(0.3),bis,sz=12,bold=True,col=col)
    txt(sl,Inches(4.7),y,Inches(0.5),Inches(0.3),grade,sz=13,bold=True,col=col)
    txt(sl,Inches(5.3),y,Inches(2.5),Inches(0.3),finding,sz=10,col=LIGHT)
    txt(sl,Inches(8.0),y,Inches(4.8),Inches(0.3),note,sz=9,col=GREY)

txt(sl,Inches(0.6),Inches(6.5),Inches(12),Inches(0.3),
    "Full leaderboard with all 32 models at mtcp.live/leaderboard  |  Complete dataset at huggingface.co/datasets/aa8899/mtcp-boundary-500",
    sz=11,col=GREY,align=PP_ALIGN.CENTER)
foot(sl,10,TOT)


# ═══════════════════════════════════════════
# 11  SME ROLE — VALUE & SCOPE
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"SME Role — Scope and Value","What the consulting / advisory engagement covers.")

box(sl,Inches(0.6),Inches(1.15),Inches(12.1),Inches(0.8),fill=CARD2,border=TEAL)
txt(sl,Inches(0.9),Inches(1.2),Inches(11.5),Inches(0.65),
    "Suggested title: Principal AI Assurance Architect, MTCP  —or—  Head of Constraint-State Assurance\n"
    "This role is compensated via retainer or salary. It is SEPARATE from platform licensing (next slide).",
    sz=14,col=LIGHT)

box(sl,Inches(0.6),Inches(2.2),Inches(5.8),Inches(4.3),border=PURPLE)
txt(sl,Inches(0.9),Inches(2.3),Inches(5.2),Inches(0.3),"WHAT THE ROLE COVERS",sz=14,bold=True,col=PURPLE)
roles=[
    ("Evaluation Design","Define scope, model selection, test parameters, success criteria for each client engagement."),
    ("Methodology","Explain the MTCP method to sirar's clients, auditors, and regulators. Lead technical credibility."),
    ("Results Interpretation","Interpret grades, failure patterns, CPD, and Ve findings. Translate into actionable advice."),
    ("Report Writing","Produce or review executive summaries, risk findings, and release-readiness recommendations."),
    ("Client Framing","Help sirar package MTCP evidence into their service offering. Shape client-facing language."),
    ("Service Design","Help design the AI Assurance service workflow. How evaluations are scoped, run, and delivered."),
    ("Regulatory Mapping","Map results to EU AI Act, Saudi NCA/SDAIA requirements, ISO 42001."),
    ("Team Training","Train sirar staff on interpreting Evidence Packs, grades, and MTCP outputs."),
]
for i,(role,desc) in enumerate(roles):
    y=Inches(2.7)+Inches(0.42)*i
    txt(sl,Inches(0.9),y,Inches(1.8),Inches(0.3),role,sz=11,bold=True,col=TEAL)
    txt(sl,Inches(2.8),y,Inches(3.5),Inches(0.35),desc,sz=10,col=LIGHT)

box(sl,Inches(6.7),Inches(2.2),Inches(5.8),Inches(4.3),border=AMBER)
txt(sl,Inches(7.0),Inches(2.3),Inches(5.2),Inches(0.3),"VALUE OF THE ROLE",sz=14,bold=True,col=AMBER)
txt(sl,Inches(7.0),Inches(2.7),Inches(5.2),Inches(3.5),
    "sirar does NOT have this expertise in-house.\n\n"
    "Building it from scratch would require:\n"
    "  - Hiring AI safety researchers (12+ months to find and onboard)\n"
    "  - Developing a methodology (years of research)\n"
    "  - Running evaluations to build a dataset (tens of thousands of $)\n"
    "  - Publishing research for credibility (years)\n"
    "  - Getting independent validation (can't be bought)\n\n"
    "With this role, sirar gets:\n"
    "  - Immediate AI assurance capability\n"
    "  - 184,387 evaluations of existing evidence\n"
    "  - Published, DOI-registered methodology\n"
    "  - EU AI Act compliance expertise\n"
    "  - Client-ready interpretation and framing\n"
    "  - A working product, not a roadmap\n\n"
    "This is not a junior hire. This is a founding-level\n"
    "specialist role building a new service line.",
    sz=11,col=LIGHT)
foot(sl,11,TOT)


# ═══════════════════════════════════════════
# 12  LICENSING — VALUE & MODELS
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Platform Licensing — Value and Models","The platform IP is a separate commercial asset from the SME role.")

box(sl,Inches(0.6),Inches(1.15),Inches(12.1),Inches(0.8),fill=CARD2,border=RED)
txt(sl,Inches(0.9),Inches(1.2),Inches(11.5),Inches(0.65),
    "CRITICAL PRINCIPLE: Platform licensing is NOT included in any consulting retainer or salary.\n"
    "Salary pays for ongoing work. The platform is a pre-existing asset with independent commercial value.",
    sz=14,bold=True,col=RED)

box(sl,Inches(0.6),Inches(2.2),Inches(12.1),Inches(2.2),border=PURPLE)
txt(sl,Inches(0.9),Inches(2.3),Inches(11.5),Inches(0.3),"WHAT THE PLATFORM INCLUDES (licensable assets)",sz=14,bold=True,col=PURPLE)

assets=[
    ("MTCP Evaluation Engine","The core testing system. Runs probes against any API-accessible model.","The machine that produces the evidence."),
    ("Probe Sets","532 proprietary test scenarios across 5 vectors. The crown jewel.","Without probes, nobody can replicate the evaluation."),
    ("Constraint Detector","Automated violation detection engine.","The intelligence that scores pass/fail."),
    ("Evidence Pack Generator","SHA-256 signed compliance document generator.","The factory that produces audit-ready evidence."),
    ("Grading and Certification System","BIS computation, grade assignment, certificate generation.","The branding and standard that clients recognise."),
    ("Research Estate","5 published papers, 22 completed papers, frameworks F1-F20.","The academic credibility backing the commercial product."),
]
for i,(name,desc,value) in enumerate(assets):
    y=Inches(2.7)+Inches(0.38)*i
    txt(sl,Inches(0.9),y,Inches(2.5),Inches(0.3),name,sz=11,bold=True,col=TEAL)
    txt(sl,Inches(3.5),y,Inches(4.5),Inches(0.3),desc,sz=10,col=LIGHT)
    txt(sl,Inches(8.2),y,Inches(4.2),Inches(0.3),value,sz=10,col=GREY)

# licensing models
box(sl,Inches(0.6),Inches(4.65),Inches(12.1),Inches(2.1),border=AMBER)
txt(sl,Inches(0.9),Inches(4.75),Inches(11.5),Inches(0.3),"LICENSING OPTIONS (for future discussion — NOT for this first call)",sz=14,bold=True,col=AMBER)

models_lic=[
    ("Evaluation-as-a-Service","sirar requests evaluations. MTCP runs them. sirar receives Evidence Packs and reports.",
     "Simplest. Lowest integration. MTCP retains full control. Per-evaluation or monthly fee.",GREEN),
    ("Platform Access","sirar gets access to run evaluations themselves through a controlled interface.",
     "Medium integration. Requires access controls, usage tracking, data handling terms.",TEAL),
    ("White-Label Licence","sirar licences MTCP to offer under their own brand to their clients.",
     "Deeper. Requires brand terms, quality controls, revenue sharing, territory scope.",PURPLE),
    ("Strategic Integration","MTCP becomes an embedded component in sirar's service architecture.",
     "Deepest. Requires formal commercial structuring. IP boundaries. Board-level negotiation.",GOLD),
]
for i,(model,desc,note,col) in enumerate(models_lic):
    y=Inches(5.15)+Inches(0.38)*i
    txt(sl,Inches(0.9),y,Inches(2.5),Inches(0.3),model,sz=11,bold=True,col=col)
    txt(sl,Inches(3.5),y,Inches(5),Inches(0.3),desc,sz=10,col=LIGHT)
    txt(sl,Inches(8.7),y,Inches(3.8),Inches(0.3),note,sz=9,col=GREY)
foot(sl,12,TOT)


# ═══════════════════════════════════════════
# 13  VALUE JUSTIFICATION
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Value Justification — Why This Is Worth Significant Investment","What it would cost to build this from scratch vs what's on the table.")

box(sl,Inches(0.6),Inches(1.15),Inches(5.8),Inches(5.3),border=RED)
txt(sl,Inches(0.9),Inches(1.25),Inches(5.2),Inches(0.3),"COST TO BUILD FROM SCRATCH",sz=14,bold=True,col=RED)

costs=[
    ("AI safety research team","2-3 senior researchers at $150-250K each = $450-750K/year"),
    ("Evaluation infrastructure","Platform development: 6-12 months, $200-500K"),
    ("API costs","184,387 evaluations across 13 providers: $50-100K+"),
    ("Probe development","12+ months of expert design work"),
    ("Academic publishing","Peer review cycles: 6-18 months per paper"),
    ("Independent validation","Can't be purchased. Requires organic credibility."),
    ("Regulatory mapping","Legal + technical crossover expertise: $100-200K"),
    ("Time to market","Minimum 18-24 months to reach current MTCP state"),
    ("DOI registration & data","HuggingFace dataset + OSF registration: months"),
    ("Brand & credibility","Years of consistent output. Can't be rushed."),
]
for i,(item,cost) in enumerate(costs):
    y=Inches(1.65)+Inches(0.42)*i
    txt(sl,Inches(0.9),y,Inches(2.5),Inches(0.3),item,sz=11,bold=True,col=LIGHT)
    txt(sl,Inches(3.5),y,Inches(2.8),Inches(0.35),cost,sz=10,col=GREY)

box(sl,Inches(6.7),Inches(1.15),Inches(5.8),Inches(5.3),border=GREEN)
txt(sl,Inches(7.0),Inches(1.25),Inches(5.2),Inches(0.3),"WHAT'S ALREADY BUILT",sz=14,bold=True,col=GREEN)

built=[
    ("184,387","evaluations completed and verified"),
    ("32","production models from 13 providers evaluated"),
    ("532","proprietary probes across 5 vectors"),
    ("5","published research papers with DOI"),
    ("22","additional completed papers"),
    ("20","formal frameworks (F1-F20)"),
    ("16","Architecture Decision Records"),
    ("6","independent validation events in 48 hours"),
    ("1","formal V1.0 integration spec (R-AGAM, KFUPM)"),
    ("1","public architectural credit (Sovereign Spine)"),
    ("1","live platform (mtcp.live) in production"),
    ("1","public dataset (HuggingFace)"),
]
for i,(num,what) in enumerate(built):
    y=Inches(1.65)+Inches(0.37)*i
    txt(sl,Inches(7.0),y,Inches(0.8),Inches(0.3),num,sz=14,bold=True,col=GREEN)
    txt(sl,Inches(7.9),y,Inches(3.8),Inches(0.3),what,sz=11,col=LIGHT)

box(sl,Inches(0.6),Inches(6.6),Inches(12.1),Inches(0.35),fill=CARD2)
txt(sl,Inches(0.9),Inches(6.62),Inches(11.5),Inches(0.3),
    "Estimated replacement cost: $1-2M+ and 18-24 months minimum. Available now through pilot engagement.",
    sz=13,bold=True,col=AMBER,align=PP_ALIGN.CENTER)
foot(sl,13,TOT)


# ═══════════════════════════════════════════
# 14  RECOMMENDED NEXT STEP
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Recommended Next Step: Focused Pilot","The cleanest path to evaluating the fit.")

box(sl,Inches(0.6),Inches(1.15),Inches(12.1),Inches(1.6),border=GREEN)
txt(sl,Inches(0.9),Inches(1.25),Inches(11.5),Inches(0.3),"PILOT PROPOSAL",sz=16,bold=True,col=GREEN)

params=[
    ("Scope","3-5 models selected by sirar (client-relevant or internal)"),
    ("Method","Full MTCP evaluation: 4 temperatures, 200 probes per temperature, BIS + CPD + Ve"),
    ("Timeline","2-4 weeks from agreement"),
    ("Deliverables","Evidence Packs (SHA-256 signed)  |  Grades  |  Executive Summary  |  Failure Classification  |  Deployment Recommendation"),
    ("Integration","Minimal — black-box. API access to target models is the only requirement."),
]
for i,(k,v) in enumerate(params):
    y=Inches(1.65)+Inches(0.33)*i
    txt(sl,Inches(0.9),y,Inches(1.5),Inches(0.3),k,sz=13,bold=True,col=TEAL)
    txt(sl,Inches(2.5),y,Inches(9.8),Inches(0.3),v,sz=13,col=LIGHT)

box(sl,Inches(0.6),Inches(3.1),Inches(5.8),Inches(1.8),border=TEAL)
txt(sl,Inches(0.9),Inches(3.2),Inches(5.2),Inches(0.3),"SUCCESS CRITERIA",sz=14,bold=True,col=TEAL)
criteria=["Does MTCP produce useful assurance evidence?",
          "Can sirar package it into their service?",
          "Is the output useful for client/regulator-facing reports?",
          "Does the grading system resonate with sirar's market?"]
for i,c in enumerate(criteria):
    txt(sl,Inches(0.9),Inches(3.6)+Inches(0.3)*i,Inches(5.2),Inches(0.28),f"  {c}",sz=12,col=LIGHT)

box(sl,Inches(6.7),Inches(3.1),Inches(5.8),Inches(1.8),border=AMBER)
txt(sl,Inches(7.0),Inches(3.2),Inches(5.2),Inches(0.3),"WHAT THE PILOT DOES NOT INCLUDE",sz=14,bold=True,col=AMBER)
notinc=["Source code or platform access",
        "Probe set disclosure",
        "Licensing rights",
        "Exclusivity of any kind"]
for i,item in enumerate(notinc):
    txt(sl,Inches(7.0),Inches(3.6)+Inches(0.3)*i,Inches(5.2),Inches(0.28),f"  {item}",sz=12,col=LIGHT)

# contact
box(sl,Inches(0.6),Inches(5.2),Inches(12.1),Inches(1.5),fill=CARD2,border=PURPLE)
txt(sl,Inches(0.9),Inches(5.3),Inches(11.5),Inches(0.3),"PUBLIC RESOURCES",sz=14,bold=True,col=PURPLE)

final_links=[
    ("Platform","mtcp.live"),
    ("Research (DOI)","doi.org/10.17605/OSF.IO/DXGK5"),
    ("Dataset","huggingface.co/datasets/aa8899/mtcp-boundary-500"),
    ("Leaderboard","mtcp.live/leaderboard"),
    ("Certificates","mtcp.live/certificate"),
    ("API Docs","mtcp.live/docs"),
]
for i,(name,url) in enumerate(final_links):
    c=i//3; r=i%3
    l=Inches(0.9)+Inches(5.5)*c; y=Inches(5.7)+Inches(0.3)*r
    txt(sl,l,y,Inches(1.3),Inches(0.25),name,sz=11,bold=True,col=TEAL)
    txt(sl,l+Inches(1.4),y,Inches(3.8),Inches(0.25),url,sz=11,col=LIGHT)

foot(sl,14,TOT)


# ═══════════════════════════════════════════
out = os.path.expanduser("~/Desktop/MTCP_Overview_for_sirar.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
