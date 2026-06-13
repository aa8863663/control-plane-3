from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG       = RGBColor(0x0D, 0x0D, 0x12)
CARD     = RGBColor(0x16, 0x16, 0x1F)
CARD2    = RGBColor(0x1C, 0x1C, 0x2A)
PURPLE   = RGBColor(0x6C, 0x5C, 0xE7)
TEAL     = RGBColor(0x00, 0xD2, 0xD3)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREY     = RGBColor(0x9B, 0x9B, 0xAA)
LIGHT    = RGBColor(0xCC, 0xCC, 0xDD)
RED      = RGBColor(0xE1, 0x50, 0x50)
GREEN    = RGBColor(0x00, 0xE6, 0x96)
AMBER    = RGBColor(0xFF, 0xB8, 0x42)
PINK     = RGBColor(0xFD, 0x79, 0xA8)
FONT = "Calibri"; MONO = "Courier New"

def bg(s): f=s.background.fill; f.solid(); f.fore_color.rgb=BG
def box(s,l,t,w,h,fill=CARD,border=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if border: sh.line.color.rgb=border; sh.line.width=Pt(1)
    else: sh.line.fill.background()
    sh.adjustments[0]=0.04; return sh
def txt(s,l,t,w,h,text,sz=18,col=WHITE,bold=False,font=FONT,align=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(sz); p.font.color.rgb=col
    p.font.bold=bold; p.font.name=font; p.alignment=align; return tb
def title(s,t,sub=""):
    txt(s,Inches(0.6),Inches(0.35),Inches(12),Inches(0.6),t,sz=34,bold=True)
    if sub: txt(s,Inches(0.6),Inches(0.95),Inches(12),Inches(0.45),sub,sz=16,col=GREY)
def foot(s,n,tot=12):
    txt(s,Inches(0.6),Inches(7.05),Inches(6),Inches(0.3),"YOUR EYES ONLY — sirar intel & fit analysis",sz=9,col=RGBColor(0x55,0x55,0x66))
    txt(s,Inches(11.8),Inches(7.05),Inches(1.2),Inches(0.3),f"{n}/{tot}",sz=9,col=GREY,align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
# 1  COVER
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
box(sl,Inches(0.6),Inches(1.2),Inches(0.1),Inches(2.0),fill=TEAL)
txt(sl,Inches(1.0),Inches(1.2),Inches(11),Inches(0.7),"SIRAR — WHO THEY ARE & HOW YOU FIT",sz=44,bold=True)
txt(sl,Inches(1.0),Inches(2.0),Inches(11),Inches(0.6),"Intel sheet for your call. Know them before they know you.",sz=22,col=GREY)

box(sl,Inches(0.6),Inches(3.2),Inches(12.1),Inches(3.3),border=TEAL)
txt(sl,Inches(0.9),Inches(3.35),Inches(11),Inches(0.4),"THE 30-SECOND VERSION",sz=18,bold=True,col=TEAL)
txt(sl,Inches(0.9),Inches(3.85),Inches(11.5),Inches(2.5),
    "sirar by stc is Saudi Arabia's #1 cybersecurity company.\n"
    "They're owned by STC Group (Saudi Telecom) — a $50B+ mega-corporation.\n"
    "~500 employees across 3 continents. Ranked #1 MSSP in the Middle East, #11 globally.\n\n"
    "They protect Saudi critical infrastructure: railways, airports, water, banking, even Hajj.\n"
    "They are part of Saudi Vision 2030 — the country's plan to diversify beyond oil.\n\n"
    "They're now building an AI service. They don't have AI testing expertise yet.\n"
    "That's why they're talking to you. You have what they need and can't build quickly.",
    sz=17,col=LIGHT)
foot(sl,1)


# ═══════════════════════════════════════════
# 2  WHAT SIRAR IS
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title(sl,"What sirar Actually Is","The basics — so you don't get caught out.")

facts=[
    ("Full name","sirar by stc (always lowercase — that's their style)"),
    ("Parent company","STC Group — Saudi Telecom Company. Biggest telecom in the Middle East. ~$50B+ market cap."),
    ("HQ","Riyadh, Saudi Arabia"),
    ("Size","~500 employees across 3 continents"),
    ("CEO","Yasser Alswailem"),
    ("What they do","Cybersecurity. Detection, response, compliance, advisory, managed security."),
    ("Rankings","#1 MSSP in Middle East & North Africa. #11 globally. Top 3 workplace in Saudi."),
    ("NCA Tier 1","One of only SIX companies with a Tier 1 managed SOC licence from Saudi's National Cybersecurity Authority."),
    ("Website","sirar.com.sa"),
    ("Vibe","Professional, national-pride, trust-focused. They see themselves as protecting Saudi's digital future."),
]
for i,(k,v) in enumerate(facts):
    y=Inches(1.35)+Inches(0.54)*i
    box(sl,Inches(0.6),y,Inches(2.5),Inches(0.45),fill=CARD,border=TEAL)
    txt(sl,Inches(0.7),y+Inches(0.05),Inches(2.3),Inches(0.35),k,sz=13,bold=True,col=TEAL)
    txt(sl,Inches(3.3),y+Inches(0.05),Inches(9.3),Inches(0.4),v,sz=13,col=LIGHT)

foot(sl,2)


# ═══════════════════════════════════════════
# 3  THEIR SERVICES
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title(sl,"What sirar Sells Today","Five pillars. Your stuff would be a sixth — or plug into PLAN.")

pillars=[
    ("DETECT","DDoS protection, SOC, threat intel, pen testing, SIEM",PURPLE),
    ("RESPOND","Incident response, managed detection, cyber resilience",TEAL),
    ("PLAN","GRC assessments, compliance advisory, strategy, training",AMBER),
    ("BUILD","Firewalls, IAM, encryption, endpoint, data classification",GREEN),
    ("TRUST","Their own products: Athar (screen watermarks), Sayen (digital signatures)",PINK),
]
for i,(name,desc,col) in enumerate(pillars):
    y=Inches(1.4)+Inches(0.95)*i
    box(sl,Inches(0.6),y,Inches(1.5),Inches(0.8),fill=col)
    txt(sl,Inches(0.65),y+Inches(0.15),Inches(1.4),Inches(0.5),name,sz=20,bold=True,col=BG,align=PP_ALIGN.CENTER)
    txt(sl,Inches(2.3),y+Inches(0.15),Inches(10),Inches(0.5),desc,sz=15,col=LIGHT)

box(sl,Inches(0.6),Inches(6.15),Inches(12.1),Inches(0.6),fill=CARD2,border=AMBER)
txt(sl,Inches(0.9),Inches(6.22),Inches(11.5),Inches(0.45),
    "Where MTCP fits: most naturally under PLAN (advisory/compliance) or as a new AI-specific pillar entirely.\n"
    "If they position it under TRUST (their own products), that's a licensing conversation — different commercial model.",
    sz=13,col=AMBER)
foot(sl,3)


# ═══════════════════════════════════════════
# 4  THEIR CLIENTS
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title(sl,"Who sirar's Clients Are","These are the people who'd eventually see YOUR work.")

box(sl,Inches(0.6),Inches(1.4),Inches(5.8),Inches(4.0),border=PURPLE)
txt(sl,Inches(0.9),Inches(1.55),Inches(5),Inches(0.4),"CONFIRMED CLIENTS (public info)",sz=16,bold=True,col=PURPLE)
clients=[
    ("Saudi Railway Company (SAR)","National rail network cybersecurity"),
    ("Riyadh Airports Company","Airport digital infrastructure"),
    ("National Water Company","OT infrastructure security"),
    ("Red Sea Global","Vision 2030 tourism mega-project"),
    ("KAFD (King Abdullah Financial District)","AI-powered cybersecurity"),
    ("Hajj digital infrastructure","3 consecutive years protecting pilgrimage systems"),
    ("Kafalah","Saudi SME guarantee programme"),
]
for i,(client,what) in enumerate(clients):
    y=Inches(2.05)+Inches(0.45)*i
    txt(sl,Inches(0.9),y,Inches(2.8),Inches(0.3),client,sz=12,bold=True,col=LIGHT)
    txt(sl,Inches(3.8),y,Inches(2.5),Inches(0.3),what,sz=11,col=GREY)

box(sl,Inches(6.7),Inches(1.4),Inches(5.8),Inches(4.0),border=TEAL)
txt(sl,Inches(7.0),Inches(1.55),Inches(5),Inches(0.4),"SECTORS THEY SERVE",sz=16,bold=True,col=TEAL)
sectors=[
    ("Banking & financial services","Where compliance pressure is highest"),
    ("Energy & utilities","Aramco, national grid, oil & gas"),
    ("Aerospace, defence, intelligence","Government / military adjacent"),
    ("Healthcare","Sensitive data, high regulation"),
    ("Telecom","STC itself + other operators"),
    ("Government","Vision 2030 mega-projects"),
]
for i,(sector,note) in enumerate(sectors):
    y=Inches(2.05)+Inches(0.45)*i
    txt(sl,Inches(7.0),y,Inches(2.8),Inches(0.3),sector,sz=12,bold=True,col=LIGHT)
    txt(sl,Inches(9.8),y,Inches(2.5),Inches(0.3),note,sz=11,col=GREY)

box(sl,Inches(0.6),Inches(5.65),Inches(12.1),Inches(0.7),fill=CARD2)
txt(sl,Inches(0.9),Inches(5.72),Inches(11.5),Inches(0.55),
    "What this means for you: sirar's clients are EXACTLY the ones who'll need AI assurance evidence.\n"
    "Banks, government, railways, airports — these are regulated, high-stakes. They'll need MTCP grades.",
    sz=14,bold=True,col=GREEN)
foot(sl,4)


# ═══════════════════════════════════════════
# 5  WHO IS SYED
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title(sl,"Who Is Syed?","Your best guess for who you're talking to.")

box(sl,Inches(0.6),Inches(1.4),Inches(12.1),Inches(2.2),border=PURPLE)
txt(sl,Inches(0.9),Inches(1.55),Inches(11),Inches(0.4),"MOST LIKELY: Mudassir Syed",sz=20,bold=True,col=PURPLE)
syed_info=[
    ("Title","Senior Manager, Product Management & Marketing — Managed Security Services"),
    ("At sirar since","August 2021 (~5 years)"),
    ("Background","Previously at Etisalat (big UAE telecom) and Integrated Telecom Company"),
    ("Education","M.Eng, MBA, Certified Product Manager"),
    ("What he does","Develops product strategies, manages service portfolios, leads new cybersecurity service launches"),
    ("His team","Part of a 3-person Product Management and Innovation unit"),
]
for i,(k,v) in enumerate(syed_info):
    y=Inches(2.1)+Inches(0.38)*i
    txt(sl,Inches(0.9),y,Inches(2.8),Inches(0.3),k,sz=13,bold=True,col=TEAL)
    txt(sl,Inches(3.8),y,Inches(8.5),Inches(0.3),v,sz=13,col=LIGHT)

box(sl,Inches(0.6),Inches(3.9),Inches(12.1),Inches(2.5),fill=CARD2)
txt(sl,Inches(0.9),Inches(4.05),Inches(11.5),Inches(0.4),
    "WHAT THIS TELLS YOU ABOUT THE CALL",sz=16,bold=True,col=AMBER)

insights=[
    ("He's a product person, not a tech person",
     "He thinks in terms of services, clients, and market positioning — not algorithms. Speak his language."),
    ("He's been there 5 years",
     "He's established and trusted. If he's bringing you in, he's already sold it internally to some degree."),
    ("He manages NEW service launches",
     "This is literally his job. AI Assurance is a new service. He's building it. He needs what you have."),
    ("He came from telecom, not security research",
     "Don't geek out on Ve metrics. Talk about what the service looks like to a client and how to package it."),
]
for i,(point,detail) in enumerate(insights):
    y=Inches(4.55)+Inches(0.48)*i
    txt(sl,Inches(0.9),y,Inches(5),Inches(0.3),point,sz=13,bold=True,col=AMBER)
    txt(sl,Inches(6.0),y,Inches(6.5),Inches(0.35),detail,sz=12,col=LIGHT)

foot(sl,5)


# ═══════════════════════════════════════════
# 6  THEIR AI MOVES SO FAR
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title(sl,"sirar's AI Moves So Far","They're building toward AI, but they don't have the testing piece yet.")

moves=[
    ("Dec 2025","/ai/ page published on sirar.com.sa","Landing page with a demo request form. No actual service described yet.\nThis means the AI service is in development — not launched.","Their AI effort is EARLY. You're not late — you might be the first serious AI testing partner they talk to.",PURPLE),
    ("Feb 2024","IBM partnership for GenAI + data security","Announced at LEAP 2024. Co-providing data security with generative AI.","They're partnering, not building in-house. That's the model: bring in specialists. That's you.",TEAL),
    ("Sep 2025","KAFD partnership for AI-powered cybersecurity","King Abdullah Financial District. AI for cyber defence.","Shows client demand is real — financial district wants AI security.", AMBER),
    ("Dec 2024","SpiderSilk partnership — GCC's first cyber-AI startup","Partnered at Black Hat MEA. SpiderSilk does attack surface management with AI.","They partner with AI companies. They don't try to build it all themselves.",GREEN),
]
for i,(date,what,detail,insight,col) in enumerate(moves):
    y=Inches(1.35)+Inches(1.4)*i
    box(sl,Inches(0.6),y,Inches(1.6),Inches(1.2),fill=col)
    txt(sl,Inches(0.65),y+Inches(0.35),Inches(1.5),Inches(0.5),date,sz=16,bold=True,col=BG,align=PP_ALIGN.CENTER)
    txt(sl,Inches(2.4),y+Inches(0.05),Inches(4.5),Inches(0.3),what,sz=14,bold=True,col=col)
    txt(sl,Inches(2.4),y+Inches(0.38),Inches(4.5),Inches(0.7),detail,sz=11,col=GREY)
    txt(sl,Inches(7.2),y+Inches(0.05),Inches(5.5),Inches(1.1),insight,sz=13,col=LIGHT)

foot(sl,6)


# ═══════════════════════════════════════════
# 7  REGULATORY LANDSCAPE
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title(sl,"The Regulatory Landscape","Compliance is their bread and butter. This is how MTCP connects.")

box(sl,Inches(0.6),Inches(1.4),Inches(5.8),Inches(3.2),border=AMBER)
txt(sl,Inches(0.9),Inches(1.55),Inches(5),Inches(0.4),"WHAT SIRAR ALREADY DOES",sz=16,bold=True,col=AMBER)
regs=[
    "NCA ECC — National Cybersecurity Authority controls",
    "NCA CSCC — Critical systems cybersecurity controls",
    "NCA CCC — Cloud computing cybersecurity controls",
    "SAMA CSF — Saudi central bank cybersecurity framework",
    "ISO 27001 — Information security standard",
    "Aramco 3rd-party audit requirements",
]
for i,r in enumerate(regs):
    txt(sl,Inches(0.9),Inches(2.05)+Inches(0.37)*i,Inches(5.2),Inches(0.3),f"  {r}",sz=13,col=LIGHT)

box(sl,Inches(6.7),Inches(1.4),Inches(5.8),Inches(3.2),border=GREEN)
txt(sl,Inches(7.0),Inches(1.55),Inches(5),Inches(0.4),"WHAT'S COMING (THE GAP)",sz=16,bold=True,col=GREEN)
coming=[
    "EU AI Act (Aug 2026) — affects anyone selling to Europe",
    "Saudi NDMO AI governance — being developed now",
    "SDAIA national AI framework — Saudi's AI authority",
    "ISO 42001 — AI management systems (new standard)",
    "NIST AI RMF — US risk management framework",
    "No current standard for post-correction AI testing",
]
for i,r in enumerate(coming):
    txt(sl,Inches(7.0),Inches(2.05)+Inches(0.37)*i,Inches(5.2),Inches(0.3),f"  {r}",sz=13,col=LIGHT)

box(sl,Inches(0.6),Inches(4.85),Inches(12.1),Inches(1.6),fill=CARD2,border=PURPLE)
txt(sl,Inches(0.9),Inches(5.0),Inches(11.5),Inches(0.3),"THE OPPORTUNITY",sz=14,bold=True,col=PURPLE)
txt(sl,Inches(0.9),Inches(5.4),Inches(11.5),Inches(0.9),
    "sirar knows how to do cybersecurity compliance. They DON'T know how to do AI compliance.\n"
    "AI regulation is coming fast — EU AI Act in 4 months, Saudi frameworks being written right now.\n"
    "sirar's clients (banks, government, critical infra) will need AI assurance evidence.\n"
    "MTCP is the ONLY platform that produces this evidence with 184,387 evaluations behind it.\n"
    "You're not asking them to take a risk. You're offering them a head start.",
    sz=14,col=LIGHT)
foot(sl,7)


# ═══════════════════════════════════════════
# 8  HOW YOU FIT — THE THREE OPTIONS
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title(sl,"How You Fit Into sirar","Three possible shapes. The pilot tests which one works.")

options=[
    ("Option A: Testing Partner","sirar runs their AI Assurance service.\nYou provide MTCP evaluations behind the scenes.\nThey package your Evidence Packs for their clients.",
     "Like a lab that does the blood tests for a hospital.\nThe hospital talks to the patient. The lab does the science.",
     "MTCP stays yours. sirar pays per evaluation or monthly access.\nYou're a supplier. Clean and simple.",GREEN),
    ("Option B: Embedded SME + Platform","You sit inside sirar as a consultant/advisor.\nYou run evaluations AND interpret results for clients.\nYou help design their AI Assurance service.",
     "Like a specialist doctor who also consults at a hospital.\nYou bring your tools AND your expertise.",
     "Retainer/salary for your time. Platform access is separate.\nDO NOT let them merge these into one payment.",AMBER),
    ("Option C: Licence / White-Label","sirar licences MTCP to run it themselves.\nYour name may or may not be on it.\nThey pay a licence fee for access.",
     "Like licensing a pharmaceutical drug to a distributor.\nThey sell it. You get royalties.",
     "This is the BIG deal. Don't discuss numbers on this call.\nJust acknowledge it exists as an option for later.",PURPLE),
]
for i,(name,desc,analogy,commercial,col) in enumerate(options):
    l=Inches(0.6)+Inches(4.2)*i
    box(sl,l,Inches(1.35),Inches(3.9),Inches(5.2),border=col)
    txt(sl,l+Inches(0.2),Inches(1.5),Inches(3.5),Inches(0.4),name,sz=17,bold=True,col=col)
    txt(sl,l+Inches(0.2),Inches(2.0),Inches(3.5),Inches(1.2),desc,sz=12,col=LIGHT)
    txt(sl,l+Inches(0.2),Inches(3.3),Inches(3.5),Inches(0.3),"Plain English:",sz=11,bold=True,col=GREY)
    txt(sl,l+Inches(0.2),Inches(3.6),Inches(3.5),Inches(0.8),analogy,sz=11,col=GREY)
    txt(sl,l+Inches(0.2),Inches(4.6),Inches(3.5),Inches(0.3),"Commercial:",sz=11,bold=True,col=col)
    txt(sl,l+Inches(0.2),Inches(4.9),Inches(3.5),Inches(0.8),commercial,sz=11,col=LIGHT)

foot(sl,8)


# ═══════════════════════════════════════════
# 9  WHY THEY NEED YOU (NOT THE OTHER WAY AROUND)
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title(sl,"Why They Need You (Not the Other Way Around)","Read this before the call. Know your leverage.")

box(sl,Inches(0.6),Inches(1.4),Inches(5.8),Inches(5.0),border=GREEN)
txt(sl,Inches(0.9),Inches(1.55),Inches(5),Inches(0.4),"WHAT YOU HAVE THAT THEY DON'T",sz=16,bold=True,col=GREEN)
youhave=[
    "184,387 real evaluations — you can't fake this",
    "32 models tested from 13 providers",
    "A working platform (mtcp.live) — not a slide deck",
    "Published research with a DOI",
    "A grading system that makes sense to non-technical buyers",
    "Evidence Packs designed for EU AI Act compliance",
    "6 independent experts who validated your work",
    "The only post-correction persistence benchmark in existence",
    "Results that show flagship models fail (provocative, attention-grabbing)",
    "Mohamed Rihan at KFUPM already integrating your data into Saudi sovereign AI",
]
for i,item in enumerate(youhave):
    txt(sl,Inches(0.9),Inches(2.05)+Inches(0.4)*i,Inches(5.2),Inches(0.35),f"  {item}",sz=12,col=LIGHT)

box(sl,Inches(6.7),Inches(1.4),Inches(5.8),Inches(5.0),border=RED)
txt(sl,Inches(7.0),Inches(1.55),Inches(5),Inches(0.4),"WHAT THEY'D HAVE TO DO WITHOUT YOU",sz=16,bold=True,col=RED)
without=[
    "Build an AI testing methodology from scratch — 12+ months",
    "Run tens of thousands of evaluations — massive cost",
    "Publish research to establish credibility — years",
    "Get independent validation — can't be bought",
    "Build a platform — engineering team needed",
    "Design a grading system — requires deep expertise",
    "Map to EU AI Act — legal + technical crossover",
    "Build probe sets — the hardest part, takes months",
    "They'd be years behind where you are RIGHT NOW",
    "And they'd still not have your dataset or validation",
]
for i,item in enumerate(without):
    txt(sl,Inches(7.0),Inches(2.05)+Inches(0.4)*i,Inches(5.2),Inches(0.35),f"  {item}",sz=12,col=LIGHT)

foot(sl,9)


# ═══════════════════════════════════════════
# 10  THEIR LANGUAGE — SPEAK LIKE THEY DO
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title(sl,"Speak Their Language","Mirror their words. It builds instant rapport.")

box(sl,Inches(0.6),Inches(1.4),Inches(12.1),Inches(1.2),fill=CARD2)
txt(sl,Inches(0.9),Inches(1.55),Inches(11),Inches(0.3),"THEIR BRAND VALUES (D3)",sz=14,bold=True,col=TEAL)
txt(sl,Inches(0.9),Inches(1.9),Inches(11),Inches(0.6),
    "Dynamism (agile, collaborative)   |   Devotion (customer-centric, trustworthy)   |   Drive (proactive, inventive)",
    sz=16,col=LIGHT)

box(sl,Inches(0.6),Inches(2.85),Inches(12.1),Inches(3.6),border=PURPLE)
txt(sl,Inches(0.9),Inches(3.0),Inches(11.5),Inches(0.4),"REFRAME YOUR LANGUAGE FOR THEM",sz=16,bold=True,col=PURPLE)

reframes=[
    ("Instead of...","Say...","Why it works for them"),
    ('"My platform"','"The MTCP assurance capability"',"Sounds like a service, not a personal project"),
    ('"I tested 32 models"','"The programme has evaluated 32 production models"',"Programme sounds established, not solo"),
    ('"Nobody passed"','"No model meets the 90% threshold for deployment certification"',"Certification language = their world"),
    ('"Here\'s the data"','"Here\'s the assurance evidence"',"Evidence is what regulators and clients want"),
    ('"I can help"','"MTCP can strengthen the service launch"',"Makes it about the platform, not dependency on you"),
    ('"It\'s cheap to run"','"The pilot is low-integration"',"They care about effort, not just money"),
    ('"I\'ll send you stuff"','"I\'ll prepare a pilot proposal"',"Professional, structured, business-like"),
]
for i,(old,new,why) in enumerate(reframes):
    y=Inches(3.45)+Inches(0.45)*i
    col_text = GREY if i==0 else LIGHT
    b = i==0
    txt(sl,Inches(0.9),y,Inches(3.2),Inches(0.35),old,sz=13,bold=b,col=RED if i>0 else GREY)
    txt(sl,Inches(4.3),y,Inches(3.8),Inches(0.35),new,sz=13,bold=b,col=GREEN if i>0 else GREY)
    txt(sl,Inches(8.3),y,Inches(4),Inches(0.35),why,sz=11,col=GREY)

foot(sl,10)


# ═══════════════════════════════════════════
# 11  THE SAUDI CONTEXT
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title(sl,"The Saudi Context — What You Should Know","Background that helps you understand why this matters to them.")

box(sl,Inches(0.6),Inches(1.35),Inches(5.8),Inches(3.0),border=AMBER)
txt(sl,Inches(0.9),Inches(1.5),Inches(5),Inches(0.4),"VISION 2030",sz=16,bold=True,col=AMBER)
txt(sl,Inches(0.9),Inches(1.95),Inches(5.2),Inches(2.3),
    "Saudi Arabia's plan to diversify beyond oil.\n"
    "Massive investment in tech, AI, smart cities.\n\n"
    "NEOM — $500B smart city project.\n"
    "Riyadh — becoming a global tech hub.\n"
    "SDAIA — Saudi Data & AI Authority.\n"
    "NDMO — National Data Management Office.\n\n"
    "AI governance is a national priority.\n"
    "They want to lead, not follow.",
    sz=13,col=LIGHT)

box(sl,Inches(6.7),Inches(1.35),Inches(5.8),Inches(3.0),border=TEAL)
txt(sl,Inches(7.0),Inches(1.5),Inches(5),Inches(0.4),"THE CYBER MARKET",sz=16,bold=True,col=TEAL)
txt(sl,Inches(7.0),Inches(1.95),Inches(5.2),Inches(2.3),
    "Saudi cybersecurity market: ~$3.7B (2024).\n"
    "Growing 13% per year.\n\n"
    "NCA (National Cybersecurity Authority) regulates.\n"
    "sirar has Tier 1 licence — one of only six.\n\n"
    "Every major project needs cyber compliance.\n"
    "AI compliance is next. It's coming fast.\n\n"
    "Whoever gets there first owns the market.",
    sz=13,col=LIGHT)

box(sl,Inches(0.6),Inches(4.6),Inches(12.1),Inches(1.8),fill=CARD2,border=GREEN)
txt(sl,Inches(0.9),Inches(4.75),Inches(11.5),Inches(0.3),"YOUR EXISTING SAUDI CONNECTION",sz=14,bold=True,col=GREEN)
txt(sl,Inches(0.9),Inches(5.15),Inches(11.5),Inches(1.1),
    "Mohamed Rihan — Sovereign Intelligence Systems Architect at KFUPM (King Fahd University).\n"
    "He's already integrating MTCP data into R-AGAM, his sovereign AI governance system.\n"
    "He sent you a formal V1.0 interface spec making MTCP a required input.\n\n"
    "You can mention this to Syed: \"MTCP is already being integrated into sovereign AI governance work in Saudi Arabia.\"\n"
    "Don't name Mohamed unless he's OK with it. Just reference the fact.",
    sz=13,col=LIGHT)
foot(sl,11)


# ═══════════════════════════════════════════
# 12  SMARTER QUESTIONS
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title(sl,"Smarter Questions — Now You Know Who They Are","These questions show you did your homework.")

qs=[
    ("\"I noticed sirar has strong GRC and advisory capabilities already.\nHow does the AI Assurance service relate to the existing PLAN pillar?\"",
     "Shows you researched their service structure. Makes him explain how AI fits in.","Use this early.",PURPLE),
    ("\"Are the first AI Assurance clients likely to come from your existing\ncyber clients — banking, government, critical infrastructure?\"",
     "Shows you know their client base. Tells you who your Evidence Packs would be for.","Very useful for scoping the pilot.",TEAL),
    ("\"Is there pressure from NCA or SDAIA on AI-specific compliance\nrequirements, or is this more market-driven right now?\"",
     "Shows you understand Saudi regulatory landscape. Tells you whether they need compliance evidence NOW or are getting ahead.","Sets up the EU AI Act conversation.",AMBER),
    ("\"Would the AI Assurance service be positioned for Saudi domestic\nclients, international clients, or both?\"",
     "If international: EU AI Act matters. If domestic only: Saudi frameworks matter more.","Changes how you pitch MTCP.",GREEN),
    ("\"Is sirar looking to build AI testing capability in-house, or partner\nwith a specialist who already has the platform and data?\"",
     "THE most important question. His answer tells you exactly what commercial model fits.","The answer shapes everything.",PINK),
    ("\"What would success look like for the AI Assurance service in its\nfirst year?\"",
     "Makes him paint the picture. Then you show how MTCP accelerates getting there.","Use near the end.",GREEN),
]
for i,(q,why,when,col) in enumerate(qs):
    y=Inches(1.3)+Inches(0.95)*i
    box(sl,Inches(0.6),y,Inches(8.2),Inches(0.85),border=col)
    txt(sl,Inches(0.8),y+Inches(0.05),Inches(7.8),Inches(0.75),q,sz=12,col=WHITE)
    txt(sl,Inches(9.0),y+Inches(0.05),Inches(3.8),Inches(0.4),why,sz=11,col=GREY)
    txt(sl,Inches(9.0),y+Inches(0.55),Inches(3.8),Inches(0.25),when,sz=11,bold=True,col=col)

foot(sl,12)


# ═══════════════════════════════════════════
out=os.path.expanduser("~/Desktop/SIRAR_Intel_and_Fit.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
