from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(8.5)
prs.slide_height = Inches(11)  # portrait A4-ish — reads like a document

BG=RGBColor(0xFF,0xFF,0xFF)
DARK=RGBColor(0x0D,0x0D,0x12)
CARD=RGBColor(0xF4,0xF4,0xF8)
CARD2=RGBColor(0xE8,0xE8,0xF0)
ACCENT=RGBColor(0x1A,0x1A,0x2E)
PURPLE=RGBColor(0x6C,0x5C,0xE7)
TEAL=RGBColor(0x00,0x96,0x96)
BLACK=RGBColor(0x1A,0x1A,0x1A)
BODY=RGBColor(0x33,0x33,0x44)
GREY=RGBColor(0x66,0x66,0x77)
LGREY=RGBColor(0x99,0x99,0xAA)
RED=RGBColor(0xC0,0x39,0x2B)
GREEN=RGBColor(0x27,0xAE,0x60)
AMBER=RGBColor(0xD4,0x8A,0x0B)
WHITE=RGBColor(0xFF,0xFF,0xFF)

FONT="Calibri"
MONO="Courier New"

def bgw(s): f=s.background.fill; f.solid(); f.fore_color.rgb=BG
def bgd(s): f=s.background.fill; f.solid(); f.fore_color.rgb=DARK
def box(s,l,t,w,h,fill=CARD,border=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if border: sh.line.color.rgb=border; sh.line.width=Pt(0.75)
    else: sh.line.fill.background()
    sh.adjustments[0]=0.03; return sh
def rect(s,l,t,w,h,fill=CARD):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.fill.background(); return sh
def txt(s,l,t,w,h,text,sz=11,col=BODY,bold=False,align=PP_ALIGN.LEFT,font=FONT):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(sz); p.font.color.rgb=col
    p.font.bold=bold; p.font.name=font; p.alignment=align; return tb
def line(s,l,t,w,col=PURPLE):
    rect(s,l,t,w,Pt(2),fill=col)
def footer(s):
    line(s,Inches(0.6),Inches(10.3),Inches(7.3),CARD2)
    txt(s,Inches(0.6),Inches(10.35),Inches(5),Inches(0.3),
        "MTCP Research Programme  |  A. Abby  |  mtcp.live",sz=7,col=LGREY)
    txt(s,Inches(5.5),Inches(10.35),Inches(2.4),Inches(0.3),
        "DOI: 10.17605/OSF.IO/DXGK5",sz=7,col=LGREY,align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
# PAGE 1 — COVER
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bgd(sl)
rect(sl,Inches(0),Inches(0),Inches(8.5),Inches(11),fill=DARK)
line(sl,Inches(0.8),Inches(2.5),Inches(0.8),PURPLE)

txt(sl,Inches(0.8),Inches(3.0),Inches(6.5),Inches(0.5),
    "MTCP",sz=48,bold=True,col=WHITE)
txt(sl,Inches(0.8),Inches(3.7),Inches(6.5),Inches(0.4),
    "Multi-Turn Constraint Persistence",sz=18,col=PURPLE)
txt(sl,Inches(0.8),Inches(4.3),Inches(6.5),Inches(0.4),
    "Platform Overview and Engagement Proposal",sz=14,col=RGBColor(0x99,0x99,0xAA))

txt(sl,Inches(0.8),Inches(5.5),Inches(6.5),Inches(0.3),
    "Prepared for sirar by stc",sz=12,col=WHITE)
txt(sl,Inches(0.8),Inches(5.9),Inches(6.5),Inches(0.3),
    "April 2026",sz=11,col=LGREY)

# stats strip
rect(sl,Inches(0.8),Inches(7.0),Inches(6.9),Inches(1.6),fill=RGBColor(0x16,0x16,0x1F))
stats=[("184,387","evaluations"),("32","models"),("13","providers"),("532","probes"),("0","models passed")]
for i,(num,label) in enumerate(stats):
    x=Inches(0.9)+Inches(1.38)*i
    txt(sl,x,Inches(7.15),Inches(1.2),Inches(0.4),num,sz=18,bold=True,col=PURPLE,font=MONO,align=PP_ALIGN.CENTER)
    txt(sl,x,Inches(7.6),Inches(1.2),Inches(0.3),label,sz=8,col=LGREY,align=PP_ALIGN.CENTER)

txt(sl,Inches(0.8),Inches(9.0),Inches(6.5),Inches(0.25),
    "mtcp.live  |  doi.org/10.17605/OSF.IO/DXGK5",sz=9,col=LGREY)
txt(sl,Inches(0.8),Inches(9.3),Inches(6.5),Inches(0.25),
    "A. Abby  |  MTCP Research Programme  |  All Rights Reserved 2026",sz=8,col=RGBColor(0x55,0x55,0x66))


# ═══════════════════════════════════════════
# PAGE 2 — WHAT MTCP IS
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bgw(sl)

txt(sl,Inches(0.6),Inches(0.5),Inches(7),Inches(0.35),
    "What Is MTCP?",sz=22,bold=True,col=BLACK)
line(sl,Inches(0.6),Inches(0.9),Inches(1.5),PURPLE)

txt(sl,Inches(0.6),Inches(1.1),Inches(7.3),Inches(1.2),
    "MTCP (Multi-Turn Constraint Persistence) is a black-box AI assurance platform that measures "
    "whether large language models maintain compliance with stated behavioural constraints after "
    "explicit correction.\n\n"
    "Unlike red-teaming (which tests whether a model can fail) or monitoring (which watches after "
    "deployment), MTCP tests whether correction actually works: if an AI breaks a rule and is told "
    "to stop, does it actually stop — or does the failure persist?",
    sz=11,col=BODY)

txt(sl,Inches(0.6),Inches(2.5),Inches(7),Inches(0.3),
    "How It Works",sz=16,bold=True,col=BLACK)

steps=[
    ("1. Set constraint","Define a behavioural rule the model must follow."),
    ("2. Apply pressure","Conduct a conversation that tests the rule under realistic conditions."),
    ("3. Detect violation","The platform automatically identifies when the model breaks the rule."),
    ("4. Correct","Issue an explicit correction: \"You broke the rule. Return to compliance.\""),
    ("5. Test persistence","Observe whether the correction holds or the model reverts to non-compliance."),
    ("6. Produce evidence","Generate a SHA-256 signed Evidence Pack documenting the full evaluation."),
]
for i,(step,desc) in enumerate(steps):
    y=Inches(2.9)+Inches(0.32)*i
    txt(sl,Inches(0.6),y,Inches(1.5),Inches(0.25),step,sz=10,bold=True,col=PURPLE)
    txt(sl,Inches(2.2),y,Inches(5.7),Inches(0.25),desc,sz=10,col=BODY)

txt(sl,Inches(0.6),Inches(4.9),Inches(7),Inches(0.3),
    "Key Findings",sz=16,bold=True,col=BLACK)

findings=[
    "No model achieves Grade A (90%+). The best performer, grok-3-mini (xAI), reaches 88.7% (Grade B).",
    "GPT-4o scores 16.2 percentage points below GPT-4o-mini — flagship models are not necessarily safer.",
    "14 of 32 models exhibit architectural failure: temperature controls and operational settings cannot improve their performance.",
    "Open-source models outperform commercial alternatives on constraint persistence.",
    "Safety-tuned models (Claude, Gemini) show systematic, temperature-invariant failure — the issue is structural, not stochastic.",
]
for i,f in enumerate(findings):
    y=Inches(5.3)+Inches(0.34)*i
    txt(sl,Inches(0.6),y,Inches(7.3),Inches(0.3),f"  {f}",sz=10,col=BODY)

txt(sl,Inches(0.6),Inches(7.2),Inches(7),Inches(0.3),
    "Research Credentials",sz=16,bold=True,col=BLACK)

txt(sl,Inches(0.6),Inches(7.6),Inches(7.3),Inches(1.5),
    "5 published research papers on OSF with registered DOI (10.17605/OSF.IO/DXGK5). "
    "22 additional completed extension papers covering failure taxonomy, regulatory mapping, "
    "audit methodology, and theoretical foundations.\n\n"
    "The methodology has been independently validated by six practitioners within 48 hours of "
    "publication, including formal citation in an enterprise Architecture Decision Record and "
    "integration into a sovereign AI governance system via formal interface specification (V1.0).\n\n"
    "The MTCP dataset (500 boundary evaluations) is publicly available on HuggingFace for independent verification.",
    sz=10,col=BODY)

footer(sl)


# ═══════════════════════════════════════════
# PAGE 3 — SCORING METHODOLOGY
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bgw(sl)

txt(sl,Inches(0.6),Inches(0.5),Inches(7),Inches(0.35),
    "Scoring Methodology",sz=22,bold=True,col=BLACK)
line(sl,Inches(0.6),Inches(0.9),Inches(1.5),PURPLE)

# BIS
txt(sl,Inches(0.6),Inches(1.1),Inches(7),Inches(0.3),
    "BIS — Benchmark Index Score",sz=14,bold=True,col=PURPLE)
txt(sl,Inches(0.6),Inches(1.45),Inches(7.3),Inches(0.9),
    "BIS is the primary scoring metric. It is the mean pass rate across evaluations at four "
    "temperature settings: T=0.0, T=0.2, T=0.5, and T=0.8. At each temperature, 200 probes are "
    "administered. Each probe is a structured three-turn conversation: Turn 1 establishes a "
    "constrained task, Turn 2 issues correction if violated, Turn 3 reinforces if still violated. "
    "A probe passes when the model achieves compliance at the terminal turn. BIS determines the grade.",
    sz=10,col=BODY)

# Grade table
txt(sl,Inches(0.6),Inches(2.5),Inches(7),Inches(0.25),
    "MTCP Grading Scale",sz=12,bold=True,col=BLACK)

box(sl,Inches(0.6),Inches(2.8),Inches(7.3),Inches(0.3),fill=ACCENT)
txt(sl,Inches(0.7),Inches(2.83),Inches(0.5),Inches(0.22),"Grade",sz=9,bold=True,col=WHITE)
txt(sl,Inches(1.3),Inches(2.83),Inches(1),Inches(0.22),"BIS",sz=9,bold=True,col=WHITE)
txt(sl,Inches(2.4),Inches(2.83),Inches(2.2),Inches(0.22),"Deployment Readiness",sz=9,bold=True,col=WHITE)
txt(sl,Inches(4.8),Inches(2.83),Inches(3),Inches(0.22),"Current Dataset",sz=9,bold=True,col=WHITE)

grades=[
    ("A+","95.0%+","Deployment-ready","No model achieves this grade",GREEN),
    ("A","90.0-94.9%","Strong reliability","No model achieves this grade",GREEN),
    ("B","80.0-89.9%","Certified with monitoring","grok-3-mini (88.7%), GPT-3.5 (83.2%), GPT-4o-mini (81.1%)",TEAL),
    ("C","70.0-79.9%","Review required","LLaMA 3.3 70B (75.6%)",AMBER),
    ("D","60.0-69.9%","Not recommended","GPT-4o (64.9%), DeepSeek-R1 (62.4%) — most common grade",RED),
    ("F","<60.0%","Failing","Claude Sonnet 4.5 (59.4%), Gemini Flash (59.5%)",RED),
]
for i,(g,bis,ready,dataset,col) in enumerate(grades):
    y=Inches(3.12)+Inches(0.28)*i
    if i%2==0: box(sl,Inches(0.6),y-Inches(0.02),Inches(7.3),Inches(0.28),fill=CARD)
    txt(sl,Inches(0.7),y,Inches(0.5),Inches(0.22),g,sz=10,bold=True,col=col)
    txt(sl,Inches(1.3),y,Inches(1),Inches(0.22),bis,sz=9,col=BODY)
    txt(sl,Inches(2.4),y,Inches(2.2),Inches(0.22),ready,sz=9,col=BODY)
    txt(sl,Inches(4.8),y,Inches(3),Inches(0.22),dataset,sz=8,col=GREY)

# Temperature
txt(sl,Inches(0.6),Inches(5.0),Inches(7),Inches(0.3),
    "Temperature and Failure Classification",sz=14,bold=True,col=PURPLE)
txt(sl,Inches(0.6),Inches(5.35),Inches(7.3),Inches(0.7),
    "Temperature controls the randomness of AI responses (0.0 = deterministic, 0.8 = high variance). "
    "By testing at four settings, MTCP classifies each model's failure pattern:",
    sz=10,col=BODY)

patterns=[
    ("Stochastic (12 models, 38%)","Performance degrades with higher temperature. Fixable with operational controls."),
    ("Architectural (14 models, 44%)","Performance unchanged across temperatures. Unfixable without retraining. The failure is structural."),
    ("Genuine Persistence (2 models, 6%)","Low temperature variance and low CPD. The model genuinely maintains constraints."),
    ("Atypical (2 models, 6%)","Irregular behaviour requiring individual investigation."),
]
for i,(name,desc) in enumerate(patterns):
    y=Inches(6.15)+Inches(0.35)*i
    txt(sl,Inches(0.6),y,Inches(2.5),Inches(0.25),name,sz=9,bold=True,col=PURPLE)
    txt(sl,Inches(3.2),y,Inches(4.7),Inches(0.3),desc,sz=9,col=BODY)

# CPD
txt(sl,Inches(0.6),Inches(7.7),Inches(7),Inches(0.3),
    "CPD — Control Performance Degradation",sz=14,bold=True,col=PURPLE)
txt(sl,Inches(0.6),Inches(8.05),Inches(7.3),Inches(1.0),
    "CPD measures whether a model's score reflects genuine capability or test familiarity. A separate "
    "set of 20 control probes (novel content the model has not seen) is administered. "
    "CPD = Primary score minus Control score. A large negative CPD indicates the primary score is "
    "inflated by benchmark familiarity.\n\n"
    "Ranges: Minimal (0 to -10pp, score reliable), Moderate (-10 to -30pp, requires disclosure), "
    "Substantial (-30 to -50pp, unreliable alone), Severe (below -50pp, not valid compliance evidence).\n\n"
    "Example: grok-3-mini achieves Grade B (88.7%) but has CPD of -58.7pp (Severe). "
    "DeepSeek-R1 achieves Grade D (62.4%) with CPD of -3.7pp (Minimal) — the most honest score in the dataset.",
    sz=10,col=BODY)

# Ve
txt(sl,Inches(0.6),Inches(9.3),Inches(7),Inches(0.3),
    "Ve — Veterance Metric",sz=14,bold=True,col=PURPLE)
txt(sl,Inches(0.6),Inches(9.6),Inches(7.3),Inches(0.6),
    "Ve counts consecutive corrections the model ignores. Ve=0 (corrected on first attempt), "
    "Ve=1 (one ignored correction), Ve=2 (hard stop — all corrections ignored, complete failure). "
    "Computed entirely from transcripts. Produces a single integer usable as a real-time monitoring "
    "signal for post-deployment behavioural oversight.",
    sz=10,col=BODY)

footer(sl)


# ═══════════════════════════════════════════
# PAGE 4 — EVALUATION VECTORS + EU AI ACT
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bgw(sl)

txt(sl,Inches(0.6),Inches(0.5),Inches(7),Inches(0.35),
    "Evaluation Scope",sz=22,bold=True,col=BLACK)
line(sl,Inches(0.6),Inches(0.9),Inches(1.5),PURPLE)

txt(sl,Inches(0.6),Inches(1.1),Inches(7),Inches(0.25),
    "Five Evaluation Vectors",sz=14,bold=True,col=PURPLE)
txt(sl,Inches(0.6),Inches(1.4),Inches(7.3),Inches(0.4),
    "MTCP evaluates constraint persistence across five distinct categories of behavioural rules, "
    "ensuring coverage of the main failure surfaces in production AI deployment:",
    sz=10,col=BODY)

vectors=[
    ("NCA","Negative Constraint Adherence","Content exclusion rules — \"do not mention X\", \"avoid topic Y\""),
    ("SFC","Structural Format Compliance","Output format rules — bullet points, numbered lists, word limits"),
    ("IDL","Information Density and Length","Scope rules — \"brief summary only\", \"stay high-level\""),
    ("CG","Contextual Grounding","Domain restriction rules — \"only discuss finance\", \"stay within expertise\""),
    ("LANG","Multilingual Consistency","Language rules — \"respond in English only\", \"maintain formal register\""),
]
for i,(code,name,desc) in enumerate(vectors):
    y=Inches(2.0)+Inches(0.3)*i
    txt(sl,Inches(0.6),y,Inches(0.6),Inches(0.22),code,sz=10,bold=True,col=PURPLE,font=MONO)
    txt(sl,Inches(1.3),y,Inches(2.3),Inches(0.22),name,sz=10,bold=True,col=BLACK)
    txt(sl,Inches(3.7),y,Inches(4.2),Inches(0.25),desc,sz=9,col=BODY)

txt(sl,Inches(0.6),Inches(3.7),Inches(7.3),Inches(0.3),
    "532 probes across these five vectors. Probe content is proprietary and not disclosed to prevent benchmark gaming.",
    sz=9,col=GREY)

# EU AI Act
txt(sl,Inches(0.6),Inches(4.3),Inches(7),Inches(0.35),
    "EU AI Act Compliance Mapping",sz=22,bold=True,col=BLACK)
line(sl,Inches(0.6),Inches(4.7),Inches(1.5),PURPLE)

txt(sl,Inches(0.6),Inches(4.9),Inches(7.3),Inches(0.5),
    "MTCP was designed with EU AI Act compliance in mind. Evidence Packs are formatted for Annex IV "
    "technical documentation requirements. The compliance deadline is August 2026.",
    sz=10,col=BODY)

box(sl,Inches(0.6),Inches(5.5),Inches(7.3),Inches(0.3),fill=ACCENT)
txt(sl,Inches(0.7),Inches(5.53),Inches(1.2),Inches(0.22),"Article",sz=9,bold=True,col=WHITE)
txt(sl,Inches(2.0),Inches(5.53),Inches(2),Inches(0.22),"Obligation",sz=9,bold=True,col=WHITE)
txt(sl,Inches(4.2),Inches(5.53),Inches(3.5),Inches(0.22),"MTCP Component",sz=9,bold=True,col=WHITE)

arts=[
    ("Article 9","Risk management","BIS grade + failure pattern classification"),
    ("Article 13","Transparency","Published methodology + public leaderboard"),
    ("Article 17","Quality management","Sigma-Forensics audit reports"),
    ("Article 61","Post-market monitoring","Ve-based real-time monitoring system"),
    ("Article 72","Serious incident reporting","Ve escalation + Sigma-Forensics incident reports"),
    ("Annex IX","Technical documentation","SHA-256 signed Evidence Packs"),
]
for i,(art,obl,comp) in enumerate(arts):
    y=Inches(5.82)+Inches(0.28)*i
    if i%2==0: box(sl,Inches(0.6),y-Inches(0.02),Inches(7.3),Inches(0.28),fill=CARD)
    txt(sl,Inches(0.7),y,Inches(1.2),Inches(0.22),art,sz=9,bold=True,col=BODY)
    txt(sl,Inches(2.0),y,Inches(2),Inches(0.22),obl,sz=9,col=BODY)
    txt(sl,Inches(4.2),y,Inches(3.5),Inches(0.22),comp,sz=9,col=TEAL)

txt(sl,Inches(0.6),Inches(7.7),Inches(7.3),Inches(0.5),
    "Saudi AI governance frameworks under development by SDAIA and NDMO draw on EU AI Act structure. "
    "MTCP compliance mapping is directly transferable to emerging Saudi regulatory requirements. "
    "MTCP measurement data is already being integrated into sovereign AI governance work in the Gulf region.",
    sz=10,col=BODY)

# DELIVERABLES
txt(sl,Inches(0.6),Inches(8.5),Inches(7),Inches(0.3),
    "Deliverables — What Clients Receive",sz=14,bold=True,col=PURPLE)

deliverables=[
    ("Evidence Pack","SHA-256 signed document per model. Contains: model ID, temperature, probe manifest hash, per-probe outcomes, BIS, grade. Tamper-evident. Audit-ready."),
    ("MTCP Grade","Single letter grade (A+ to F) based on BIS. Instantly understandable for non-technical stakeholders, procurement teams, and regulators."),
    ("Executive Summary","Plain-language report covering risk findings, model comparison, and release-readiness recommendation."),
    ("Failure Classification","Each model categorised as stochastic (fixable) or architectural (unfixable). Determines remediation pathway."),
    ("Deployment Recommendation","Per-model guidance: deploy, deploy with monitoring, or do not deploy in constrained contexts."),
]
for i,(name,desc) in enumerate(deliverables):
    y=Inches(8.85)+Inches(0.28)*i
    txt(sl,Inches(0.6),y,Inches(1.8),Inches(0.22),name,sz=9,bold=True,col=BLACK)
    txt(sl,Inches(2.5),y,Inches(5.4),Inches(0.25),desc,sz=8,col=BODY)

footer(sl)


# ═══════════════════════════════════════════
# PAGE 5 — ENGAGEMENT MODEL
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bgw(sl)

txt(sl,Inches(0.6),Inches(0.5),Inches(7),Inches(0.35),
    "Engagement Model",sz=22,bold=True,col=BLACK)
line(sl,Inches(0.6),Inches(0.9),Inches(1.5),PURPLE)

# SME
txt(sl,Inches(0.6),Inches(1.1),Inches(7),Inches(0.3),
    "1. Subject Matter Expert Role",sz=14,bold=True,col=PURPLE)
txt(sl,Inches(0.6),Inches(1.45),Inches(7.3),Inches(0.6),
    "MTCP is not only a platform — it requires expert interpretation. The SME role provides sirar "
    "with methodology design, results interpretation, client-facing framing, regulatory mapping, "
    "and service design support. This is the human expertise layer that turns platform output into "
    "a commercially deployable assurance service.",
    sz=10,col=BODY)

sme_items=[
    ("Evaluation design","Define scope, model selection, and success criteria for each engagement"),
    ("Results interpretation","Explain grades, failure patterns, CPD, and Ve findings to clients"),
    ("Report production","Executive summaries, risk assessments, and deployment recommendations"),
    ("Client framing","Help package MTCP evidence for sirar's client-facing service"),
    ("Regulatory mapping","EU AI Act, Saudi NCA/SDAIA, ISO 42001 compliance guidance"),
    ("Service design","Shape the AI Assurance service workflow and offering structure"),
    ("Team training","Train sirar staff on interpreting and presenting MTCP outputs"),
]
for i,(role,desc) in enumerate(sme_items):
    y=Inches(2.2)+Inches(0.26)*i
    txt(sl,Inches(0.6),y,Inches(2),Inches(0.2),role,sz=9,bold=True,col=TEAL)
    txt(sl,Inches(2.7),y,Inches(5.2),Inches(0.2),desc,sz=9,col=BODY)

box(sl,Inches(0.6),Inches(4.1),Inches(7.3),Inches(0.5),fill=CARD,border=PURPLE)
txt(sl,Inches(0.8),Inches(4.15),Inches(6.9),Inches(0.4),
    "The SME role is compensated via retainer or salary and covers ongoing work. "
    "It is commercially separate from platform access or licensing.",
    sz=9,bold=True,col=PURPLE)

# LICENSING
txt(sl,Inches(0.6),Inches(4.9),Inches(7),Inches(0.3),
    "2. Platform Licensing",sz=14,bold=True,col=PURPLE)
txt(sl,Inches(0.6),Inches(5.25),Inches(7.3),Inches(0.6),
    "The MTCP platform — including evaluation engine, proprietary probe sets, constraint detector, "
    "Evidence Pack generator, grading system, and research estate — is a pre-existing commercial "
    "asset with independent value. Platform access is available under defined licensing arrangements:",
    sz=10,col=BODY)

lic_models=[
    ("Evaluation-as-a-Service","sirar requests evaluations. MTCP delivers Evidence Packs and reports. "
     "Simplest model. Per-evaluation or monthly fee. MTCP retains full operational control."),
    ("Platform Access","sirar accesses MTCP through a controlled interface to run evaluations directly. "
     "Requires access controls, usage tracking, and data handling terms."),
    ("White-Label Licence","sirar licences MTCP to offer under their own brand to clients. "
     "Requires quality controls, revenue sharing, and territory scope definition."),
    ("Strategic Integration","MTCP is embedded as a component in sirar's service architecture. "
     "Deepest model. Requires formal commercial structuring and IP boundary definition."),
]
for i,(model,desc) in enumerate(lic_models):
    y=Inches(6.0)+Inches(0.52)*i
    txt(sl,Inches(0.6),y,Inches(7.3),Inches(0.2),model,sz=10,bold=True,col=BLACK)
    txt(sl,Inches(0.6),y+Inches(0.2),Inches(7.3),Inches(0.3),desc,sz=9,col=BODY)

box(sl,Inches(0.6),Inches(8.15),Inches(7.3),Inches(0.5),fill=CARD,border=AMBER)
txt(sl,Inches(0.8),Inches(8.2),Inches(6.9),Inches(0.4),
    "Licensing terms (scope, territory, duration, exclusivity, revenue) are defined through "
    "commercial agreement and are separate from the SME engagement.",
    sz=9,bold=True,col=AMBER)

# VALUE
txt(sl,Inches(0.6),Inches(9.0),Inches(7),Inches(0.3),
    "3. Asset Value Context",sz=14,bold=True,col=PURPLE)
txt(sl,Inches(0.6),Inches(9.3),Inches(7.3),Inches(0.8),
    "Estimated replacement cost to build equivalent capability from scratch: $1-2M+ and 18-24 months "
    "minimum (research team, evaluation infrastructure, API costs, probe development, academic "
    "publishing, regulatory mapping, independent validation, brand and credibility building). "
    "The MTCP research programme represents a fully operational platform with published methodology, "
    "verified dataset, live infrastructure, and independent third-party validation — available now.",
    sz=10,col=BODY)

footer(sl)


# ═══════════════════════════════════════════
# PAGE 6 — PILOT + LINKS
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bgw(sl)

txt(sl,Inches(0.6),Inches(0.5),Inches(7),Inches(0.35),
    "Recommended Next Step: Focused Pilot",sz=22,bold=True,col=BLACK)
line(sl,Inches(0.6),Inches(0.9),Inches(1.5),PURPLE)

txt(sl,Inches(0.6),Inches(1.1),Inches(7.3),Inches(0.5),
    "A controlled commercial engagement to evaluate MTCP output against sirar's service requirements "
    "before committing to broader licensing or integration.",
    sz=11,col=BODY)

box(sl,Inches(0.6),Inches(1.7),Inches(7.3),Inches(2.2),fill=CARD)
params=[
    ("Scope","3-5 models selected by sirar (client-relevant or internal systems)"),
    ("Method","Full MTCP evaluation: 4 temperatures, 200 probes per temperature, BIS + CPD + Ve analysis"),
    ("Timeline","2-4 weeks from agreement"),
    ("Deliverables","Evidence Packs  |  Grades  |  Executive Summary  |  Failure Classification  |  Deployment Recommendation"),
    ("Integration","Minimal — black-box. Requires only API access to target models. No weights, training data, or internal access."),
]
for i,(k,v) in enumerate(params):
    y=Inches(1.85)+Inches(0.33)*i
    txt(sl,Inches(0.8),y,Inches(1.2),Inches(0.22),k,sz=10,bold=True,col=PURPLE)
    txt(sl,Inches(2.1),y,Inches(5.6),Inches(0.25),v,sz=10,col=BODY)

txt(sl,Inches(0.6),Inches(4.1),Inches(7),Inches(0.25),
    "Success Criteria",sz=12,bold=True,col=BLACK)
criteria=[
    "Does MTCP produce assurance evidence that is useful for sirar's clients?",
    "Can the output be packaged into sirar's AI Assurance and Testing service?",
    "Is the evidence suitable for client-facing and regulator-facing reports?",
    "Does the grading system resonate with sirar's target market?",
]
for i,c in enumerate(criteria):
    txt(sl,Inches(0.6),Inches(4.4)+Inches(0.25)*i,Inches(7.3),Inches(0.22),
        f"  {c}",sz=10,col=BODY)

# LINKS
txt(sl,Inches(0.6),Inches(5.7),Inches(7),Inches(0.35),
    "Public Resources",sz=22,bold=True,col=BLACK)
line(sl,Inches(0.6),Inches(6.1),Inches(1.5),PURPLE)

links=[
    ("mtcp.live","Live Platform","Public leaderboard, dashboard, key findings, and model rankings across all 32 evaluated models."),
    ("mtcp.live/landing","Overview","Non-technical overview of MTCP for enterprise decision-makers. Results summary and enquiry form."),
    ("mtcp.live/leaderboard","Leaderboard","Ranked table of all models by MTCP grade, BIS score, and hard stop count."),
    ("mtcp.live/certificate","Certificates","Individual MTCP compliance certificates per model. Downloadable as PDF. SHA-256 verified."),
    ("mtcp.live/docs","API Documentation","Full OpenAPI specification. Demonstrates enterprise-grade technical architecture."),
    ("doi.org/10.17605/OSF.IO/DXGK5","Research Papers (DOI)","Published Papers I-V: Benchmark methodology, IGS theory, Sigma-Forensics audit framework, Evidence Pack spec, dataset documentation."),
    ("huggingface.co/datasets/aa8899/mtcp-boundary-500","Public Dataset","500 boundary evaluations. Downloadable and independently queryable. Used for external validation."),
]
for i,(url,name,desc) in enumerate(links):
    y=Inches(6.3)+Inches(0.42)*i
    txt(sl,Inches(0.6),y,Inches(3.5),Inches(0.2),url,sz=9,bold=True,col=PURPLE)
    txt(sl,Inches(0.6),y+Inches(0.18),Inches(3.5),Inches(0.2),name,sz=8,col=GREY)
    txt(sl,Inches(4.3),y,Inches(3.6),Inches(0.38),desc,sz=8,col=BODY)

# contact
box(sl,Inches(0.6),Inches(9.4),Inches(7.3),Inches(0.7),fill=ACCENT)
txt(sl,Inches(0.8),Inches(9.45),Inches(6.9),Inches(0.2),
    "MTCP Research Programme  |  A. Abby",sz=10,bold=True,col=WHITE)
txt(sl,Inches(0.8),Inches(9.7),Inches(6.9),Inches(0.3),
    "mtcp.live  |  DOI: 10.17605/OSF.IO/DXGK5  |  All Rights Reserved 2026",sz=9,col=RGBColor(0x99,0x99,0xAA))

footer(sl)


# ═══════════════════════════════════════════
out = os.path.expanduser("~/Desktop/MTCP_Overview_for_sirar.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
