from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG       = RGBColor(0x0D, 0x0D, 0x12)
CARD     = RGBColor(0x16, 0x16, 0x1F)
CARD2    = RGBColor(0x1C, 0x1C, 0x2A)
PURPLE   = RGBColor(0x6C, 0x5C, 0xE7)
TEAL     = RGBColor(0x00, 0xD2, 0xD3)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREY     = RGBColor(0x9B, 0x9B, 0xAA)
LIGHT    = RGBColor(0xCC, 0xCC, 0xDD)
RED      = RGBColor(0xE1, 0x50, 0x50)
GREEN    = RGBColor(0x00, 0xE6, 0x96)
AMBER    = RGBColor(0xFF, 0xB8, 0x42)
PINK     = RGBColor(0xFD, 0x79, 0xA8)
FONT = "Calibri"

def bg(s): f=s.background.fill; f.solid(); f.fore_color.rgb=BG
def box(s,l,t,w,h,fill=CARD,border=None):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if border: sh.line.color.rgb=border; sh.line.width=Pt(1)
    else: sh.line.fill.background()
    sh.adjustments[0]=0.04; return sh
def txt(s,l,t,w,h,text,sz=18,col=WHITE,bold=False,align=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(sz); p.font.color.rgb=col
    p.font.bold=bold; p.font.name=FONT; p.alignment=align; return tb
def head(s,t,sub=""):
    txt(s,Inches(0.6),Inches(0.3),Inches(12),Inches(0.55),t,sz=32,bold=True)
    if sub: txt(s,Inches(0.6),Inches(0.85),Inches(12),Inches(0.4),sub,sz=15,col=GREY)
def foot(s,n,tot):
    txt(s,Inches(0.6),Inches(7.1),Inches(4),Inches(0.25),"PRINT THIS OUT — keep it next to you on the call",sz=8,col=RGBColor(0x44,0x44,0x55))
    txt(s,Inches(12.0),Inches(7.1),Inches(1),Inches(0.25),f"{n}/{tot}",sz=8,col=GREY,align=PP_ALIGN.RIGHT)

TOT = 18

# ═══════════════════════════════════════════
# 1  COVER
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
box(sl,Inches(0.6),Inches(0.8),Inches(0.1),Inches(1.5),fill=PINK)
txt(sl,Inches(1.0),Inches(0.8),Inches(11),Inches(0.7),"YOUR CALL CHEAT SHEET",sz=46,bold=True)
txt(sl,Inches(1.0),Inches(1.6),Inches(11),Inches(0.5),"Everything about MTCP, sirar, and the call — from scratch.",sz=22,col=GREY)

txt(sl,Inches(1.0),Inches(2.5),Inches(11),Inches(0.4),"What's in here:",sz=18,bold=True,col=TEAL)
toc = [
    "Slides 2-5     MTCP explained from zero — what it is, what it does, what it found",
    "Slides 6-7     sirar explained — who they are, who Syed is, what they want",
    "Slides 8-9     How you fit together — why this makes sense",
    "Slides 10-14   Every question Syed will ask — with your exact answer",
    "Slides 15-16   Questions you ask HIM — and why",
    "Slides 17-18   Opening, closing, and if things go wrong",
]
for i, item in enumerate(toc):
    txt(sl, Inches(1.0), Inches(3.1)+Inches(0.42)*i, Inches(11), Inches(0.4), item, sz=16, col=LIGHT)

box(sl,Inches(1.0),Inches(5.8),Inches(11),Inches(0.6),fill=CARD2)
txt(sl,Inches(1.3),Inches(5.85),Inches(10),Inches(0.5),
    "Print this. Put it on the desk. Glance at it during the call. Nobody will know.",sz=15,bold=True,col=PINK)
foot(sl,1,TOT)


# ═══════════════════════════════════════════
# 2  WHAT IS MTCP — FROM ZERO
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"What Is MTCP?","Starting from scratch. Forget everything else. This is the simple version.")

box(sl,Inches(0.6),Inches(1.3),Inches(12.1),Inches(1.5),border=GREEN)
txt(sl,Inches(0.9),Inches(1.4),Inches(11.5),Inches(0.3),"THE ONE-SENTENCE VERSION",sz=14,bold=True,col=GREEN)
txt(sl,Inches(0.9),Inches(1.8),Inches(11.5),Inches(0.8),
    "MTCP is a test for AI models. It checks: if an AI breaks a rule and you tell it to stop,\n"
    "does it actually stop — or does it just pretend to stop and then break the rule again?",
    sz=20,col=WHITE)

box(sl,Inches(0.6),Inches(3.1),Inches(12.1),Inches(1.6),fill=CARD2)
txt(sl,Inches(0.9),Inches(3.2),Inches(11.5),Inches(0.3),"THE FULL NAME",sz=14,bold=True,col=PURPLE)
txt(sl,Inches(0.9),Inches(3.55),Inches(11.5),Inches(1.0),
    "MTCP = Multi-Turn Constraint Persistence\n\n"
    "Multi-Turn    = across a conversation (not just one question)\n"
    "Constraint    = a rule the AI must follow (e.g. \"don't use jargon\" or \"don't share personal data\")\n"
    "Persistence   = does the rule STICK after you correct the AI?",
    sz=16,col=LIGHT)

box(sl,Inches(0.6),Inches(5.0),Inches(12.1),Inches(1.7),border=AMBER)
txt(sl,Inches(0.9),Inches(5.1),Inches(11.5),Inches(0.3),"THE ANALOGY (use this on the call if you want)",sz=14,bold=True,col=AMBER)
txt(sl,Inches(0.9),Inches(5.5),Inches(11.5),Inches(1.0),
    "Imagine you hire a new employee. You give them a rule: \"Never share client data on the phone.\"\n"
    "They do it anyway. You correct them. They say sorry.\n"
    "The question is: did they actually learn? Or will they do it again tomorrow?\n\n"
    "MTCP answers that question — but for AI models instead of employees.",
    sz=16,col=WHITE)
foot(sl,2,TOT)


# ═══════════════════════════════════════════
# 3  WHAT DOES MTCP DO — STEP BY STEP
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"What Does MTCP Actually Do?","The 6 steps. Know these — Syed will ask.")

steps = [
    ("STEP 1: Give the AI a rule",
     "We tell the AI: \"In this conversation, you must always do X\" or \"you must never do Y.\"\n"
     "Example: \"Always respond in formal English\" or \"Never reveal you are an AI.\"",
     "This is the CONSTRAINT.",GREEN),
    ("STEP 2: Have a conversation that pushes on that rule",
     "We send messages designed to make the rule hard to follow.\n"
     "Not tricks or hacks — just realistic pressure, like a real user would create.",
     "This is the PRESSURE.",PURPLE),
    ("STEP 3: The AI breaks the rule (most do)",
     "Our system automatically detects when the AI stops following the rule.\n"
     "We have 532 different test scenarios (called PROBES) that do this.",
     "These probes are SECRET — our most valuable IP.",AMBER),
    ("STEP 4: We correct the AI",
     "We tell the AI: \"You just broke the rule. Go back to following it.\"\n"
     "Like a manager correcting an employee.",
     "This is the CORRECTION.",TEAL),
    ("STEP 5: We watch what happens next",
     "Does the AI actually go back to following the rule? Or does it slip again?\n"
     "THIS IS THE KEY BIT. This is what nobody else tests. This is the whole point of MTCP.",
     "This is the PERSISTENCE test.",RED),
    ("STEP 6: We write it all up",
     "Every test produces a signed, tamper-proof document called an Evidence Pack.\n"
     "It records exactly what happened, pass or fail, with cryptographic proof it wasn't edited.",
     "This is the EVIDENCE.",GREEN),
]
for i,(title,desc,note,col) in enumerate(steps):
    r = i // 2; c = i % 2
    l = Inches(0.6) + Inches(6.3)*c
    t = Inches(1.2) + Inches(2.05)*r
    box(sl,l,t,Inches(6.0),Inches(1.85),border=col)
    txt(sl,l+Inches(0.2),t+Inches(0.08),Inches(5.6),Inches(0.35),title,sz=15,bold=True,col=col)
    txt(sl,l+Inches(0.2),t+Inches(0.45),Inches(5.6),Inches(0.7),desc,sz=12,col=LIGHT)
    txt(sl,l+Inches(0.2),t+Inches(1.35),Inches(5.6),Inches(0.4),note,sz=11,bold=True,col=col)
foot(sl,3,TOT)


# ═══════════════════════════════════════════
# 4  WHAT DID MTCP FIND — THE RESULTS
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"What Did MTCP Find?","These are your headline results. Memorise the bold bits.")

box(sl,Inches(0.6),Inches(1.2),Inches(12.1),Inches(0.8),fill=CARD2)
txt(sl,Inches(0.9),Inches(1.25),Inches(11.5),Inches(0.7),
    "We ran 184,387 tests on 32 real AI models from 13 companies (Google, OpenAI, Anthropic, Meta, xAI, etc.)",
    sz=18,bold=True,col=WHITE)

box(sl,Inches(0.6),Inches(2.2),Inches(12.1),Inches(4.5),border=PURPLE)
txt(sl,Inches(0.9),Inches(2.3),Inches(11.5),Inches(0.3),"THE BIG FINDINGS — in order of how impressive they sound",sz=14,bold=True,col=PURPLE)

findings = [
    ("NOT ONE MODEL PASSED.",
     "The highest score was 88.7% out of 100%. You need 90% to pass. Nobody got there.",RED),
    ("The most expensive models are NOT the best.",
     "GPT-4o (OpenAI's top model) scored 64.9% — that's a D. Its CHEAPER little brother GPT-4o-mini scored 81.1% — a B.\n"
     "That means companies paying top price for flagship models are getting WORSE safety performance.",AMBER),
    ("14 out of 32 models have unfixable problems.",
     "These models fail no matter what temperature or settings you use. The failure is baked into the model itself.\n"
     "You can't fix it with guardrails, prompt engineering, or operational controls. It's architectural.",RED),
    ("Open-source (free) models beat expensive commercial ones.",
     "Meta's LLaMA models outperformed Anthropic's Claude and competed with OpenAI's GPT.\n"
     "Companies are paying for a false sense of security.",GREEN),
    ("Safety-tuned models barely change across settings.",
     "Claude (Anthropic) scored almost exactly the same whether you turned the temperature up or down.\n"
     "That means the failure isn't random — it's systematic. The model literally can't do this reliably.",TEAL),
    ("We can tell which failures are fixable and which aren't.",
     "MTCP classifies failures into types. Some models fail randomly (fixable with settings). Some fail architecturally (unfixable).\n"
     "This is the thing enterprises need to know before they deploy.",PURPLE),
]
for i,(finding,detail,col) in enumerate(findings):
    y = Inches(2.7) + Inches(0.65)*i
    txt(sl,Inches(0.9),y,Inches(11.5),Inches(0.3),f"{i+1}.  {finding}",sz=14,bold=True,col=col)
    txt(sl,Inches(1.3),y+Inches(0.28),Inches(11),Inches(0.35),detail,sz=11,col=GREY)
foot(sl,4,TOT)


# ═══════════════════════════════════════════
# 5  THE GRADING SYSTEM + WHAT YOU PRODUCE
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Grades and Outputs","What MTCP produces — the thing sirar would actually sell.")

# grading
box(sl,Inches(0.6),Inches(1.2),Inches(6.0),Inches(3.5),border=PURPLE)
txt(sl,Inches(0.9),Inches(1.35),Inches(5.5),Inches(0.3),"THE GRADING SYSTEM (like school grades)",sz=14,bold=True,col=PURPLE)
grades=[
    ("A+","95%+","No model gets this. It's the gold standard.",GREY),
    ("A","90-94%","No model gets this either. Nobody passes.",GREY),
    ("B","80-89%","Only 3 models. Best available right now.",GREEN),
    ("C","70-79%","Needs a risk review before deployment.",AMBER),
    ("D","60-69%","Most common grade. Not safe to deploy.",RED),
    ("F","Below 60%","Failing. Should not be used in regulated settings.",RED),
]
for i,(g,pct,note,col) in enumerate(grades):
    y = Inches(1.75) + Inches(0.42)*i
    txt(sl,Inches(0.9),y,Inches(0.5),Inches(0.3),g,sz=16,bold=True,col=col)
    txt(sl,Inches(1.5),y,Inches(1.2),Inches(0.3),pct,sz=13,col=LIGHT)
    txt(sl,Inches(2.8),y,Inches(3.5),Inches(0.3),note,sz=12,col=GREY)

txt(sl,Inches(0.9),Inches(4.3),Inches(5.5),Inches(0.3),
    "Why this matters: a GRADE is something a buyer understands.",sz=12,bold=True,col=AMBER)

# outputs
box(sl,Inches(6.9),Inches(1.2),Inches(5.8),Inches(3.5),border=TEAL)
txt(sl,Inches(7.2),Inches(1.35),Inches(5.2),Inches(0.3),"WHAT THE CLIENT GETS",sz=14,bold=True,col=TEAL)

outputs = [
    ("Evidence Pack","A signed document showing exactly what the AI\ndid in each test. Tamper-proof. Audit-ready."),
    ("Grade","A simple letter grade (A-F) for the model.\nInstantly understandable."),
    ("Executive Summary","A short report explaining results in\nplain language for leadership."),
    ("Risk Classification","Which failures are fixable (settings)\nvs unfixable (architectural)."),
    ("Release Recommendation","\"Safe to deploy\" / \"Deploy with monitoring\" /\n\"Do not deploy\" for each model."),
]
for i,(name,desc) in enumerate(outputs):
    y = Inches(1.75) + Inches(0.52)*i
    txt(sl,Inches(7.2),y,Inches(2.0),Inches(0.3),name,sz=13,bold=True,col=GREEN)
    txt(sl,Inches(9.3),y,Inches(3.2),Inches(0.45),desc,sz=10,col=LIGHT)

# the key line
box(sl,Inches(0.6),Inches(5.0),Inches(12.1),Inches(1.0),fill=CARD2,border=GREEN)
txt(sl,Inches(0.9),Inches(5.05),Inches(11.5),Inches(0.9),
    "The point: MTCP doesn't just say \"this model is bad.\" It produces EVIDENCE that a company can use\n"
    "in reports, audits, regulatory filings, and client conversations. That evidence is what sirar would sell.\n"
    "It's like a lab report — not an opinion, not a guess — documented proof of what the AI actually did.",
    sz=15,col=WHITE)

# EU AI Act
box(sl,Inches(0.6),Inches(6.2),Inches(12.1),Inches(0.6),border=AMBER)
txt(sl,Inches(0.9),Inches(6.25),Inches(11.5),Inches(0.5),
    "EU AI Act deadline: August 2026 (4 months away). Companies deploying AI in Europe need this evidence. Most don't have it yet.",
    sz=14,bold=True,col=AMBER)
foot(sl,5,TOT)


# ═══════════════════════════════════════════
# 6  WHO IS SIRAR — SIMPLE VERSION
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Who Is sirar?","Know them before they know you.")

box(sl,Inches(0.6),Inches(1.2),Inches(12.1),Inches(2.5),border=TEAL)
txt(sl,Inches(0.9),Inches(1.3),Inches(11.5),Inches(0.3),"THE BASICS",sz=14,bold=True,col=TEAL)
basics=[
    "sirar by stc — a cybersecurity company based in Riyadh, Saudi Arabia.",
    "Owned by STC Group (Saudi Telecom Company) — the biggest telecom in the Middle East. Worth $50B+.",
    "About 500 employees across 3 continents. #1 cybersecurity company in the Middle East.",
    "They protect Saudi critical infrastructure — railways, airports, water systems, banks, even the Hajj pilgrimage.",
    "They do: threat detection, incident response, compliance advisory, managed security, and their own security products.",
    "They're part of Saudi Vision 2030 — the country's plan to become a global tech hub.",
    "They already do compliance work (Saudi regulations, ISO standards, banking frameworks).",
    "They are NOW expanding into AI services. Their website has an /ai/ page but it's basically empty.",
    "That means their AI service is NEW and being built. They don't have AI testing capability yet. That's you.",
]
for i,item in enumerate(basics):
    txt(sl,Inches(0.9),Inches(1.7)+Inches(0.33)*i,Inches(11.5),Inches(0.3),f"  {item}",sz=13,col=LIGHT)

box(sl,Inches(0.6),Inches(4.0),Inches(12.1),Inches(2.5),border=AMBER)
txt(sl,Inches(0.9),Inches(4.1),Inches(11.5),Inches(0.3),"WHO IS SYED?",sz=14,bold=True,col=AMBER)
txt(sl,Inches(0.9),Inches(4.5),Inches(11.5),Inches(1.8),
    "Most likely: Mudassir Syed — Senior Manager for Product Management & Marketing at sirar.\n\n"
    "He's a PRODUCT person, not a deep technical person. He builds new services and takes them to market.\n"
    "He's been at sirar for about 5 years. Before that he worked at big telecoms (Etisalat in UAE).\n"
    "He has an MBA and product management certifications.\n\n"
    "What this means: He thinks in terms of \"what can we sell\" and \"how do we package this.\"\n"
    "Don't go deep on maths or theory. Talk about what the service looks like, who buys it, and what they get.",
    sz=14,col=LIGHT)

foot(sl,6,TOT)


# ═══════════════════════════════════════════
# 7  WHAT SIRAR WANTS FROM YOU
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"What sirar Wants From You","Three things. This is what the whole call is about.")

wants = [
    ("1. YOUR PLATFORM",
     "They want MTCP's testing capability as part of their AI Assurance service.\n"
     "They want to be able to say to THEIR clients: \"We tested your AI and here are the results.\"\n"
     "The Evidence Packs, the grades, the reports — that's what they'd use.",
     "Think of it like: they're a hospital, they want your lab equipment to run blood tests.",
     PURPLE),
    ("2. YOUR BRAIN",
     "They don't have anyone who understands AI safety testing.\n"
     "They need someone who can: explain the method, interpret results, design test scopes,\n"
     "and help write reports their clients will understand.\n"
     "That person is you.",
     "Think of it like: they want the specialist doctor, not just the medical equipment.",
     TEAL),
    ("3. A COMMERCIAL STRUCTURE",
     "They need to figure out: how do we work together?\n"
     "Is it a pilot first? A licence? Do you join as a consultant?\n"
     "This call is about figuring out which shape makes sense.",
     "You DON'T need to decide this today. Just understand what they're looking for.",
     AMBER),
]
for i,(title,desc,analogy,col) in enumerate(wants):
    l = Inches(0.6) + Inches(4.2)*i
    box(sl,l,Inches(1.2),Inches(3.9),Inches(4.2),border=col)
    txt(sl,l+Inches(0.2),Inches(1.35),Inches(3.5),Inches(0.4),title,sz=18,bold=True,col=col)
    txt(sl,l+Inches(0.2),Inches(1.85),Inches(3.5),Inches(1.8),desc,sz=13,col=LIGHT)
    txt(sl,l+Inches(0.2),Inches(3.9),Inches(3.5),Inches(0.8),analogy,sz=12,col=GREY)

box(sl,Inches(0.6),Inches(5.7),Inches(12.1),Inches(1.0),fill=CARD2,border=RED)
txt(sl,Inches(0.9),Inches(5.8),Inches(11.5),Inches(0.8),
    "THE IMPORTANT BIT: Your brain (consulting/SME work) and your platform (the IP) are TWO SEPARATE THINGS.\n"
    "If someone asks you to consult AND gives you a salary, that does NOT include them getting your platform.\n"
    "The platform is a separate commercial asset. Never bundle them. Say: \"Those are separate conversations.\"",
    sz=14,bold=True,col=RED)
foot(sl,7,TOT)


# ═══════════════════════════════════════════
# 8  WHY THIS MAKES SENSE — FOR BOTH SIDES
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Why This Makes Sense — For Both Sides","Understand the mutual value so you can talk about it naturally.")

box(sl,Inches(0.6),Inches(1.2),Inches(5.8),Inches(5.0),border=GREEN)
txt(sl,Inches(0.9),Inches(1.35),Inches(5.2),Inches(0.3),"WHAT SIRAR GETS FROM YOU",sz=16,bold=True,col=GREEN)
them=[
    "A ready-made AI testing platform — no need to build one",
    "184,387 evaluations already done — instant credibility",
    "An expert who can explain results to their clients",
    "Published research (5 papers, DOI) — academic legitimacy",
    "Evidence Packs designed for regulatory compliance",
    "Grades that non-technical buyers instantly understand",
    "EU AI Act readiness — their clients will need this soon",
    "A head start — nobody else in the Gulf has this",
    "Already being integrated into Saudi sovereign AI work (KFUPM)",
    "Something their competitors don't have",
]
for i,item in enumerate(them):
    txt(sl,Inches(0.9),Inches(1.75)+Inches(0.4)*i,Inches(5.2),Inches(0.35),f"  {item}",sz=13,col=LIGHT)

box(sl,Inches(6.7),Inches(1.2),Inches(5.8),Inches(5.0),border=PURPLE)
txt(sl,Inches(7.0),Inches(1.35),Inches(5.2),Inches(0.3),"WHAT YOU GET FROM SIRAR",sz=16,bold=True,col=PURPLE)
you=[
    "Access to Saudi/Gulf enterprise clients you can't reach alone",
    "A sales and delivery team (500 people, 3 continents)",
    "Credibility of being partnered with #1 MSSP in Middle East",
    "Revenue — pilot fees, then ongoing access or licence fees",
    "Real-world deployment of MTCP in a commercial service",
    "A reference case for future deals",
    "Connection to Vision 2030 mega-projects",
    "Entry into banking, government, critical infrastructure",
    "Your IP stays yours — they use outputs, not the engine",
    "A path from solo researcher to commercial platform",
]
for i,item in enumerate(you):
    txt(sl,Inches(7.0),Inches(1.75)+Inches(0.4)*i,Inches(5.2),Inches(0.35),f"  {item}",sz=13,col=LIGHT)

foot(sl,8,TOT)


# ═══════════════════════════════════════════
# 9  THE PILOT — WHAT YOU'RE PROPOSING
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"The Pilot — What You're Proposing","This is your main offer. Everything should steer toward this.")

box(sl,Inches(0.6),Inches(1.2),Inches(12.1),Inches(1.2),fill=CARD2,border=GREEN)
txt(sl,Inches(0.9),Inches(1.3),Inches(11.5),Inches(1.0),
    "\"Let's do a small test run first. I'll evaluate 3 to 5 AI models that matter to sirar or their clients.\n"
    "It takes 2 to 4 weeks. You get Evidence Packs, grades, an executive summary, and a recommendation.\n"
    "That way sirar can see exactly what the output looks like before we talk about anything bigger.\"",
    sz=17,col=WHITE)

box(sl,Inches(0.6),Inches(2.7),Inches(5.8),Inches(2.0),border=TEAL)
txt(sl,Inches(0.9),Inches(2.85),Inches(5.2),Inches(0.3),"WHAT'S IN THE PILOT",sz=14,bold=True,col=TEAL)
inc=["3-5 models tested (they pick which ones matter to them)",
     "Full MTCP evaluation at multiple temperature settings",
     "Evidence Packs for each model (signed, tamper-proof)",
     "Failure classification (fixable vs unfixable)",
     "Executive summary for their leadership",
     "Release-readiness recommendation per model"]
for i,item in enumerate(inc):
    txt(sl,Inches(0.9),Inches(3.25)+Inches(0.33)*i,Inches(5.2),Inches(0.3),f"  {item}",sz=12,col=LIGHT)

box(sl,Inches(6.7),Inches(2.7),Inches(5.8),Inches(2.0),border=RED)
txt(sl,Inches(7.0),Inches(2.85),Inches(5.2),Inches(0.3),"WHAT'S NOT IN THE PILOT",sz=14,bold=True,col=RED)
notinc=["Source code or platform access",
        "The probe sets (the test scenarios) — these are secret",
        "The detector logic (how it scores)",
        "A licence to run MTCP themselves",
        "Exclusivity of any kind",
        "Unlimited rights — that's a different conversation"]
for i,item in enumerate(notinc):
    txt(sl,Inches(7.0),Inches(3.25)+Inches(0.33)*i,Inches(5.2),Inches(0.3),f"  {item}",sz=12,col=LIGHT)

box(sl,Inches(0.6),Inches(5.0),Inches(12.1),Inches(1.5),border=AMBER)
txt(sl,Inches(0.9),Inches(5.1),Inches(11.5),Inches(0.3),"WHY A PILOT IS THE RIGHT MOVE",sz=14,bold=True,col=AMBER)
txt(sl,Inches(0.9),Inches(5.5),Inches(11.5),Inches(0.9),
    "1. It's low risk for both sides — small scope, short timeline, clear deliverables.\n"
    "2. sirar gets to see real output before committing to anything bigger.\n"
    "3. You don't have to hand over any IP — you deliver results, not the engine.\n"
    "4. It creates the basis for pricing: once they see what it produces, you can talk numbers properly.\n"
    "5. If it doesn't work out, nobody's lost much. If it does, you have a real partnership.",
    sz=13,col=LIGHT)
foot(sl,9,TOT)


# ═══════════════════════════════════════════
# 10  QUESTIONS HE'LL ASK — ABOUT MTCP
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Questions Syed Will Ask — About MTCP","He needs to understand what he's buying. These are the obvious ones.")

qas = [
    ("\"So what exactly is MTCP?\"",
     "It's a testing platform for AI models. We check whether an AI actually follows rules after you correct it.\n"
     "We've tested 32 real models from 13 providers — 181,000+ tests. Not one model passed.",TEAL),
    ("\"How does it work?\"",
     "We give the AI a rule, have a conversation that puts pressure on that rule, and when the AI breaks it,\n"
     "we correct it and watch whether the correction actually holds. Then we grade it and produce evidence.",PURPLE),
    ("\"What models have you tested?\"",
     "32 production models — GPT-4o, Claude, LLaMA, Gemini, Mistral, DeepSeek, and more.\n"
     "The full leaderboard is on mtcp.live. The best model scored 88.7%, worst was 37.3%.",GREEN),
    ("\"Can it test any model?\"",
     "Any model we can access through an API or a controlled interface. It's completely black-box —\n"
     "we don't need access to the model's code, weights, or training data. Just the ability to talk to it.",AMBER),
    ("\"How long does a test take?\"",
     "For the pilot, 2 to 4 weeks depending on how many models and how deep we go.\n"
     "A single model at one temperature setting takes a few hours to run.",TEAL),
    ("\"What do we actually get at the end?\"",
     "An Evidence Pack for each model (signed, tamper-proof, audit-ready), a grade (A to F),\n"
     "an executive summary, failure classification, and a deployment recommendation.",GREEN),
]
for i,(q,a,col) in enumerate(qas):
    y = Inches(1.15) + Inches(1.0)*i
    box(sl,Inches(0.6),y,Inches(12.1),Inches(0.9),border=col)
    txt(sl,Inches(0.9),y+Inches(0.03),Inches(11.5),Inches(0.3),q,sz=14,bold=True,col=col)
    txt(sl,Inches(0.9),y+Inches(0.35),Inches(11.5),Inches(0.5),a,sz=12,col=LIGHT)
foot(sl,10,TOT)


# ═══════════════════════════════════════════
# 11  QUESTIONS — ABOUT DIFFERENTIATION
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Questions Syed Will Ask — How Is This Different?","He'll compare you to things he already knows. Be ready.")

qas2 = [
    ("\"How is this different from red teaming?\"",
     "Red teaming tries to break AI (can it fail?). We test what happens AFTER it fails and you correct it (does the fix stick?).\n"
     "Red teaming is offence. MTCP is assurance. They're complementary — not competing.",PURPLE),
    ("\"How is this different from monitoring tools?\"",
     "Monitoring watches a live system after deployment. MTCP gives you evidence BEFORE deployment —\n"
     "and you can re-run it later as a regression check. It's proactive, not reactive.",TEAL),
    ("\"How is this different from benchmarks like MMLU or HellaSwag?\"",
     "Those test general knowledge — can the AI answer questions correctly? MTCP tests behaviour —\n"
     "can the AI follow rules reliably? Totally different thing. A model can be smart but ungovernable.",AMBER),
    ("\"Isn't this just prompt injection testing?\"",
     "No. Prompt injection is about tricking the AI. MTCP is about normal conversation pressure.\n"
     "We're not attacking the model. We're checking if corrections actually work in realistic use.",GREEN),
    ("\"Are there other tools that do this?\"",
     "Not for post-correction persistence specifically. There are red-teaming tools, monitoring tools,\n"
     "and general benchmarks — but nobody else measures whether correction holds across a conversation.\n"
     "That's the gap. That's why we have 181,000+ evaluations and published research — we defined this space.",PINK),
    ("\"Is this peer-reviewed?\"",
     "We have 5 published papers on OSF with a DOI (10.17605/OSF.IO/DXGK5) and 22 more completed.\n"
     "Six independent practitioners validated our findings within 48 hours. One formally cited us in a company document.",PURPLE),
]
for i,(q,a,col) in enumerate(qas2):
    y = Inches(1.15) + Inches(1.0)*i
    box(sl,Inches(0.6),y,Inches(12.1),Inches(0.9),border=col)
    txt(sl,Inches(0.9),y+Inches(0.03),Inches(11.5),Inches(0.3),q,sz=14,bold=True,col=col)
    txt(sl,Inches(0.9),y+Inches(0.35),Inches(11.5),Inches(0.5),a,sz=12,col=LIGHT)
foot(sl,11,TOT)


# ═══════════════════════════════════════════
# 12  QUESTIONS — ABOUT THE BUSINESS
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Questions Syed Will Ask — About Business","These are the tricky ones. Know your lines.")

qas3 = [
    ("\"How much does it cost?\"",
     "\"It depends on the scope — how many models, how deep, and whether this is a one-off pilot\n"
     "or ongoing access. The cleanest thing is to define the pilot scope first, then I can put a number to it.\"",
     "NEVER name a price first. Make him define what he wants, then you price it.",RED),
    ("\"Can we run MTCP ourselves?\"",
     "\"That's a licensing model — definitely possible, but it's a different conversation from a pilot.\n"
     "The pilot lets you see the output first. Licensing is a bigger commercial discussion for later.\"",
     "Don't say no. Don't say yes. Say it's a separate discussion.",AMBER),
    ("\"Can we white-label it? / Put our name on it?\"",
     "\"Absolutely possible as a commercial arrangement. Let's see how the pilot goes first\n"
     "and then we can scope out what a white-label model looks like.\"",
     "Same tactic: yes in principle, but after the pilot.",AMBER),
    ("\"Who else is using this?\"",
     "\"MTCP has been validated by independent practitioners in sovereign AI governance, enterprise AI adoption,\n"
     "and health informatics. The methodology is being integrated into sovereign AI governance systems in the Gulf.\n"
     "I can share more detail once we have a framework for the engagement.\"",
     "Be vague-but-impressive. Don't name names without permission.",TEAL),
    ("\"What's your team like?\"",
     "\"The MTCP research programme produces the methodology, evaluations, and platform.\n"
     "I lead the methodology and interpretation. For the pilot, I'd be your direct point of contact.\"",
     "Don't say \"it's just me.\" Say \"I lead\" or \"the programme.\" Sound established.",PURPLE),
    ("\"Can we have an exclusive deal for the Gulf?\"",
     "\"Exclusivity is definitely on the table as a commercial option — but that's a bigger conversation\n"
     "with different commercial terms. Let's prove the value with a pilot first.\"",
     "Exclusivity = expensive. Don't give it away. Don't shut the door either.",PINK),
]
for i,(q,a,note,col) in enumerate(qas3):
    y = Inches(1.1) + Inches(1.02)*i
    box(sl,Inches(0.6),y,Inches(9.0),Inches(0.9),border=col)
    txt(sl,Inches(0.9),y+Inches(0.03),Inches(8.5),Inches(0.3),q,sz=13,bold=True,col=col)
    txt(sl,Inches(0.9),y+Inches(0.33),Inches(8.5),Inches(0.55),a,sz=11,col=LIGHT)
    txt(sl,Inches(9.8),y+Inches(0.1),Inches(2.8),Inches(0.7),note,sz=11,bold=True,col=col)
foot(sl,12,TOT)


# ═══════════════════════════════════════════
# 13  QUESTIONS — ABOUT COMPLIANCE
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Questions Syed Will Ask — Compliance & Regulation","sirar lives in compliance world. Be ready for these.")

qas4 = [
    ("\"Does this map to the EU AI Act?\"",
     "\"Yes. MTCP was designed with the EU AI Act in mind. Our Evidence Packs are formatted for Annex IV\n"
     "technical documentation. We map to Article 9 (risk management), Article 61 (post-market monitoring),\n"
     "and Article 72 (incident reporting). The deadline is August 2026 — 4 months away.\"",GREEN),
    ("\"What about Saudi regulations? NCA? SDAIA?\"",
     "\"Saudi AI governance frameworks are being developed right now by SDAIA and NDMO. They're drawing on\n"
     "EU AI Act structure. MTCP's compliance mapping is directly transferable. We're also already being\n"
     "integrated into sovereign AI governance work at KFUPM in Saudi Arabia.\"",TEAL),
    ("\"Does it cover ISO 42001?\"",
     "\"ISO 42001 is the new AI management systems standard. MTCP provides the behavioural evidence\n"
     "that feeds into ISO 42001 compliance — specifically the monitoring and measurement requirements.\n"
     "It doesn't replace the standard, but it generates the evidence the standard requires.\"",PURPLE),
    ("\"Can the Evidence Packs be used in audits?\"",
     "\"That's exactly what they're designed for. Every Evidence Pack is cryptographically signed (SHA-256),\n"
     "tamper-proof, and contains the full test record. They can be filed directly as technical documentation\n"
     "for regulatory submissions.\"",AMBER),
    ("\"What if a model gets updated after we test it?\"",
     "\"The certificate is valid for the specific model version tested. If the provider updates the model,\n"
     "we recommend re-testing. That's actually a revenue opportunity for sirar — ongoing assurance monitoring.\"",GREEN),
]
for i,(q,a,col) in enumerate(qas4):
    y = Inches(1.1) + Inches(1.15)*i
    box(sl,Inches(0.6),y,Inches(12.1),Inches(1.05),border=col)
    txt(sl,Inches(0.9),y+Inches(0.03),Inches(11.5),Inches(0.3),q,sz=14,bold=True,col=col)
    txt(sl,Inches(0.9),y+Inches(0.35),Inches(11.5),Inches(0.65),a,sz=12,col=LIGHT)
foot(sl,13,TOT)


# ═══════════════════════════════════════════
# 14  QUESTIONS — AWKWARD ONES
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"The Awkward Questions","The ones that might catch you off guard. Rehearse these.")

awks = [
    ("\"Can you show us the test scenarios / probes?\"",
     "\"The methodology is fully published — 5 papers on OSF. But the actual test scenarios are proprietary.\n"
     "That's the core IP. Sharing those would be like a drug company sharing its formula.\n"
     "I can show you what the OUTPUT looks like — the Evidence Packs and reports. That's what matters for the service.\"",RED),
    ("\"Why hasn't this been adopted by big companies yet?\"",
     "\"The research was published this year. Six independent experts validated it within 48 hours.\n"
     "We're in early conversations with multiple parties. sirar would be an early mover — which is an advantage,\n"
     "not a risk. The companies that adopt AI assurance first will own the market.\"",AMBER),
    ("\"Are you a one-person operation?\"",
     "\"I lead the MTCP research programme. The platform, methodology, and 181,000+ evaluations are the work.\n"
     "The value is in what's been built and validated, not headcount. If we work together, we'd discuss\n"
     "what support structure makes sense for the engagement.\"",
     PURPLE),
    ("\"What if our clients' models aren't in your database?\"",
     "\"We can test any model with API access. The 32 models in the database are the published benchmark.\n"
     "The pilot would specifically test models that matter to sirar and your clients.\"",TEAL),
    ("\"Can you just send us the platform / code?\"",
     "\"The platform is a commercial asset with its own value. Access to it is a licensing conversation.\n"
     "For the pilot, I deliver the outputs — Evidence Packs, grades, reports. You don't need the platform\n"
     "to use the results. That's the whole point of the service model.\"",RED),
]
for i,(q,a,col) in enumerate(awks):
    y = Inches(1.1) + Inches(1.2)*i
    box(sl,Inches(0.6),y,Inches(12.1),Inches(1.1),border=col)
    txt(sl,Inches(0.9),y+Inches(0.03),Inches(11.5),Inches(0.3),q,sz=14,bold=True,col=col)
    txt(sl,Inches(0.9),y+Inches(0.38),Inches(11.5),Inches(0.65),a,sz=12,col=LIGHT)
foot(sl,14,TOT)


# ═══════════════════════════════════════════
# 15  YOUR QUESTIONS FOR HIM
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Questions YOU Ask Him","Asking good questions makes you look smart and gets you the info you need.")

yourqs = [
    ("\"Tell me about the AI Assurance service sirar is building — where are you in the launch?\"",
     "Opens the conversation. Makes him talk first. You learn what stage they're at.",
     "Ask this FIRST.",GREEN),
    ("\"Are the first clients likely to be your existing cyber clients — banks, government, infrastructure?\"",
     "Tells you who the Evidence Packs would be for. Shapes the pilot scope.",
     "Ask early.",TEAL),
    ("\"Is sirar looking to build AI testing in-house, or partner with someone who already has the platform?\"",
     "THE most important question. If in-house: they want to buy your IP. If partner: they want your service.\n"
     "Completely different commercial models.",
     "Ask this when you feel the vibe.",AMBER),
    ("\"Is there pressure from NCA or SDAIA for AI-specific compliance, or is this more market-driven?\"",
     "Shows you know Saudi regulatory landscape. Tells you if they need evidence NOW or are getting ahead.",
     "Shows you did homework.",PURPLE),
    ("\"What would success look like for this service in its first year?\"",
     "Makes him paint the picture. Then you show how MTCP accelerates getting there.",
     "Ask midway through.",GREEN),
    ("\"What would make a pilot successful from sirar's perspective?\"",
     "Gets him to define success criteria — which you then use to write the pilot proposal.\n"
     "This is the most important question in the whole call.",
     "Ask near the end. THE question.",PINK),
]
for i,(q,why,when,col) in enumerate(yourqs):
    y = Inches(1.1) + Inches(1.0)*i
    box(sl,Inches(0.6),y,Inches(9.0),Inches(0.9),border=col)
    txt(sl,Inches(0.9),y+Inches(0.03),Inches(8.5),Inches(0.3),q,sz=13,bold=True,col=WHITE)
    txt(sl,Inches(0.9),y+Inches(0.35),Inches(8.5),Inches(0.5),why,sz=11,col=GREY)
    box(sl,Inches(9.8),y+Inches(0.2),Inches(2.8),Inches(0.5),fill=col)
    txt(sl,Inches(9.85),y+Inches(0.25),Inches(2.7),Inches(0.4),when,sz=12,bold=True,col=BG,align=PP_ALIGN.CENTER)
foot(sl,15,TOT)


# ═══════════════════════════════════════════
# 16  MORE QUESTIONS + FOLLOW-UP QUESTIONS
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Follow-Up Questions (Based on What He Says)","Listen to his answers, then use these to go deeper.")

box(sl,Inches(0.6),Inches(1.15),Inches(12.1),Inches(5.5),border=PURPLE)
txt(sl,Inches(0.9),Inches(1.25),Inches(11.5),Inches(0.3),"IF HE SAYS THIS → YOU ASK THIS",sz=16,bold=True,col=PURPLE)

followups = [
    ("\"We're targeting enterprise clients\"",
     "\"Which industries? Banking, healthcare, government? That helps me think about which models to test in the pilot.\"",TEAL),
    ("\"We want to offer this to government\"",
     "\"Is that Saudi government specifically, or GCC-wide? And is there a specific regulation driving the timeline?\"",AMBER),
    ("\"We need something we can run ourselves\"",
     "\"Understood. That's a licensing model. Let's prove the value with a pilot first, then we can scope licensing properly.\"",RED),
    ("\"We already have some AI tools\"",
     "\"Interesting — what are they focused on? Detection? Monitoring? MTCP specifically covers post-correction\n"
     "persistence, which is a different layer. They'd likely complement each other.\"",GREEN),
    ("\"Our clients are asking about EU AI Act\"",
     "\"Perfect — MTCP Evidence Packs are designed for exactly that. The deadline is August 2026.\n"
     "A pilot now means sirar could offer compliance evidence to those clients within weeks.\"",GREEN),
    ("\"We're early — still figuring this out\"",
     "\"That's actually ideal. A small pilot gives you something concrete to build the service around.\n"
     "Better to shape the service around real output than design it in theory and hope it works.\"",TEAL),
    ("\"What would the pricing look like?\"",
     "\"It depends on scope. Can I ask — what would make a pilot successful from sirar's perspective?\n"
     "Once I know that, I can put a proper proposal together with pricing.\"",AMBER),
    ("\"Can you send us some materials?\"",
     "\"Absolutely. I'll put together a pilot proposal — scope, timeline, deliverables. Would a one-pager work?\n"
     "I can also share the public research overview and an example of what an Evidence Pack looks like.\"",GREEN),
]
for i,(hesays,youask,col) in enumerate(followups):
    y = Inches(1.65) + Inches(0.62)*i
    txt(sl,Inches(0.9),y,Inches(4.8),Inches(0.3),hesays,sz=12,bold=True,col=col)
    txt(sl,Inches(5.9),y,Inches(6.5),Inches(0.55),youask,sz=11,col=LIGHT)
foot(sl,16,TOT)


# ═══════════════════════════════════════════
# 17  OPENING AND CLOSING — EXACT WORDS
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"Opening and Closing — Word for Word","Start strong. End with a next step. Everything else is conversation.")

box(sl,Inches(0.6),Inches(1.15),Inches(12.1),Inches(1.6),border=GREEN)
txt(sl,Inches(0.9),Inches(1.25),Inches(11.5),Inches(0.3),"YOUR OPENING (after hello / how are you / pleasantries)",sz=14,bold=True,col=GREEN)
txt(sl,Inches(0.9),Inches(1.6),Inches(11.5),Inches(1.0),
    "\"Thanks for setting this up, Syed. I've been looking forward to it.\n"
    "Before I go into MTCP, I'd love to hear from you — tell me about the AI Assurance service\n"
    "sirar is building. What's the vision?\"",
    sz=18,col=WHITE)

box(sl,Inches(0.6),Inches(3.0),Inches(12.1),Inches(1.6),border=TEAL)
txt(sl,Inches(0.9),Inches(3.1),Inches(11.5),Inches(0.3),"CLOSING IF HE'S INTERESTED",sz=14,bold=True,col=TEAL)
txt(sl,Inches(0.9),Inches(3.45),Inches(11.5),Inches(1.0),
    "\"Great — I think there's a real fit here. The cleanest next step from my side is a short pilot proposal:\n"
    "scope, timeline, deliverables, and commercial structure. I'll get that to you this week.\n"
    "Does that work?\"",
    sz=18,col=WHITE)

box(sl,Inches(0.6),Inches(4.85),Inches(12.1),Inches(1.0),border=AMBER)
txt(sl,Inches(0.9),Inches(4.95),Inches(11.5),Inches(0.3),"CLOSING IF HE'S VAGUE",sz=14,bold=True,col=AMBER)
txt(sl,Inches(0.9),Inches(5.3),Inches(11.5),Inches(0.5),
    "\"No pressure at all. I'll send you a short summary of what MTCP produces and how a pilot would work.\n"
    "You can share it internally and we can pick up when the timing's right.\"",
    sz=18,col=WHITE)

box(sl,Inches(0.6),Inches(6.1),Inches(12.1),Inches(0.7),border=RED)
txt(sl,Inches(0.9),Inches(6.15),Inches(11.5),Inches(0.3),"CLOSING IF HE PUSHES FOR FREE STUFF",sz=14,bold=True,col=RED)
txt(sl,Inches(0.9),Inches(6.45),Inches(11.5),Inches(0.3),
    "\"The pilot is designed to be the low-commitment proof of concept. That's the cleanest starting point.\"",
    sz=18,col=WHITE)
foot(sl,17,TOT)


# ═══════════════════════════════════════════
# 18  IF THINGS GO WRONG
# ═══════════════════════════════════════════
sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
head(sl,"If Things Go Wrong","Rescue lines. Memorise 2 or 3. They'll save you.")

rescues = [
    ("You lose your train of thought",
     "\"Let me come back to the key point...\"",
     "Then say the ONE-LINER: \"MTCP turns model behaviour into assurance evidence.\"",TEAL),
    ("He asks something you don't know",
     "\"That's a really good question. I want to give you a proper answer — let me follow up in writing.\"",
     "This is what confident people say. Never bluff.",PURPLE),
    ("He gets too technical / deep",
     "\"I'd love to go deeper on that — and we can, once we've scoped the pilot. For today the key thing is what you'd receive.\"",
     "Steers back to the output, not the engine.",AMBER),
    ("Awkward silence",
     "\"What are your initial thoughts on how this could fit into what sirar is building?\"",
     "Turns the silence into his turn to talk.",GREEN),
    ("He's rushing / distracted",
     "\"I know you're busy. The one thing I'd love to leave you with: we've tested 32 AI models, none passed,\n"
     "and the evidence we produce is exactly what the EU AI Act requires. Shall I send you a pilot proposal?\"",
     "Compress everything into 15 seconds.",PINK),
    ("You blank completely",
     "\"Sorry — I want to make sure I explain this clearly. The core of it is: MTCP gives companies proof\n"
     "of whether their AI models are actually safe. Not opinion — evidence. That's what we'd provide to sirar.\"",
     "This works no matter what the question was.",RED),
    ("He asks for something free",
     "\"The pilot is designed to be exactly that — a low-commitment way to see the value.\n"
     "It's the cleanest starting point for both sides.\"",
     "Always redirect to the pilot. Every time.",AMBER),
    ("He says \"we'll think about it\" (vague brush-off)",
     "\"Of course. I'll send the pilot proposal so you have something concrete to look at.\n"
     "Is email the best way to reach you?\"",
     "Always leave with a concrete next action, even if small.",GREEN),
]
for i,(situation,say,note,col) in enumerate(rescues):
    y = Inches(1.05) + Inches(0.77)*i
    txt(sl,Inches(0.6),y,Inches(3.0),Inches(0.3),situation,sz=12,bold=True,col=col)
    txt(sl,Inches(3.7),y,Inches(5.5),Inches(0.65),say,sz=11,col=WHITE)
    txt(sl,Inches(9.5),y+Inches(0.05),Inches(3.2),Inches(0.65),note,sz=10,col=GREY)
foot(sl,18,TOT)


# ═══════════════════════════════════════════
out = os.path.expanduser("~/Desktop/SIRAR_Call_Cheat_Sheet.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
